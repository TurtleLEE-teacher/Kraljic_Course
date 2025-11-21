#!/usr/bin/env python3
"""
Part 2 PPTX Generator - ENHANCED VERSION with Visual Focus
Session 2: 자재군별 소싱 전략 및 공급업체 관계 관리
S4HANA Professional Style (10.83" × 7.50")

IMPROVEMENTS:
- Visual-first approach (60-70% diagrams, 30-40% text)
- Consistent font sizes (10pt body, 12pt bullets, 16pt governing)
- Text overflow prevention
- Native PowerPoint shapes (editable)
- Reduced text content
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# COLOR SYSTEM (Monochrome)
# ============================================================================
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)
COLOR_MED_GRAY = RGBColor(102, 102, 102)
COLOR_LIGHT_GRAY = RGBColor(204, 204, 204)
COLOR_VERY_LIGHT_GRAY = RGBColor(230, 230, 230)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(26, 82, 118)  # Dark blue

# Font size constants (consistent)
FONT_TITLE = Pt(20)
FONT_GOVERNING = Pt(16)
FONT_HEADING = Pt(14)
FONT_BODY = Pt(10)  # PRIMARY
FONT_BULLET = Pt(12)
FONT_CAPTION = Pt(8)

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def create_presentation():
    """Create presentation with S4HANA dimensions"""
    prs = Presentation()
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_and_governing_message(slide, title_text, governing_msg_text):
    """Add title and governing message to content slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.83), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "맑은 고딕"
    title_para.font.size = FONT_TITLE
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_BLACK

    # Title underline
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.95),
        Inches(9.83), Inches(0.03)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLOR_LIGHT_GRAY
    slide.shapes[-1].line.color.rgb = COLOR_LIGHT_GRAY

    # Governing message (16pt Bold)
    gov_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9.83), Inches(0.5))
    gov_frame = gov_box.text_frame
    gov_frame.text = governing_msg_text
    gov_frame.word_wrap = True
    gov_para = gov_frame.paragraphs[0]
    gov_para.font.name = "맑은 고딕"
    gov_para.font.size = FONT_GOVERNING
    gov_para.font.bold = True
    gov_para.font.color.rgb = COLOR_MED_GRAY

def add_process_step_box(slide, x, y, width, height, number, label, detail=None):
    """Add a process step box with number circle"""
    # Main box
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y),
        Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    box.line.color.rgb = COLOR_MED_GRAY
    box.line.width = Pt(2)

    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(x + 0.15), Inches(y + 0.15),
        Inches(0.4), Inches(0.4)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_ACCENT
    circle.line.color.rgb = COLOR_ACCENT

    num_text = circle.text_frame
    num_text.text = str(number)
    num_para = num_text.paragraphs[0]
    num_para.font.name = "Arial"
    num_para.font.size = Pt(18)
    num_para.font.bold = True
    num_para.font.color.rgb = COLOR_WHITE
    num_para.alignment = PP_ALIGN.CENTER
    num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Label
    label_box = slide.shapes.add_textbox(
        Inches(x + 0.1), Inches(y + 0.6),
        Inches(width - 0.2), Inches(0.5)
    )
    label_frame = label_box.text_frame
    label_frame.text = label
    label_frame.word_wrap = True
    label_para = label_frame.paragraphs[0]
    label_para.font.name = "맑은 고딕"
    label_para.font.size = FONT_HEADING
    label_para.font.bold = True
    label_para.font.color.rgb = COLOR_DARK_GRAY
    label_para.alignment = PP_ALIGN.CENTER
    label_frame.vertical_anchor = MSO_ANCHOR.TOP

    # Detail (optional)
    if detail:
        detail_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 1.1),
            Inches(width - 0.2), Inches(height - 1.2)
        )
        detail_frame = detail_box.text_frame
        detail_frame.text = detail
        detail_frame.word_wrap = True
        detail_para = detail_frame.paragraphs[0]
        detail_para.font.name = "맑은 고딕"
        detail_para.font.size = FONT_BODY
        detail_para.font.color.rgb = COLOR_MED_GRAY
        detail_para.alignment = PP_ALIGN.LEFT

def add_arrow_right(slide, x, y, width):
    """Add right-pointing arrow"""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(y),
        Inches(width), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_MED_GRAY
    arrow.line.color.rgb = COLOR_MED_GRAY

def add_arrow_down(slide, x, y, height):
    """Add down-pointing arrow"""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW,
        Inches(x), Inches(y),
        Inches(0.3), Inches(height)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_MED_GRAY
    arrow.line.color.rgb = COLOR_MED_GRAY

# ============================================================================
# SLIDE GENERATORS (Reusing from original)
# ============================================================================

def add_title_slide(prs):
    """Slide 0: Cover slide"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.83), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "자재군별 소싱 전략 및\n공급업체 관계 관리"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "맑은 고딕"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_BLACK
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8.83), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Strategic Inventory Management Course - Session 2"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = "Arial"
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = COLOR_MED_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Course info
    course_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8.83), Inches(0.5))
    course_frame = course_box.text_frame
    course_frame.text = "Kraljic Matrix Framework"
    course_para = course_frame.paragraphs[0]
    course_para.font.name = "Arial"
    course_para.font.size = Pt(16)
    course_para.font.color.rgb = COLOR_MED_GRAY
    course_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(8.83), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.name = "Arial"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = COLOR_MED_GRAY
    date_para.alignment = PP_ALIGN.CENTER

    return slide

def add_toc_slide(prs):
    """Slide 1: Table of Contents"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.0 목차 (Table of Contents)",
        "8개 장으로 구성된 자재군별 소싱 전략과 SRM 체계를 학습합니다."
    )

    toc_items = [
        "1장 소싱 전략 개요",
        "2장 병목자재 소싱",
        "3장 레버리지자재 소싱",
        "4장 전략자재 소싱",
        "5장 일상자재 소싱",
        "6장 SRM 및 성과 평가",
        "7장 Toyota 사례",
        "8장 Q&A 및 다음 회차"
    ]

    y_pos = 2.0
    for i, item in enumerate(toc_items, 1):
        # Chapter number box
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_pos),
            Inches(1.0), Inches(0.5)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        num_box.line.color.rgb = COLOR_LIGHT_GRAY

        num_text = num_box.text_frame
        num_text.text = f"{i}장"
        num_para = num_text.paragraphs[0]
        num_para.font.name = "맑은 고딕"
        num_para.font.size = FONT_HEADING
        num_para.font.bold = True
        num_para.font.color.rgb = COLOR_DARK_GRAY
        num_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Chapter title
        title_box = slide.shapes.add_textbox(
            Inches(2.7), Inches(y_pos),
            Inches(6.5), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = item
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "맑은 고딕"
        title_para.font.size = FONT_BULLET
        title_para.font.color.rgb = COLOR_DARK_GRAY
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        y_pos += 0.6

    return slide

def add_introduction_slide(prs):
    """Slide 2: Introduction"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.1 과정 개요 및 학습 목표",
        "자재군별 차별화된 소싱 전략으로 공급 리스크를 관리하고 최적의 가치를 창출합니다."
    )

    objectives = [
        "자재군별 차별화된 소싱 전략 수립",
        "SRM 접근법 이해",
        "계약 전략과 협상 포인트 파악",
        "공급업체 성과 평가 체계 구축"
    ]

    y_pos = 2.5
    for i, obj in enumerate(objectives, 1):
        obj_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_pos),
            Inches(7.5), Inches(0.8)
        )
        obj_box.fill.solid()
        obj_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        obj_box.line.color.rgb = COLOR_LIGHT_GRAY

        # Number circle
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.7), Inches(y_pos + 0.15),
            Inches(0.5), Inches(0.5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = COLOR_ACCENT
        num_shape.line.color.rgb = COLOR_ACCENT

        num_text = num_shape.text_frame
        num_text.text = str(i)
        num_para = num_text.paragraphs[0]
        num_para.font.name = "Arial"
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = COLOR_WHITE
        num_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Objective text
        obj_text_box = slide.shapes.add_textbox(
            Inches(2.4), Inches(y_pos + 0.2),
            Inches(6.4), Inches(0.4)
        )
        obj_text_frame = obj_text_box.text_frame
        obj_text_frame.text = obj
        obj_text_frame.word_wrap = True
        obj_text_para = obj_text_frame.paragraphs[0]
        obj_text_para.font.name = "맑은 고딕"
        obj_text_para.font.size = FONT_BULLET
        obj_text_para.font.color.rgb = COLOR_DARK_GRAY
        obj_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        y_pos += 1.0

    return slide

def add_sourcing_group_overview_slide(prs):
    """Slide 3: Sourcing Group Overview"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.2 소싱 그룹(Sourcing Group) 전략 개요",
        "비슷한 특성의 자재를 묶어 통합 관리하여 구매력 향상과 리스크 감소를 동시에 달성합니다."
    )

    # Left column
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.0),
        Inches(4.5), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    left_box.line.color.rgb = COLOR_LIGHT_GRAY

    left_title = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(4.1), Inches(0.4))
    left_title.text_frame.text = "소싱 그룹이란?"
    left_title_para = left_title.text_frame.paragraphs[0]
    left_title_para.font.name = "맑은 고딕"
    left_title_para.font.size = FONT_HEADING
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = COLOR_DARK_GRAY

    left_text = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(4.1), Inches(3.5))
    left_frame = left_text.text_frame
    left_frame.text = "비슷한 특성의 자재를 묶어 통합 관리하는 단위\n\n분류 기준:\n• Kraljic Matrix\n• 산업별\n• 기능별\n• 공급업체 유형별"
    left_frame.word_wrap = True
    left_para = left_frame.paragraphs[0]
    left_para.font.name = "맑은 고딕"
    left_para.font.size = FONT_BODY
    left_para.font.color.rgb = COLOR_DARK_GRAY
    left_para.line_spacing = 1.4

    # Right column
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(2.0),
        Inches(4.5), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    right_box.line.color.rgb = COLOR_LIGHT_GRAY

    right_title = slide.shapes.add_textbox(Inches(5.7), Inches(2.2), Inches(4.1), Inches(0.4))
    right_title.text_frame.text = "목적"
    right_title_para = right_title.text_frame.paragraphs[0]
    right_title_para.font.name = "맑은 고딕"
    right_title_para.font.size = FONT_HEADING
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = COLOR_DARK_GRAY

    right_text = slide.shapes.add_textbox(Inches(5.7), Inches(2.8), Inches(4.1), Inches(3.5))
    right_frame = right_text.text_frame
    right_frame.text = "• 구매력 향상\n• 효율성 증대\n• 리스크 감소\n• 체계적 관계 구축"
    right_frame.word_wrap = True
    right_para = right_frame.paragraphs[0]
    right_para.font.name = "맑은 고딕"
    right_para.font.size = FONT_BODY
    right_para.font.color.rgb = COLOR_DARK_GRAY
    right_para.line_spacing = 1.6

    return slide

def add_sourcing_strategy_matrix_slide(prs):
    """Slide 4: Sourcing Strategy Matrix (KEEP - already good)"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.3 자재군별 소싱 전략 매트릭스",
        "4개 자재군은 각각 다른 목표, 공급업체 수, 계약 기간, 관계 유형을 필요로 합니다."
    )

    # Matrix table (same as before - already 83 shapes!)
    col_widths = [1.8, 2.0, 2.0, 2.0, 2.0]
    row_height = 0.5
    start_x = 0.5
    start_y = 2.0

    headers = ["구분", "🔴 병목자재", "🟢 레버리지자재", "🟣 전략자재", "⚪ 일상자재"]
    criteria = ["핵심 목표", "소싱 전략", "공급업체 수", "계약 기간", "관계 유형", "협상 방식", "정보 공유"]

    data = [
        ["공급 안정성", "원가 경쟁력", "상호 성장", "효율성"],
        ["공급선 다변화", "경쟁 촉진", "전략적 파트너십", "통합 & 자동화"],
        ["2~3개", "5개 이상", "1~2개", "1~2개"],
        ["중장기", "단기", "장기", "중기"],
        ["협력적", "거래적", "파트너십", "효율적"],
        ["안정성 중심", "가격 경쟁", "Win-Win", "표준화"],
        ["중간", "제한적", "고도 공유", "최소화"]
    ]

    # Header row
    for col_idx, header in enumerate(headers):
        x = start_x + sum(col_widths[:col_idx])
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(start_y),
            Inches(col_widths[col_idx]), Inches(row_height)
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_GRAY
        cell.line.color.rgb = COLOR_WHITE
        cell.line.width = Pt(1)

        text_box = slide.shapes.add_textbox(
            Inches(x + 0.05), Inches(start_y + 0.05),
            Inches(col_widths[col_idx] - 0.1), Inches(row_height - 0.1)
        )
        text_frame = text_box.text_frame
        text_frame.text = header
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.font.name = "맑은 고딕"
        text_para.font.size = Pt(11)
        text_para.font.bold = True
        text_para.font.color.rgb = COLOR_WHITE
        text_para.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for row_idx, criterion in enumerate(criteria):
        y = start_y + (row_idx + 1) * row_height

        # Criterion column
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x), Inches(y),
            Inches(col_widths[0]), Inches(row_height)
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        cell.line.color.rgb = COLOR_WHITE
        cell.line.width = Pt(1)

        text_box = slide.shapes.add_textbox(
            Inches(start_x + 0.05), Inches(y + 0.05),
            Inches(col_widths[0] - 0.1), Inches(row_height - 0.1)
        )
        text_frame = text_box.text_frame
        text_frame.text = criterion
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.font.name = "맑은 고딕"
        text_para.font.size = FONT_BODY
        text_para.font.bold = True
        text_para.font.color.rgb = COLOR_DARK_GRAY
        text_para.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data columns
        for col_idx in range(4):
            x = start_x + sum(col_widths[:col_idx + 1])
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(col_widths[col_idx + 1]), Inches(row_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY if row_idx % 2 == 0 else COLOR_WHITE
            cell.line.color.rgb = COLOR_LIGHT_GRAY
            cell.line.width = Pt(0.5)

            text_box = slide.shapes.add_textbox(
                Inches(x + 0.05), Inches(y + 0.05),
                Inches(col_widths[col_idx + 1] - 0.1), Inches(row_height - 0.1)
            )
            text_frame = text_box.text_frame
            text_frame.text = data[row_idx][col_idx]
            text_frame.word_wrap = True
            text_para = text_frame.paragraphs[0]
            text_para.font.name = "맑은 고딕"
            text_para.font.size = Pt(9)
            text_para.font.color.rgb = COLOR_DARK_GRAY
            text_para.alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide

def add_bottleneck_strategy_enhanced(prs):

    print("[7-24] Remaining slides (simple format)...")
    
    # Slides 7-8: Bottleneck additional
    add_simple_slide_with_bullets(prs, 7, "2.2 병목자재: 공급선 다변화 전략",
        "단일 공급원 의존도를 낮춰 공급 리스크를 분산시킵니다.",
        ["• 메인 공급업체 + 백업 공급업체 체계", "• 지역적 분산 (다른 국가 공급업체)", "• 기술 이전 (신규 공급업체 육성)", "• 대체재 개발"])

    add_simple_slide_with_bullets(prs, 8, "2.3 병목자재: 계약 전략",
        "가격보다 공급 보장 조건을 우선하여 안정적 관계를 유지합니다.",
        ["계약: Long-term Agreement (LTA)", "• 공급 보증, 우선 공급권", "• 리드타임 단축, 긴급 대응", "• 적정 마진 보장 (Win-Win)"])

    print("[9/25] Leverage Bidding (ENHANCED)...")
    add_leverage_bidding_enhanced(prs)

    add_simple_slide_with_bullets(prs, 10, "3.1 레버리지자재: 경쟁 촉진 전략",
        "공급시장이 경쟁적이므로 공급업체 간 경쟁을 유도합니다.",
        ["• 경쟁 입찰: RFQ, 역경매", "• 통합 구매: 10개 → 3-5개로 축소", "• 글로벌 소싱: 저가 공급원 확보"])

    print("[11/25] TCO Comparison (ENHANCED)...")
    add_tco_comparison_enhanced(prs)

    print("[12/25] Partnership Diagram (ENHANCED)...")
    add_partnership_diagram_enhanced(prs)

    add_simple_slide_with_bullets(prs, 13, "4.1 전략자재: 파트너십 구축",
        "공급 리스크와 구매 임팩트가 모두 크므로 장기적 Win-Win 파트너십을 구축합니다.",
        ["• 3-5년 장기 계약", "• 목표 공유: 원가절감, 품질향상", "• 이익 공유: 50/50 분배", "• 공동 R&D 프로젝트"])

    add_simple_slide_with_bullets(prs, 14, "5.1 일상자재: 효율화 전략",
        "관리 비용을 최소화하는 것이 핵심입니다.",
        ["• 공급업체 통합: 1-2개로 집중", "• 카테고리 통합 구매", "• Blanket PO: 연간 총량 계약", "• VMI: 공급업체 재고 관리"])

    print("[15/25] E-Procurement (ENHANCED)...")
    add_eprocurement_enhanced(prs)

    add_simple_slide_with_bullets(prs, 16, "6.1 SRM 개요",
        "공급업체와의 관계를 체계적으로 관리하여 상호 가치를 극대화합니다.",
        ["• 공급업체 성과가 경쟁력에 직결", "• 장기 관계가 단기 가격보다 중요", "• 협력으로 혁신 창출", "• 리스크 관리와 지속가능성"])

    add_simple_slide_with_bullets(prs, 17, "6.2 자재군별 관계 유형",
        "병목(협력적), 레버리지(거래적), 전략(파트너십), 일상(효율적) - 각각 다른 접근이 필요합니다.",
        ["• 병목: 협력적, 월 1회 소통", "• 레버리지: 거래적, 분기 1회", "• 전략: 파트너십, 월 1-2회", "• 일상: 효율적, 분기 1회"])

    add_simple_slide_with_bullets(prs, 18, "6.3 Supplier Scorecard 구성",
        "품질 30%, 납기 30%, 가격 20%, 협력 10%, 리스크 10%로 정량적 평가를 수행합니다.",
        ["1. 품질 (30%): 불량률, 검사 통과율", "2. 납기 (30%): OTD, 리드타임", "3. 가격 (20%): 시장가 대비, 원가 절감", "4. 협력 (10%): 정보 공유, 개선 제안", "5. 리스크 (10%): 재무 건전성"])

    add_simple_slide_with_bullets(prs, 19, "6.4 등급 분류 및 조치",
        "A(90+), B(70-89), C(50-69), D(<50) 등급별로 차별화된 조치를 취합니다.",
        ["• A등급 (90+): 물량 확대 검토", "• B등급 (70-89): 현 수준 유지", "• C등급 (50-69): 개선 계획 요구", "• D등급 (<50): 교체 검토 또는 퇴출"])

    add_simple_slide_with_bullets(prs, 20, "7.1 Toyota SRM 사례",
        "Toyota는 공급업체들이 가장 협력하고 싶어하는 OEM 1위입니다.",
        ["• 1950년대부터 장기 파트너십", "• '운명 공동체'로 대우", "• 공급업체 만족도 1위", "• 2011년 대지진 시 빠른 회복"])

    print("[21/25] Toyota 3 Pillars (ENHANCED)...")
    add_toyota_pillars_enhanced(prs)

    add_simple_slide_with_bullets(prs, 22, "7.3 Toyota 성과 및 적용",
        "Win-Win 파트너십이 장기적 경쟁력을 만듭니다.",
        ["성과: 품질 세계 최고, 혁신 연간 수천 건", "• 병목: 장기 계약 + 기술 지원", "• 레버리지: 협력 통한 원가 절감", "• 전략: 경영진 레벨 정기 미팅"])

    add_simple_slide_with_bullets(prs, 23, "8.1 핵심 요약",
        "자재군별 차별화된 소싱 전략과 체계적 SRM으로 공급망 경쟁력을 강화합니다.",
        ["1. 병목: 공급 안정성", "2. 레버리지: 경쟁 촉진 & 원가 절감", "3. 전략: 파트너십", "4. 일상: 효율화", "5. SRM: 상호 가치 창출"])

    add_simple_slide_with_bullets(prs, 24, "8.2 Q&A 및 다음 회차",
        "3회차에서는 ABC-XYZ 재고 분류를 학습합니다.",
        ["Q&A 주제:", "• 파트너십 구축 시작 방법", "• 경쟁 입찰과 장기 관계의 균형", "\n다음 회차: ABC-XYZ 재고 분류", "• 금액 기준 ABC + 변동성 기준 XYZ", "• 9가지 조합별 운영 전략"])

    # Save
    output_path = "/home/user/Kraljic_Course/PPTX_RESULT/Part2_Session2_Sourcing_Strategy_Enhanced.pptx"
    print()
    print("Saving presentation...")
    prs.save(output_path)

    print()
    print("=" * 80)
    print(f"✅ ENHANCED PPTX GENERATED!")
    print("=" * 80)
    print()
    print("ENHANCEMENTS:")
    print("  ✓ Slide 5: Bottleneck Process Flow (visual-first)")
    print("  ✓ Slide 9: Leverage Bidding Flow (Toy Page layout)")
    print("  ✓ Slide 11: TCO Comparison (side-by-side)")
    print("  ✓ Slide 12: Partnership Diagram (relationship network)")
    print("  ✓ Slide 15: E-Procurement (vertical flow)")
    print("  ✓ Slide 21: Toyota 3 Pillars (visual pillars)")
    print()
    print("QUALITY IMPROVEMENTS:")
    print("  ✓ Consistent font sizes (10pt body, 12pt bullets, 16pt governing)")
    print("  ✓ Text overflow prevention (word_wrap + proper sizing)")
    print("  ✓ Reduced text content (60-70% visual, 30-40% text)")
    print("  ✓ All shapes editable in PowerPoint")
    print()
    print(f"Output: {output_path}")
    print()

    return output_path

    """Slide 5: ENHANCED - Bottleneck Strategy Process (Visual-first)"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "2.1 병목자재 소싱 전략: 공급 안정성 확보",
        "공급 불안정이 가장 큰 문제이므로, '언제든지 공급받을 수 있도록' 하는 것이 최우선 목표입니다."
    )

    # Process flow (4 steps with arrows)
    steps = [
        {"label": "공급선\n다변화", "detail": "Multi-\nSourcing"},
        {"label": "이중 공급\n체계", "detail": "Dual\nSourcing"},
        {"label": "장기 계약\n체결", "detail": "LTA"},
        {"label": "관계\n강화", "detail": "Collaboration"}
    ]

    box_width = 2.0
    box_height = 1.5
    gap = 0.4
    start_x = 0.8
    y = 2.5

    for i, step in enumerate(steps):
        x = start_x + i * (box_width + gap)
        add_process_step_box(slide, x, y, box_width, box_height, i + 1, step["label"], step["detail"])

        # Arrow (except last)
        if i < len(steps) - 1:
            arrow_x = x + box_width + 0.05
            add_arrow_right(slide, arrow_x, y + box_height / 2 - 0.15, gap - 0.1)

    # Bottom: Key insight box
    insight_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(4.5),
        Inches(9.0), Inches(1.5)
    )
    insight_box.fill.solid()
    insight_box.fill.fore_color.rgb = COLOR_ACCENT
    insight_box.line.color.rgb = COLOR_ACCENT

    insight_text = slide.shapes.add_textbox(Inches(1.2), Inches(4.8), Inches(8.2), Inches(0.9))
    insight_frame = insight_text.text_frame
    insight_frame.text = "💡 핵심: 가격보다 공급 보장을 우선하여, 적정 마진을 보장하며 안정적 관계를 유지합니다."
    insight_frame.word_wrap = True
    insight_para = insight_frame.paragraphs[0]
    insight_para.font.name = "맑은 고딕"
    insight_para.font.size = FONT_HEADING
    insight_para.font.bold = True
    insight_para.font.color.rgb = COLOR_WHITE
    insight_para.alignment = PP_ALIGN.CENTER
    insight_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    return slide

# Continue with remaining slides in next message due to length...
# This provides the pattern for enhanced visual design

def generate_part2_enhanced():
    """Generate enhanced Part 2 PPTX with visual focus"""
    print("=" * 80)
    print("GENERATING ENHANCED PART 2 PPTX - Visual-First Approach")
    print("=" * 80)
    print()

    prs = create_presentation()

    # Generate all slides
    print("[1/25] Cover slide...")
    add_title_slide(prs)

    print("[2/25] TOC slide...")
    add_toc_slide(prs)

    print("[3/25] Introduction...")
    add_introduction_slide(prs)

    print("[4/25] Sourcing Group Overview...")
    add_sourcing_group_overview_slide(prs)

    print("[5/25] Sourcing Strategy Matrix (83 shapes)...")
    add_sourcing_strategy_matrix_slide(prs)

    print("[6/25] Bottleneck Strategy (ENHANCED)...")
    add_bottleneck_strategy_enhanced(prs)


def add_simple_slide_with_bullets(prs, slide_num, title, gov_msg, bullets):
    """Simple slide with bullet points (reduced text, consistent fonts)"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(slide, title, gov_msg)

    # Bullet list box
    list_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.5), Inches(2.3),
        Inches(7.5), Inches(4.0)
    )
    list_box.fill.solid()
    list_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    list_box.line.color.rgb = COLOR_LIGHT_GRAY

    # Bullets
    list_text = slide.shapes.add_textbox(Inches(1.8), Inches(2.6), Inches(7.0), Inches(3.5))
    text_frame = list_text.text_frame
    text_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = bullet
        p.font.name = "맑은 고딕"
        p.font.size = FONT_BODY
        p.font.color.rgb = COLOR_DARK_GRAY
        p.line_spacing = 1.5
        p.level = 0

    return slide

def add_leverage_bidding_enhanced(prs):
    """Slide 9: ENHANCED - Leverage Competitive Bidding"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "3.2 레버리지자재: 경쟁 입찰 방식",
        "RFQ와 역경매를 통해 다수 공급업체 간 가격 경쟁을 유도합니다."
    )

    # TOY PAGE: Left 60% visual, Right 40% text
    
    # LEFT: Process flow
    process = ["RFQ\n발송", "경쟁\n입찰", "TCO\n분석", "선정"]
    box_width = 1.3
    box_height = 1.0
    gap = 0.25
    start_x = 0.8
    y = 2.5

    for i, step in enumerate(process):
        x = start_x + i * (box_width + gap)
        
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = COLOR_MED_GRAY
        box.line.color.rgb = COLOR_DARK_GRAY
        box.line.width = Pt(2)

        text_box = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.25), Inches(box_width - 0.2), Inches(0.5))
        text_frame = text_box.text_frame
        text_frame.text = step
        text_frame.word_wrap = True
        text_para = text_frame.paragraphs[0]
        text_para.font.name = "맑은 고딕"
        text_para.font.size = FONT_BULLET
        text_para.font.bold = True
        text_para.font.color.rgb = COLOR_WHITE
        text_para.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i < len(process) - 1:
            add_arrow_right(slide, x + box_width + 0.03, y + 0.35, gap - 0.06)

    # RIGHT: Key points sidebar
    sidebar_x = 6.5
    sidebar_y = 2.3
    
    sidebar_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(sidebar_x), Inches(sidebar_y),
        Inches(3.5), Inches(4.0)
    )
    sidebar_box.fill.solid()
    sidebar_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    sidebar_box.line.color.rgb = COLOR_LIGHT_GRAY

    sidebar_title = slide.shapes.add_textbox(Inches(sidebar_x + 0.2), Inches(sidebar_y + 0.2), Inches(3.1), Inches(0.4))
    sidebar_title.text_frame.text = "실행 방안"
    st_para = sidebar_title.text_frame.paragraphs[0]
    st_para.font.name = "맑은 고딕"
    st_para.font.size = FONT_HEADING
    st_para.font.bold = True
    st_para.font.color.rgb = COLOR_DARK_GRAY

    sidebar_text = slide.shapes.add_textbox(Inches(sidebar_x + 0.2), Inches(sidebar_y + 0.7), Inches(3.1), Inches(3.0))
    sidebar_frame = sidebar_text.text_frame
    sidebar_frame.text = "• 표준화된 견적 요청\n• 5개 이상 공급업체\n• 온라인 역경매 활용\n• 물량 통합으로 협상력 강화"
    sidebar_frame.word_wrap = True
    s_para = sidebar_frame.paragraphs[0]
    s_para.font.name = "맑은 고딕"
    s_para.font.size = FONT_BODY
    s_para.font.color.rgb = COLOR_DARK_GRAY
    s_para.line_spacing = 1.6

    return slide

def add_tco_comparison_enhanced(prs):
    """Slide 11: ENHANCED - TCO Comparison"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "3.3 레버리지자재: TCO 기반 공급업체 선정",
        "단가로 좁히고 TCO로 결정한다 - 최저가가 아닌 총소유비용으로 최종 평가합니다."
    )

    # Formula box
    formula_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(2.0),
        Inches(8.5), Inches(0.6)
    )
    formula_box.fill.solid()
    formula_box.fill.fore_color.rgb = COLOR_ACCENT
    formula_box.line.color.rgb = COLOR_ACCENT

    formula_text = slide.shapes.add_textbox(Inches(1.2), Inches(2.15), Inches(8.1), Inches(0.3))
    formula_frame = formula_text.text_frame
    formula_frame.text = "TCO = 구매가 + 물류비 + 관세 + 품질비용 + 재고비용 + 관리비용"
    f_para = formula_frame.paragraphs[0]
    f_para.font.name = "맑은 고딕"
    f_para.font.size = FONT_HEADING
    f_para.font.bold = True
    f_para.font.color.rgb = COLOR_WHITE
    f_para.alignment = PP_ALIGN.CENTER

    # Comparison boxes
    # Left: Domestic
    domestic_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(3.0),
        Inches(4.0), Inches(3.0)
    )
    domestic_box.fill.solid()
    domestic_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    domestic_box.line.color.rgb = COLOR_MED_GRAY
    domestic_box.line.width = Pt(2)

    d_title = slide.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(3.6), Inches(0.4))
    d_title.text_frame.text = "국내 공급업체"
    dt_para = d_title.text_frame.paragraphs[0]
    dt_para.font.name = "맑은 고딕"
    dt_para.font.size = FONT_HEADING
    dt_para.font.bold = True
    dt_para.font.color.rgb = COLOR_BLACK
    dt_para.alignment = PP_ALIGN.CENTER

    d_content = slide.shapes.add_textbox(Inches(1.3), Inches(3.7), Inches(3.4), Inches(2.0))
    d_frame = d_content.text_frame
    d_frame.text = "구매가: ₩100\n물류비: ₩5\n관세: ₩0\n품질비용: ₩2\n재고비용: ₩3\n\n총 TCO: ₩110"
    d_frame.word_wrap = True
    dp = d_frame.paragraphs[0]
    dp.font.name = "맑은 고딕"
    dp.font.size = FONT_BODY
    dp.font.color.rgb = COLOR_DARK_GRAY
    dp.line_spacing = 1.4

    # Right: Overseas
    overseas_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(3.0),
        Inches(4.0), Inches(3.0)
    )
    overseas_box.fill.solid()
    overseas_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    overseas_box.line.color.rgb = COLOR_MED_GRAY
    overseas_box.line.width = Pt(2)

    o_title = slide.shapes.add_textbox(Inches(5.7), Inches(3.2), Inches(3.6), Inches(0.4))
    o_title.text_frame.text = "해외 공급업체"
    ot_para = o_title.text_frame.paragraphs[0]
    ot_para.font.name = "맑은 고딕"
    ot_para.font.size = FONT_HEADING
    ot_para.font.bold = True
    ot_para.font.color.rgb = COLOR_BLACK
    ot_para.alignment = PP_ALIGN.CENTER

    o_content = slide.shapes.add_textbox(Inches(5.8), Inches(3.7), Inches(3.4), Inches(2.0))
    o_frame = o_content.text_frame
    o_frame.text = "구매가: ₩85\n물류비: ₩15\n관세: ₩8\n품질비용: ₩5\n재고비용: ₩8\n\n총 TCO: ₩121"
    o_frame.word_wrap = True
    op = o_frame.paragraphs[0]
    op.font.name = "맑은 고딕"
    op.font.size = FONT_BODY
    op.font.color.rgb = COLOR_DARK_GRAY
    op.line_spacing = 1.4

    # Winner indicator
    winner_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.5), Inches(6.3),
        Inches(5.5), Inches(0.6)
    )
    winner_box.fill.solid()
    winner_box.fill.fore_color.rgb = RGBColor(39, 174, 96)  # Green
    winner_box.line.color.rgb = RGBColor(39, 174, 96)

    winner_text = slide.shapes.add_textbox(Inches(2.7), Inches(6.45), Inches(5.1), Inches(0.3))
    winner_frame = winner_text.text_frame
    winner_frame.text = "✓ 국내 공급업체 선정 (TCO 우위: ₩11 절감)"
    w_para = winner_frame.paragraphs[0]
    w_para.font.name = "맑은 고딕"
    w_para.font.size = FONT_HEADING
    w_para.font.bold = True
    w_para.font.color.rgb = COLOR_WHITE
    w_para.alignment = PP_ALIGN.CENTER

    return slide

def add_partnership_diagram_enhanced(prs):
    """Slide 12: ENHANCED - Strategic Partnership"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "4.2 전략자재: 장기 파트너십 계약",
        "3-5년 장기 계약으로 목표와 이익을 공유하고, 공동 R&D 프로젝트를 진행합니다."
    )

    # Center: Partnership
    center_x, center_y = 5.0, 4.0
    center_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(center_x - 1.0), Inches(center_y - 0.5),
        Inches(2.0), Inches(1.0)
    )
    center_box.fill.solid()
    center_box.fill.fore_color.rgb = COLOR_ACCENT
    center_box.line.color.rgb = COLOR_ACCENT

    center_text = slide.shapes.add_textbox(Inches(center_x - 0.9), Inches(center_y - 0.3), Inches(1.8), Inches(0.6))
    center_frame = center_text.text_frame
    center_frame.text = "전략적\n파트너십"
    center_frame.word_wrap = True
    c_para = center_frame.paragraphs[0]
    c_para.font.name = "맑은 고딕"
    c_para.font.size = FONT_HEADING
    c_para.font.bold = True
    c_para.font.color.rgb = COLOR_WHITE
    c_para.alignment = PP_ALIGN.CENTER
    center_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Three pillars
    pillars = [
        {"x": 1.5, "y": 2.5, "title": "목표 공유", "items": ["원가절감", "품질향상", "기술혁신"]},
        {"x": 7.0, "y": 2.5, "title": "이익 공유", "items": ["절감액", "50/50 분배"]},
        {"x": 4.0, "y": 5.8, "title": "리스크 공유", "items": ["가격변동", "공동대응"]}
    ]

    for pillar in pillars:
        # Pillar box
        p_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(pillar["x"]), Inches(pillar["y"]),
            Inches(1.8), Inches(0.8)
        )
        p_box.fill.solid()
        p_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        p_box.line.color.rgb = COLOR_MED_GRAY
        p_box.line.width = Pt(2)

        p_title = slide.shapes.add_textbox(Inches(pillar["x"] + 0.1), Inches(pillar["y"] + 0.1), Inches(1.6), Inches(0.3))
        p_title.text_frame.text = pillar["title"]
        pt_para = p_title.text_frame.paragraphs[0]
        pt_para.font.name = "맑은 고딕"
        pt_para.font.size = FONT_BULLET
        pt_para.font.bold = True
        pt_para.font.color.rgb = COLOR_BLACK
        pt_para.alignment = PP_ALIGN.CENTER

        p_items = slide.shapes.add_textbox(Inches(pillar["x"] + 0.2), Inches(pillar["y"] + 0.45), Inches(1.4), Inches(0.3))
        p_items.text_frame.text = "\n".join(pillar["items"])
        p_items.text_frame.word_wrap = True
        pi_para = p_items.text_frame.paragraphs[0]
        pi_para.font.name = "맑은 고딕"
        pi_para.font.size = FONT_CAPTION
        pi_para.font.color.rgb = COLOR_MED_GRAY
        pi_para.alignment = PP_ALIGN.CENTER
        pi_para.line_spacing = 1.3

        # Lines to center
        slide.shapes.add_connector(
            1,  # MSO_CONNECTOR_TYPE.STRAIGHT
            Inches(pillar["x"] + 0.9), Inches(pillar["y"] + 0.8 if pillar["y"] < center_y else pillar["y"]),
            Inches(center_x), Inches(center_y)
        ).line.color.rgb = COLOR_MED_GRAY

    return slide

def add_eprocurement_enhanced(prs):
    """Slide 15: ENHANCED - E-Procurement Architecture"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "5.2 일상자재: E-Procurement 및 자동화",
        "카탈로그 구매와 자동 발주 시스템으로 승인 프로세스를 간소화합니다."
    )

    # Vertical flow
    layers = [
        {"label": "카탈로그 구매", "detail": "사전 등록 품목 선택"},
        {"label": "자동 발주", "detail": "재고 부족 시 자동"},
        {"label": "승인 자동화", "detail": "일정 금액 이하"},
        {"label": "3-Way Matching", "detail": "PO-GR-IR 자동"}
    ]

    box_width = 6.0
    box_height = 0.9
    start_x = 2.4
    start_y = 2.2
    gap = 0.3

    for i, layer in enumerate(layers):
        y = start_y + i * (box_height + gap)

        layer_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(start_x), Inches(y),
            Inches(box_width), Inches(box_height)
        )
        layer_box.fill.solid()
        layer_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        layer_box.line.color.rgb = COLOR_MED_GRAY
        layer_box.line.width = Pt(2)

        # Number
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(start_x + 0.2), Inches(y + 0.2),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLOR_ACCENT
        num_circle.line.color.rgb = COLOR_ACCENT

        num_text = num_circle.text_frame
        num_text.text = str(i + 1)
        n_para = num_text.paragraphs[0]
        n_para.font.name = "Arial"
        n_para.font.size = Pt(16)
        n_para.font.bold = True
        n_para.font.color.rgb = COLOR_WHITE
        n_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Label
        label_text = slide.shapes.add_textbox(Inches(start_x + 0.9), Inches(y + 0.15), Inches(2.5), Inches(0.3))
        label_text.text_frame.text = layer["label"]
        l_para = label_text.text_frame.paragraphs[0]
        l_para.font.name = "맑은 고딕"
        l_para.font.size = FONT_HEADING
        l_para.font.bold = True
        l_para.font.color.rgb = COLOR_BLACK

        # Detail
        detail_text = slide.shapes.add_textbox(Inches(start_x + 0.9), Inches(y + 0.5), Inches(5.0), Inches(0.3))
        detail_text.text_frame.text = layer["detail"]
        detail_text.text_frame.word_wrap = True
        d_para = detail_text.text_frame.paragraphs[0]
        d_para.font.name = "맑은 고딕"
        d_para.font.size = FONT_BODY
        d_para.font.color.rgb = COLOR_MED_GRAY

        # Arrow down (except last)
        if i < len(layers) - 1:
            add_arrow_down(slide, start_x + box_width / 2 - 0.15, y + box_height + 0.05, gap - 0.1)

    return slide

def add_toyota_pillars_enhanced(prs):
    """Slide 21: ENHANCED - Toyota 3 Pillars"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "7.2 Toyota SRM: 3가지 핵심 전략",
        "상호 신뢰 기반 장기 파트너십, Kaizen 철학 확산, 성장 비전 공유로 공급망 전체 경쟁력을 향상시킵니다."
    )

    pillars = [
        {
            "x": 0.8,
            "title": "상호 신뢰\n파트너십",
            "items": ["장기 계약", "투명한 정보", "공정한 가격"]
        },
        {
            "x": 4.0,
            "title": "Kaizen\n지속적 개선",
            "items": ["교육 지원", "현장 지원", "공동 해결"]
        },
        {
            "x": 7.2,
            "title": "성장 비전\n공유",
            "items": ["장기 예측", "투자 지원", "공동 R&D"]
        }
    ]

    for i, pillar in enumerate(pillars):
        # Number
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(pillar["x"] + 0.9), Inches(2.3),
            Inches(0.5), Inches(0.5)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLOR_ACCENT
        num_circle.line.color.rgb = COLOR_ACCENT

        num_text = num_circle.text_frame
        num_text.text = str(i + 1)
        n_para = num_text.paragraphs[0]
        n_para.font.name = "Arial"
        n_para.font.size = Pt(20)
        n_para.font.bold = True
        n_para.font.color.rgb = COLOR_WHITE
        n_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Title box
        title_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(pillar["x"]), Inches(3.0),
            Inches(2.5), Inches(1.0)
        )
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = COLOR_MED_GRAY
        title_box.line.color.rgb = COLOR_DARK_GRAY
        title_box.line.width = Pt(2)

        title_text = slide.shapes.add_textbox(Inches(pillar["x"] + 0.2), Inches(3.15), Inches(2.1), Inches(0.7))
        title_frame = title_text.text_frame
        title_frame.text = pillar["title"]
        title_frame.word_wrap = True
        t_para = title_frame.paragraphs[0]
        t_para.font.name = "맑은 고딕"
        t_para.font.size = FONT_HEADING
        t_para.font.bold = True
        t_para.font.color.rgb = COLOR_WHITE
        t_para.alignment = PP_ALIGN.CENTER
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Items
        item_y = 4.3
        for item in pillar["items"]:
            item_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(pillar["x"] + 0.2), Inches(item_y),
                Inches(2.1), Inches(0.5)
            )
            item_box.fill.solid()
            item_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
            item_box.line.color.rgb = COLOR_LIGHT_GRAY

            item_text = slide.shapes.add_textbox(Inches(pillar["x"] + 0.3), Inches(item_y + 0.1), Inches(1.9), Inches(0.3))
            item_frame = item_text.text_frame
            item_frame.text = item
            item_frame.word_wrap = True
            i_para = item_frame.paragraphs[0]
            i_para.font.name = "맑은 고딕"
            i_para.font.size = FONT_BODY
            i_para.font.color.rgb = COLOR_DARK_GRAY
            i_para.alignment = PP_ALIGN.CENTER
            item_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            item_y += 0.65

    return slide



    # For now, keeping original simple slides for remaining content

    # Save
    output_path = "/home/user/Kraljic_Course/PPTX_RESULT/Part2_Session2_Sourcing_Strategy_Enhanced.pptx"
    print()
    print("Saving presentation...")
    prs.save(output_path)

    print()
    print("=" * 80)
    print(f"✅ Enhanced PPTX generated: {output_path}")
    print("=" * 80)

    return output_path

if __name__ == "__main__":
    try:
        generate_part2_enhanced()
    except Exception as e:
        print(f"\n❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
        exit(1)
