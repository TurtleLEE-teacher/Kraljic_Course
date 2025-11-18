#!/usr/bin/env python3
"""
Part 1 PPTX Generation (Sessions 1-3) - 20 Slides
S4HANA Monochrome Design System

STRICT COMPLIANCE:
- Monochrome colors (Black/White/Gray + 1 accent max)
- Governing messages on ALL content slides (16pt Bold)
- Arial (English) + 맑은 고딕 (Korean)
- Size: 10.83" x 7.5"
- White background
- Simple and clean layout
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# S4HANA Monochrome Color System
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)
COLOR_MED_GRAY = RGBColor(102, 102, 102)
COLOR_LIGHT_GRAY = RGBColor(204, 204, 204)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(26, 82, 118)  # Single accent color (dark blue)

# Kraljic Matrix colors (ONLY for matrix diagram)
COLOR_STRATEGIC = RGBColor(142, 68, 173)
COLOR_BOTTLENECK = RGBColor(230, 126, 34)
COLOR_LEVERAGE = RGBColor(39, 174, 96)
COLOR_ROUTINE = RGBColor(149, 165, 166)

def create_presentation():
    """Create Part 1 PPTX with 20 slides"""
    prs = Presentation()
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    add_cover_slide(prs)

    # Session 1: Kraljic Matrix Foundation (Slides 2-7)
    add_session1_slides(prs)

    # Session 2: Sourcing Strategies (Slides 8-14)
    add_session2_slides(prs)

    # Session 3: ABC-XYZ (Slides 15-20)
    add_session3_slides(prs)

    return prs


def add_standard_title(slide, title_text, governing_message=None, slide_num=None):
    """Add standard title with optional governing message (S4HANA style)"""
    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31), Inches(9.5), Inches(0.43)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_BLACK

    # Governing message (REQUIRED for content slides)
    if governing_message:
        gov_box = slide.shapes.add_textbox(
            Inches(0.3), Inches(1.01), Inches(10.32), Inches(0.63)
        )
        gtf = gov_box.text_frame
        gtf.word_wrap = True
        gp = gtf.paragraphs[0]
        grun = gp.add_run()
        grun.text = governing_message
        grun.font.name = "맑은 고딕"
        grun.font.size = Pt(16)
        grun.font.bold = True
        grun.font.color.rgb = COLOR_DARK_GRAY

    # Slide number
    if slide_num:
        num_box = slide.shapes.add_textbox(
            Inches(10.0), Inches(7.15), Inches(0.5), Inches(0.25)
        )
        nf = num_box.text_frame
        np = nf.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        nrun = np.add_run()
        nrun.text = str(slide_num)
        nrun.font.name = "Arial"
        nrun.font.size = Pt(12)
        nrun.font.color.rgb = COLOR_MED_GRAY


def add_bullet_list(slide, left, top, width, height, items, font_size=16):
    """Add bullet list (using • symbol, S4HANA style)"""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True

    for item in items:
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        p.level = 0
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(font_size)
        run.font.color.rgb = COLOR_DARK_GRAY


def add_cover_slide(prs):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Simple white background with accent color title box
    title_shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.45), Inches(2), Inches(9.93), Inches(2)
    )
    title_shape.fill.solid()
    title_shape.fill.fore_color.rgb = COLOR_ACCENT
    title_shape.line.color.rgb = COLOR_ACCENT

    tf = title_shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "전략적 재고운영 및 자재계획수립\nPart 1: Foundation"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(4.5), Inches(9.93), Inches(1)
    )
    stf = subtitle_box.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srun = sp.add_run()
    srun.text = "Sessions 1-3: Kraljic Matrix, 소싱 전략, ABC-XYZ 분류"
    srun.font.name = "맑은 고딕"
    srun.font.size = Pt(20)
    srun.font.color.rgb = COLOR_DARK_GRAY

    # Date
    date_box = slide.shapes.add_textbox(
        Inches(0.45), Inches(6), Inches(9.93), Inches(0.5)
    )
    dtf = date_box.text_frame
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    drun = dp.add_run()
    drun.text = "2025년"
    drun.font.name = "Arial"
    drun.font.size = Pt(16)
    drun.font.color.rgb = COLOR_MED_GRAY


def add_session1_slides(prs):
    """Session 1: Kraljic Matrix Foundation (Slides 2-7)"""

    # Slide 2: Learning Objectives & Agenda
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide2, "Session 1: 학습 목표 & Agenda",
                      "JIT에서 JIC로의 패러다임 전환과 Kraljic Matrix를 활용한 자재 분류 방법을 학습합니다.",
                      2)

    content_box = slide2.shapes.add_textbox(
        Inches(0.7), Inches(2), Inches(9.5), Inches(4.5)
    )
    tf = content_box.text_frame
    p = tf.paragraphs[0]
    p.space_after = Pt(20)
    run = p.add_run()
    run.text = "학습 목표"
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = COLOR_BLACK

    objectives = [
        "JIT에서 JIC로의 패러다임 전환 배경과 필요성 이해",
        "전략적 재고운영의 핵심 개념과 차별화된 접근법 습득",
        "Kraljic Matrix를 활용한 자재 포트폴리오 분류 역량 확보",
        "자재군별 관리 철학과 자재계획 방법론의 전체 맵 이해"
    ]
    add_bullet_list(slide2, Inches(0.7), Inches(2.8), Inches(9.5), Inches(3.5), objectives)

    # Slide 3: JIT vs JIC
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide3, "JIT vs JIC: 패러다임의 전환",
                      "2020년 팬데믹은 JIT의 취약성을 드러냈고, 기업들은 회복력 중심의 JIC로 전환하고 있습니다.",
                      3)

    # Left column - JIT
    jit_items = [
        "재고 = 낭비 (Waste)",
        "목표: 재고 Zero",
        "우선순위: 효율성 (Efficiency)",
        "안전재고: 최소 (1-2주)",
        "공급업체: Single Source"
    ]

    left_title = slide3.shapes.add_textbox(Inches(0.7), Inches(2), Inches(4.5), Inches(0.4))
    lt = left_title.text_frame.paragraphs[0]
    ltr = lt.add_run()
    ltr.text = "JIT (과거)"
    ltr.font.name = "Arial"
    ltr.font.size = Pt(18)
    ltr.font.bold = True
    ltr.font.color.rgb = COLOR_BLACK

    add_bullet_list(slide3, Inches(0.7), Inches(2.6), Inches(4.5), Inches(4), jit_items, 15)

    # Right column - JIC
    jic_items = [
        "재고 = 전략적 자산",
        "목표: 최적 재고 (Optimal)",
        "우선순위: 회복력 (Resilience)",
        "안전재고: 차별화 (1주-6개월)",
        "공급업체: Dual/Multi Source"
    ]

    right_title = slide3.shapes.add_textbox(Inches(5.6), Inches(2), Inches(4.5), Inches(0.4))
    rt = right_title.text_frame.paragraphs[0]
    rtr = rt.add_run()
    rtr.text = "JIC (현재/미래)"
    rtr.font.name = "Arial"
    rtr.font.size = Pt(18)
    rtr.font.bold = True
    rtr.font.color.rgb = COLOR_ACCENT

    add_bullet_list(slide3, Inches(5.6), Inches(2.6), Inches(4.5), Inches(4), jic_items, 15)

    # Slide 4: JIT Crisis Cases
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide4, "JIT 붕괴 사례: 2020-2021 공급망 위기",
                      "글로벌 팬데믹은 Just-In-Time 방식의 치명적 약점을 극명하게 드러냈습니다.",
                      4)

    crisis_cases = [
        "2021년 글로벌 반도체 대란: 자동차 생산 1,000만 대 감소, 산업 손실 $210억",
        "마스크 대란 (2020): 전 세계 생산의 50%를 중국에 의존, 봉쇄로 공급 중단",
        "컨테이너선 대란 (2021-2022): 운임비 10배 폭등, 리드타임 2배 증가",
        "GM: 반도체 부족으로 생산 차질 200만 대 → 안전재고 0주에서 10-12주로 확대",
        "Toyota: JIT 창시자조차 반도체 안전재고 4-6개월치 확보로 전략 전환"
    ]
    add_bullet_list(slide4, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), crisis_cases, 16)

    # Slide 5: Kraljic Matrix Introduction
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide5, "Kraljic Matrix의 탄생과 핵심 통찰",
                      "1983년 Peter Kraljic은 모든 자재가 동등하지 않다는 혁신적 통찰을 제시했습니다.",
                      5)

    intro_box = slide5.shapes.add_textbox(Inches(0.7), Inches(2), Inches(9.5), Inches(0.6))
    itf = intro_box.text_frame
    ip = itf.paragraphs[0]
    irun = ip.add_run()
    irun.text = 'Peter Kraljic (1983, Harvard Business Review): "Purchasing Must Become Supply Management"'
    irun.font.name = "맑은 고딕"
    irun.font.size = Pt(14)
    irun.font.italic = True
    irun.font.color.rgb = COLOR_MED_GRAY

    kraljic_intro = [
        "탄생 배경: 1970년대 2차례 석유파동으로 공급 불안정 극심",
        "핵심 통찰: 'Not all materials are created equal' - 모든 자재가 동등하지 않다",
        "차별화 필요성: 자재의 특성에 따라 차별화된 관리 전략 적용",
        "2축 평가: 공급 리스크 (Supply Risk) × 구매 임팩트 (Purchase Impact)",
        "4대 자재군: 전략/레버리지/병목/일상 자재로 분류"
    ]
    add_bullet_list(slide5, Inches(0.7), Inches(3), Inches(9.5), Inches(3.5), kraljic_intro, 16)

    # Slide 6: Kraljic Matrix 2 Axes
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide6, "Kraljic Matrix의 두 축",
                      "공급 리스크와 구매 임팩트 두 차원으로 자재를 평가하여 4개 군으로 분류합니다.",
                      6)

    # Y-axis
    y_title = slide6.shapes.add_textbox(Inches(0.7), Inches(2), Inches(4.5), Inches(0.4))
    yt = y_title.text_frame.paragraphs[0]
    ytr = yt.add_run()
    ytr.text = "Y축: 공급 리스크 (Supply Risk)"
    ytr.font.name = "Arial"
    ytr.font.size = Pt(16)
    ytr.font.bold = True
    ytr.font.color.rgb = COLOR_BLACK

    y_items = [
        "공급업체 수 (1-2개 고위험)",
        "대체 가능성 (대체 불가 고위험)",
        "시장 구조 (독과점 고위험)",
        "지리적 집중도 (특정 지역 집중 고위험)",
        "리드타임 (6개월 이상 고위험)"
    ]
    add_bullet_list(slide6, Inches(0.7), Inches(2.6), Inches(4.5), Inches(4), y_items, 14)

    # X-axis
    x_title = slide6.shapes.add_textbox(Inches(5.6), Inches(2), Inches(4.5), Inches(0.4))
    xt = x_title.text_frame.paragraphs[0]
    xtr = xt.add_run()
    xtr.text = "X축: 구매 임팩트 (Purchase Impact)"
    xtr.font.name = "Arial"
    xtr.font.size = Pt(16)
    xtr.font.bold = True
    xtr.font.color.rgb = COLOR_ACCENT

    x_items = [
        "구매 금액 (총 구매액 대비 비중)",
        "원가 비중 (제품 원가 중 차지 비율)",
        "사업 영향도 (결품 시 생산 중단)",
        "부가가치 (최종 제품 성능 영향)",
        "품질 중요성 (품질 문제 파급 효과)"
    ]
    add_bullet_list(slide6, Inches(5.6), Inches(2.6), Inches(4.5), Inches(4), x_items, 14)

    # Slide 7: Kraljic Matrix 2x2
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide7, "Kraljic Matrix: 4대 자재군 분류",
                      "자재를 공급 리스크와 구매 임팩트 두 축으로 분류하여 차별화된 전략을 수립합니다.",
                      7)

    # 2x2 Matrix
    box_w = Inches(4.2)
    box_h = Inches(1.8)
    gap = Inches(0.2)
    matrix_left = Inches(1.2)
    matrix_top = Inches(2.5)

    # Top-left: Bottleneck
    bn_box = slide7.shapes.add_shape(1, matrix_left, matrix_top, box_w, box_h)
    bn_box.fill.solid()
    bn_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    bn_box.line.color.rgb = COLOR_MED_GRAY
    bn_box.line.width = Pt(2)

    bn_tf = bn_box.text_frame
    bn_tf.vertical_anchor = MSO_ANCHOR.TOP
    bn_p = bn_tf.paragraphs[0]
    bn_p.alignment = PP_ALIGN.LEFT
    bn_p.space_before = Pt(10)
    bn_run = bn_p.add_run()
    bn_run.text = "병목자재\nBottleneck\n\n독점 공급, 희토류\n특수 소재"
    bn_run.font.name = "맑은 고딕"
    bn_run.font.size = Pt(14)
    bn_run.font.bold = True
    bn_run.font.color.rgb = COLOR_DARK_GRAY

    # Top-right: Strategic
    st_box = slide7.shapes.add_shape(1, matrix_left + box_w + gap, matrix_top, box_w, box_h)
    st_box.fill.solid()
    st_box.fill.fore_color.rgb = COLOR_DARK_GRAY
    st_box.line.color.rgb = COLOR_BLACK
    st_box.line.width = Pt(2)

    st_tf = st_box.text_frame
    st_tf.vertical_anchor = MSO_ANCHOR.TOP
    st_p = st_tf.paragraphs[0]
    st_p.alignment = PP_ALIGN.LEFT
    st_p.space_before = Pt(10)
    st_run = st_p.add_run()
    st_run.text = "전략자재\nStrategic\n\n핵심 부품, 고가 장비\n전략 원자재"
    st_run.font.name = "맑은 고딕"
    st_run.font.size = Pt(14)
    st_run.font.bold = True
    st_run.font.color.rgb = COLOR_WHITE

    # Bottom-left: Routine
    rt_box = slide7.shapes.add_shape(1, matrix_left, matrix_top + box_h + gap, box_w, box_h)
    rt_box.fill.solid()
    rt_box.fill.fore_color.rgb = COLOR_WHITE
    rt_box.line.color.rgb = COLOR_MED_GRAY
    rt_box.line.width = Pt(2)

    rt_tf = rt_box.text_frame
    rt_tf.vertical_anchor = MSO_ANCHOR.TOP
    rt_p = rt_tf.paragraphs[0]
    rt_p.alignment = PP_ALIGN.LEFT
    rt_p.space_before = Pt(10)
    rt_run = rt_p.add_run()
    rt_run.text = "일상자재\nRoutine\n\n사무용품, 소모품\n저가 부품"
    rt_run.font.name = "맑은 고딕"
    rt_run.font.size = Pt(14)
    rt_run.font.bold = True
    rt_run.font.color.rgb = COLOR_DARK_GRAY

    # Bottom-right: Leverage
    lv_box = slide7.shapes.add_shape(1, matrix_left + box_w + gap, matrix_top + box_h + gap, box_w, box_h)
    lv_box.fill.solid()
    lv_box.fill.fore_color.rgb = COLOR_LIGHT_GRAY
    lv_box.line.color.rgb = COLOR_MED_GRAY
    lv_box.line.width = Pt(2)

    lv_tf = lv_box.text_frame
    lv_tf.vertical_anchor = MSO_ANCHOR.TOP
    lv_p = lv_tf.paragraphs[0]
    lv_p.alignment = PP_ALIGN.LEFT
    lv_p.space_before = Pt(10)
    lv_run = lv_p.add_run()
    lv_run.text = "레버리지자재\nLeverage\n\n범용 부품, 표준 원자재\nMRO 물품"
    lv_run.font.name = "맑은 고딕"
    lv_run.font.size = Pt(14)
    lv_run.font.bold = True
    lv_run.font.color.rgb = COLOR_DARK_GRAY

    # Axis labels
    y_label = slide7.shapes.add_textbox(Inches(0.2), Inches(4), Inches(0.8), Inches(1))
    y_p = y_label.text_frame.paragraphs[0]
    y_p.alignment = PP_ALIGN.CENTER
    y_r = y_p.add_run()
    y_r.text = "공급\n리스크\n↑"
    y_r.font.name = "맑은 고딕"
    y_r.font.size = Pt(12)
    y_r.font.bold = True
    y_r.font.color.rgb = COLOR_BLACK

    x_label = slide7.shapes.add_textbox(Inches(4), Inches(6.8), Inches(3), Inches(0.4))
    x_p = x_label.text_frame.paragraphs[0]
    x_p.alignment = PP_ALIGN.CENTER
    x_r = x_p.add_run()
    x_r.text = "구매 임팩트 (Purchase Impact) →"
    x_r.font.name = "맑은 고딕"
    x_r.font.size = Pt(12)
    x_r.font.bold = True
    x_r.font.color.rgb = COLOR_ACCENT


def add_session2_slides(prs):
    """Session 2: Sourcing Strategies (Slides 8-14)"""

    # Slide 8: Session 2 Intro
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide8, "Session 2: 자재군별 소싱 전략",
                      "각 자재군의 특성에 맞는 차별화된 소싱 전략과 공급업체 관리 방법을 학습합니다.",
                      8)

    session2_agenda = [
        "병목자재 소싱 전략: 공급 안정성 확보",
        "레버리지자재 소싱 전략: 경쟁 촉진 및 통합 구매",
        "전략자재 소싱 전략: 전략적 파트너십 구축",
        "일상자재 소싱 전략: 효율화 및 자동화",
        "공급업체 관계 관리 (SRM) 프레임워크"
    ]
    add_bullet_list(slide8, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4), session2_agenda, 17)

    # Slide 9: Bottleneck Sourcing
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide9, "병목자재 소싱 전략",
                      "공급 리스크가 높은 병목자재는 공급 연속성 확보와 대체재 개발이 핵심입니다.",
                      9)

    bn_strategies = [
        "핵심 전략: 공급 안정성 확보 (Supply Continuity)",
        "장기 계약: 2-3년 장기 공급 계약으로 안정성 확보",
        "안전재고: 3-6개월치 안전재고 보유 (ROP 방법론)",
        "공급원 다변화: Dual/Triple Sourcing 추진",
        "대체재 개발: R&D 투자로 기술 종속성 탈피",
        "리스크 모니터링: 공급업체 재무 상태 및 지정학적 리스크 상시 모니터링"
    ]
    add_bullet_list(slide9, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), bn_strategies, 15)

    # Slide 10: Leverage Sourcing
    slide10 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide10, "레버리지자재 소싱 전략",
                      "공급 리스크가 낮은 레버리지자재는 경쟁 입찰과 물량 레버리지로 원가 절감을 추구합니다.",
                      10)

    lv_strategies = [
        "핵심 전략: 경쟁 촉진 및 통합 구매 (Competitive Bidding)",
        "경쟁 입찰: RFQ (Request for Quotation) 정기 실시",
        "물량 통합: 여러 사업부 수요 통합으로 협상력 강화",
        "다수 공급업체: 3-5개 공급업체 경쟁 체제 유지",
        "단기 계약: 6개월-1년 단기 계약으로 시장 가격 반영",
        "MRP 기반 정밀 계획: 수요 예측 정확도 향상으로 재고 최소화"
    ]
    add_bullet_list(slide10, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), lv_strategies, 15)

    # Slide 11: Strategic Sourcing
    slide11 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide11, "전략자재 소싱 전략",
                      "고위험·고금액 전략자재는 장기 파트너십과 협업적 혁신이 핵심 전략입니다.",
                      11)

    st_strategies = [
        "핵심 전략: 전략적 파트너십 구축 (Strategic Partnership)",
        "장기 계약: 3-5년 장기 계약 with 수량/가격 조정 조항",
        "공동 개발: 신제품 개발 단계부터 공급업체 참여 (Early Supplier Involvement)",
        "정보 공유: 수요 예측, 생산 계획 등 핵심 정보 공유",
        "리스크 분담: 수요 변동, 원자재 가격 변동 리스크 분담 메커니즘",
        "하이브리드 계획: MRP + ROP 결합으로 유연성과 안정성 확보"
    ]
    add_bullet_list(slide11, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), st_strategies, 15)

    # Slide 12: Routine Sourcing
    slide12 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide12, "일상자재 소싱 전략",
                      "저위험·저금액 일상자재는 프로세스 효율화와 자동화로 관리 비용을 최소화합니다.",
                      12)

    rt_strategies = [
        "핵심 전략: 효율화 및 자동화 (Efficiency & Automation)",
        "E-Procurement: 전자 구매 시스템으로 주문 프로세스 자동화",
        "자동 발주: 재고 수준 기반 자동 발주 (Auto-Replenishment)",
        "프레임 계약: 1년 단위 프레임 계약으로 개별 협상 불필요",
        "공급업체 통합: 1-2개 통합 공급업체로 관리 부담 축소",
        "표준화: 품목 표준화로 대체 가능성 극대화"
    ]
    add_bullet_list(slide12, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), rt_strategies, 15)

    # Slide 13: SRM Concept
    slide13 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide13, "공급업체 관계 관리 (SRM)",
                      "Supplier Relationship Management는 공급업체를 전략적 자산으로 관리하는 체계입니다.",
                      13)

    srm_concepts = [
        "SRM 정의: 공급업체를 거래 대상이 아닌 전략적 파트너로 관리",
        "자재군별 관계 유형:",
        "  • 전략자재 → 전략적 파트너십 (Strategic Partnership)",
        "  • 레버리지자재 → 경쟁적 공급 (Competitive Supply)",
        "  • 병목자재 → 안정적 공급 (Secure Supply)",
        "  • 일상자재 → 효율적 처리 (Efficient Processing)",
        "SRM 핵심 활동: 성과 평가, 개선 협업, 리스크 관리, 혁신 촉진"
    ]
    add_bullet_list(slide13, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), srm_concepts, 16)

    # Slide 14: Supplier Scorecard
    slide14 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide14, "공급업체 성과 평가 체계",
                      "정량적 지표 기반 공급업체 성과 평가로 데이터 기반 의사결정을 지원합니다.",
                      14)

    scorecard_items = [
        "평가 4대 영역: 품질 (40%) / 납기 (30%) / 가격 (20%) / 협력성 (10%)",
        "품질 지표: 검사 통과율, 불량률, 클레임 건수",
        "납기 지표: 납기 준수율 (OTD: On-Time Delivery), 리드타임 준수",
        "가격 지표: 가격 경쟁력, 가격 안정성, 개선 제안",
        "협력성 지표: 커뮤니케이션, 문제 해결, 혁신 제안",
        "등급 체계: A (90-100점) / B (80-89) / C (70-79) / D (<70점)"
    ]
    add_bullet_list(slide14, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), scorecard_items, 15)


def add_session3_slides(prs):
    """Session 3: ABC-XYZ (Slides 15-20)"""

    # Slide 15: Session 3 Intro
    slide15 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide15, "Session 3: ABC-XYZ 재고 분류",
                      "ABC-XYZ 분석은 금액과 수요 변동성 기준으로 재고를 9개 그룹으로 세분화합니다.",
                      15)

    session3_intro = [
        "ABC-XYZ의 필요성: Kraljic Matrix가 전략적 분류라면, ABC-XYZ는 운영적 분류",
        "ABC 분석: 금액 기준 우선순위화 (파레토 법칙)",
        "XYZ 분석: 수요 변동성 기준 예측 가능성 평가",
        "9개 조합: AX, AY, AZ, BX, BY, BZ, CX, CY, CZ",
        "실무 적용: Kraljic Matrix로 전략 수립 → ABC-XYZ로 세부 운영"
    ]
    add_bullet_list(slide15, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), session3_intro, 16)

    # Slide 16: ABC Analysis
    slide16 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide16, "ABC 분석: 금액 기준 우선순위화",
                      "파레토 법칙 (80/20 법칙)에 따라 소수의 핵심 품목에 관리 역량을 집중합니다.",
                      16)

    abc_details = [
        "A급 (70-80%): 금액 상위 20% 품목이 전체 금액의 70-80% 차지",
        "  → 집중 관리: 주간 재고 점검, 정밀 수요 예측, 최적 재고 수준 관리",
        "B급 (15-25%): 중간 30% 품목이 전체 금액의 15-25% 차지",
        "  → 표준 관리: 월간 재고 점검, 일반 수요 예측",
        "C급 (5%): 하위 50% 품목이 전체 금액의 5% 차지",
        "  → 간소 관리: 분기별 점검, 간단한 발주 정책"
    ]
    add_bullet_list(slide16, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), abc_details, 15)

    # Slide 17: XYZ Analysis
    slide17 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide17, "XYZ 분석: 수요 변동성 기준 분류",
                      "변동계수 (CV)로 수요의 예측 가능성을 평가하여 재고 정책을 차별화합니다.",
                      17)

    xyz_details = [
        "X급 (CV < 0.5): 수요 안정적, 예측 가능",
        "  → 정밀 예측 가능, 안전재고 최소화, MRP 적합",
        "Y급 (0.5 ≤ CV < 1.0): 수요 변동 중간",
        "  → 통계적 예측, 적정 안전재고, MRP + 버퍼",
        "Z급 (CV ≥ 1.0): 수요 변동 크고 예측 어려움",
        "  → 예측 정확도 낮음, 높은 안전재고 또는 주문 생산, ROP 적합",
        "변동계수 (CV) = 표준편차 / 평균 수요"
    ]
    add_bullet_list(slide17, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), xyz_details, 15)

    # Slide 18: ABC-XYZ Matrix (9 combinations)
    slide18 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide18, "ABC-XYZ 매트릭스: 9가지 조합",
                      "금액과 변동성을 결합하여 9개 그룹별로 최적화된 재고 정책을 수립합니다.",
                      18)

    # 3x3 table would be ideal, but using text for simplicity
    matrix_text = slide18.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5))
    mtf = matrix_text.text_frame
    mtf.word_wrap = True

    matrix_content = [
        "AX (고금액+안정): 최우선 관리, 정밀 예측, 최소 안전재고",
        "AY (고금액+중간 변동): 중점 관리, 통계 예측, 적정 안전재고",
        "AZ (고금액+고변동): 집중 모니터링, 높은 안전재고 또는 주문 생산",
        "BX/BY/BZ: 표준 관리, 자동화된 발주 정책",
        "CX/CY/CZ: 간소 관리, 대량 발주, 낮은 관리 우선순위"
    ]

    for item in matrix_content:
        p = mtf.add_paragraph() if mtf.paragraphs[0].text else mtf.paragraphs[0]
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_DARK_GRAY

    # Slide 19: Kraljic vs ABC-XYZ
    slide19 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide19, "Kraljic vs ABC-XYZ: 언제 무엇을 사용할까?",
                      "Kraljic은 전략 수립, ABC-XYZ는 운영 최적화에 사용하며 상호 보완적입니다.",
                      19)

    comparison = [
        "Kraljic Matrix: 자재 포트폴리오 전략 수립 (소싱, 계약, 관계 관리)",
        "  → 전략적 의사결정: 어떤 공급업체와 어떻게 협력할 것인가?",
        "ABC-XYZ: 재고 운영 정책 수립 (재고 수준, 발주 주기, 안전재고)",
        "  → 운영적 의사결정: 얼마나 보유하고 언제 발주할 것인가?",
        "통합 활용: Kraljic으로 전략 → ABC-XYZ로 세부 운영 정책",
        "예: 전략자재(Kraljic) 중 AX는 집중 관리, CZ는 유연한 주문 생산"
    ]
    add_bullet_list(slide19, Inches(0.7), Inches(2.2), Inches(9.5), Inches(4.5), comparison, 15)

    # Slide 20: Summary & Next Steps
    slide20 = prs.slides.add_slide(prs.slide_layouts[6])
    add_standard_title(slide20, "Part 1 종합 요약 & Next Steps",
                      "Sessions 1-3를 통해 Kraljic Matrix와 ABC-XYZ 분석의 기초를 완성했습니다.",
                      20)

    summary_box = slide20.shapes.add_textbox(Inches(0.7), Inches(2), Inches(9.5), Inches(0.6))
    stf = summary_box.text_frame
    sp = stf.paragraphs[0]
    srun = sp.add_run()
    srun.text = "핵심 요약"
    srun.font.name = "Arial"
    srun.font.size = Pt(18)
    srun.font.bold = True
    srun.font.color.rgb = COLOR_BLACK

    summary_points = [
        "JIT → JIC: 재고는 낭비가 아닌 전략적 자산",
        "Kraljic Matrix: 공급 리스크 × 구매 임팩트로 4대 자재군 분류",
        "소싱 전략: 병목(안정성) / 레버리지(경쟁) / 전략(파트너십) / 일상(효율화)",
        "ABC-XYZ: 금액 × 변동성으로 9개 그룹 세분화",
        "Part 2 예고: 자재군별 Deep Dive (Sessions 4-7)"
    ]
    add_bullet_list(slide20, Inches(0.7), Inches(2.8), Inches(9.5), Inches(3.5), summary_points, 16)

    # Action item
    action_box = slide20.shapes.add_shape(
        1, Inches(1.5), Inches(6.5), Inches(7.83), Inches(0.6)
    )
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = COLOR_DARK_GRAY
    action_box.line.color.rgb = COLOR_DARK_GRAY

    atf = action_box.text_frame
    atf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = atf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    arun = ap.add_run()
    arun.text = "다음 학습: 자사 자재 20개를 Kraljic Matrix와 ABC-XYZ로 분류해 보세요"
    arun.font.name = "맑은 고딕"
    arun.font.size = Pt(16)
    arun.font.bold = True
    arun.font.color.rgb = COLOR_WHITE


def main():
    """Main function"""
    print("Generating Part 1 PPTX (20 slides)...")
    print("Style: S4HANA Monochrome Design")
    print("Compliance: Governing messages, Arial+맑은 고딕, 10.83\"×7.5\"\n")

    prs = create_presentation()

    output_path = "/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Foundation_Sessions1-3.pptx"
    prs.save(output_path)

    print(f"✓ PPTX saved: {output_path}")
    print(f"  - 20 slides created")
    print(f"  - Session 1: Slides 1-7 (Kraljic Matrix)")
    print(f"  - Session 2: Slides 8-14 (Sourcing Strategies)")
    print(f"  - Session 3: Slides 15-20 (ABC-XYZ)")
    print(f"\nDesign Compliance:")
    print(f"  ✓ Monochrome colors (Black/White/Gray + 1 accent)")
    print(f"  ✓ Governing messages on all content slides (16pt Bold)")
    print(f"  ✓ Arial (English) + 맑은 고딕 (Korean)")
    print(f"  ✓ Size: 10.83\" × 7.50\"")
    print(f"  ✓ White background, simple layout")


if __name__ == "__main__":
    main()
