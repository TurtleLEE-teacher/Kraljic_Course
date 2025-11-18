#!/usr/bin/env python3
"""
Session 1 PPTX Generation - Using Reference Style
전략적 재고운영 Foundation: Kraljic Matrix와 자재계획 방법론

Matches the reference file style:
- Size: 10.83" x 7.5"
- Fonts: Arial (titles), 맑은 고딕 (body)
- Background: WHITE (#FFFFFF)
- Style: Simple business document
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Create Session 1 PPTX with 9 slides"""
    prs = Presentation()

    # Match reference file dimensions exactly
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.5)

    # Define theme colors (professional business style)
    PRIMARY_COLOR = RGBColor(26, 82, 118)  # Dark Blue
    ACCENT_COLOR = RGBColor(230, 126, 34)  # Orange
    TEXT_COLOR = RGBColor(33, 37, 41)      # Dark Gray
    LIGHT_GRAY = RGBColor(248, 249, 250)   # Light Gray

    # Slide 1: Cover
    add_cover_slide(prs, PRIMARY_COLOR)

    # Slide 2: Learning Objectives
    add_learning_objectives_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    # Slide 3: JIT vs JIC Comparison
    add_jit_jic_comparison_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    # Slide 4: JIT Crisis Cases
    add_jit_crisis_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    # Slide 5: Kraljic Matrix Introduction
    add_kraljic_intro_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    # Slide 6: Kraljic Matrix 2 Axes
    add_kraljic_axes_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    # Slide 7: Kraljic Matrix Diagram (2x2)
    add_kraljic_matrix_slide(prs, PRIMARY_COLOR, ACCENT_COLOR)

    # Slide 8: 4 Material Categories Detail
    add_material_categories_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    # Slide 9: Summary and Next Steps
    add_summary_slide(prs, PRIMARY_COLOR, ACCENT_COLOR, TEXT_COLOR)

    return prs


def add_text_with_font(text_frame, text, font_name, font_size, is_bold=False, color=None):
    """Add text with specific font settings"""
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = is_bold
    if color:
        run.font.color.rgb = color
    return p


def add_cover_slide(prs, primary_color):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = primary_color

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(8.83), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "전략적 재고운영 Foundation"
    run.font.name = "Arial"
    run.font.size = Pt(44)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.7), Inches(8.83), Inches(0.8)
    )
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Kraljic Matrix와 자재계획 방법론"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(28)
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Course info
    info_box = slide.shapes.add_textbox(
        Inches(1), Inches(5), Inches(8.83), Inches(1.5)
    )
    info_frame = info_box.text_frame
    p = info_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "전략적 재고운영 및 자재계획수립 과정\n2025년 1회차"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(255, 255, 255)


def add_learning_objectives_slide(prs, primary_color, accent_color, text_color):
    """Slide 2: Learning Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(9.83), Inches(0.6)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = "학습 목표"
    run.font.name = "Arial"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = primary_color

    # Session badge
    badge_box = slide.shapes.add_textbox(
        Inches(9), Inches(0.4), Inches(1.3), Inches(0.5)
    )
    badge_frame = badge_box.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = badge_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "1회차"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Badge background
    shape = badge_box
    shape.fill.solid()
    shape.fill.fore_color.rgb = primary_color

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.3), Inches(9.43), Inches(5.5)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    objectives = [
        "JIT에서 JIC로의 패러다임 전환 배경과 필요성 이해",
        "전략적 재고운영의 핵심 개념과 차별화된 접근법 습득",
        "Kraljic Matrix를 활용한 자재 포트폴리오 분류 역량 확보",
        "자재군별 관리 철학과 자재계획 방법론의 전체 맵 이해"
    ]

    for obj in objectives:
        p = content_frame.add_paragraph()
        p.level = 0
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = obj
        run.font.name = "맑은 고딕"
        run.font.size = Pt(20)
        run.font.color.rgb = text_color

    # Footer
    add_footer(slide, "전략적 재고운영 교육", 2)


def add_jit_jic_comparison_slide(prs, primary_color, accent_color, text_color):
    """Slide 3: JIT vs JIC Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_slide_title(slide, "JIT vs JIC: 패러다임의 전환", primary_color, "1회차")

    # Left column - JIT
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(4.6), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    # Left title
    p = left_frame.paragraphs[0]
    run = p.add_run()
    run.text = "JIT (과거)"
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = primary_color

    jit_items = [
        "재고 = 낭비 (Waste)",
        "목표: 재고 Zero",
        "우선순위: 효율성 (Efficiency)",
        "안전재고: 최소 (1-2주)",
        "공급업체: Single Source",
        "의사결정: 원가 절감"
    ]

    for item in jit_items:
        p = left_frame.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(18)
        run.font.color.rgb = text_color

    # Right column - JIC
    right_box = slide.shapes.add_textbox(
        Inches(5.53), Inches(1.5), Inches(4.6), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    # Right title
    p = right_frame.paragraphs[0]
    run = p.add_run()
    run.text = "JIC (현재/미래)"
    run.font.name = "Arial"
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = accent_color

    jic_items = [
        "재고 = 전략적 자산",
        "목표: 최적 재고 (Optimal)",
        "우선순위: 회복력 (Resilience)",
        "안전재고: 차별화 (1주-6개월)",
        "공급업체: Dual/Multi Source",
        "의사결정: TCO + 리스크"
    ]

    for item in jic_items:
        p = right_frame.add_paragraph()
        p.space_before = Pt(10)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(18)
        run.font.color.rgb = text_color

    add_footer(slide, "전략적 재고운영 교육", 3)


def add_jit_crisis_slide(prs, primary_color, accent_color, text_color):
    """Slide 4: JIT Crisis Cases"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "JIT 붕괴 사례: 2020-2021 공급망 위기", primary_color, "1회차")

    # Introduction
    intro_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(9.43), Inches(0.6)
    )
    intro_frame = intro_box.text_frame
    p = intro_frame.paragraphs[0]
    run = p.add_run()
    run.text = "2020년 팬데믹은 JIT 방식의 치명적 약점을 극명하게 드러냈습니다."
    run.font.name = "맑은 고딕"
    run.font.size = Pt(18)
    run.font.italic = True
    run.font.color.rgb = text_color

    # Cases
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.3), Inches(9.43), Inches(4.3)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    cases = [
        "2021년 글로벌 반도체 대란: 자동차 생산 1,000만 대 감소, 산업 손실 $210억",
        "마스크 대란 (2020): 전 세계 생산의 50%를 중국에 의존, 봉쇄로 공급 중단",
        "컨테이너선 대란 (2021-2022): 운임비 10배 폭등, 리드타임 2배 증가",
        "GM: 반도체 부족으로 생산 차질 200만 대, 안전재고 0주 → 10-12주 확대",
        "Toyota: JIT 창시자조차 반도체 안전재고 4-6개월치 확보로 전략 전환"
    ]

    for i, case in enumerate(cases):
        p = content_frame.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = f"{i+1}. " + case
        run.font.name = "맑은 고딕"
        run.font.size = Pt(17)
        run.font.color.rgb = text_color

    add_footer(slide, "전략적 재고운영 교육", 4)


def add_kraljic_intro_slide(prs, primary_color, accent_color, text_color):
    """Slide 5: Kraljic Matrix Introduction"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "Kraljic Matrix의 탄생과 핵심 통찰", primary_color, "1회차")

    # Introduction
    intro_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(9.43), Inches(0.8)
    )
    intro_frame = intro_box.text_frame
    p = intro_frame.paragraphs[0]
    run = p.add_run()
    run.text = 'Peter Kraljic (1983, Harvard Business Review):\n"Purchasing Must Become Supply Management"'
    run.font.name = "맑은 고딕"
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.color.rgb = text_color

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.5), Inches(9.43), Inches(4.1)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    points = [
        "탄생 배경: 1970년대 2차례 석유파동으로 공급 불안정 극심",
        "핵심 통찰: \"Not all materials are created equal\" - 모든 자재가 동등하지 않다",
        "차별화 필요성: 자재의 특성에 따라 차별화된 관리 전략 적용",
        "2축 평가: 공급 리스크 (Supply Risk) × 구매 임팩트 (Purchase Impact)",
        "4대 자재군: 전략/레버리지/병목/일상 자재로 분류"
    ]

    for point in points:
        p = content_frame.add_paragraph()
        p.space_before = Pt(12)
        run = p.add_run()
        run.text = "• " + point
        run.font.name = "맑은 고딕"
        run.font.size = Pt(19)
        run.font.color.rgb = text_color

    add_footer(slide, "전략적 재고운영 교육", 5)


def add_kraljic_axes_slide(prs, primary_color, accent_color, text_color):
    """Slide 6: Kraljic Matrix Axes"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "Kraljic Matrix의 두 축", primary_color, "1회차")

    # Left column - Y Axis
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(4.6), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    run = p.add_run()
    run.text = "Y축: 공급 리스크 (Supply Risk)"
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = primary_color

    y_items = [
        "공급업체 수 (1-2개 고위험)",
        "대체 가능성 (대체 불가 고위험)",
        "시장 구조 (독과점 고위험)",
        "지리적 집중도 (특정 지역 집중 고위험)",
        "기술 복잡성 (특수 기술 고위험)",
        "리드타임 (6개월 이상 고위험)"
    ]

    for item in y_items:
        p = left_frame.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(17)
        run.font.color.rgb = text_color

    # Right column - X Axis
    right_box = slide.shapes.add_textbox(
        Inches(5.53), Inches(1.5), Inches(4.6), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    run = p.add_run()
    run.text = "X축: 구매 임팩트 (Purchase Impact)"
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = accent_color

    x_items = [
        "구매 금액 (총 구매액 대비 비중)",
        "원가 비중 (제품 원가 중 차지 비율)",
        "사업 영향도 (결품 시 생산 중단)",
        "부가가치 (최종 제품 성능 영향)",
        "품질 중요성 (품질 문제 파급 효과)"
    ]

    for item in x_items:
        p = right_frame.add_paragraph()
        p.space_before = Pt(8)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(17)
        run.font.color.rgb = text_color

    add_footer(slide, "전략적 재고운영 교육", 6)


def add_kraljic_matrix_slide(prs, primary_color, accent_color):
    """Slide 7: Kraljic Matrix 2x2 Diagram"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "Kraljic Matrix: 4대 자재군 분류", primary_color, "1회차")

    # Matrix colors
    STRATEGIC_COLOR = RGBColor(142, 68, 173)  # Purple
    BOTTLENECK_COLOR = RGBColor(230, 126, 34)  # Orange
    LEVERAGE_COLOR = RGBColor(39, 174, 96)  # Green
    ROUTINE_COLOR = RGBColor(149, 165, 166)  # Gray

    # Matrix dimensions
    matrix_left = Inches(1.5)
    matrix_top = Inches(2)
    box_width = Inches(3.6)
    box_height = Inches(2.2)
    gap = Inches(0.2)

    # Top-left: Bottleneck
    bottleneck_box = slide.shapes.add_shape(
        1,  # Rectangle
        matrix_left, matrix_top, box_width, box_height
    )
    bottleneck_box.fill.solid()
    bottleneck_box.fill.fore_color.rgb = BOTTLENECK_COLOR
    bottleneck_box.line.color.rgb = RGBColor(255, 255, 255)
    bottleneck_box.line.width = Pt(2)

    tf = bottleneck_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "병목자재\nBottleneck Materials\n\n독점 공급 부품\n희토류, 특수 소재"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Top-right: Strategic
    strategic_box = slide.shapes.add_shape(
        1,  # Rectangle
        matrix_left + box_width + gap, matrix_top, box_width, box_height
    )
    strategic_box.fill.solid()
    strategic_box.fill.fore_color.rgb = STRATEGIC_COLOR
    strategic_box.line.color.rgb = RGBColor(255, 255, 255)
    strategic_box.line.width = Pt(2)

    tf = strategic_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "전략자재\nStrategic Materials\n\n핵심 부품, 고가 장비\n전략 원자재"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Bottom-left: Routine
    routine_box = slide.shapes.add_shape(
        1,  # Rectangle
        matrix_left, matrix_top + box_height + gap, box_width, box_height
    )
    routine_box.fill.solid()
    routine_box.fill.fore_color.rgb = ROUTINE_COLOR
    routine_box.line.color.rgb = RGBColor(255, 255, 255)
    routine_box.line.width = Pt(2)

    tf = routine_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "일상자재\nRoutine Materials\n\n사무용품, 소모품\n저가 부품"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Bottom-right: Leverage
    leverage_box = slide.shapes.add_shape(
        1,  # Rectangle
        matrix_left + box_width + gap, matrix_top + box_height + gap, box_width, box_height
    )
    leverage_box.fill.solid()
    leverage_box.fill.fore_color.rgb = LEVERAGE_COLOR
    leverage_box.line.color.rgb = RGBColor(255, 255, 255)
    leverage_box.line.width = Pt(2)

    tf = leverage_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "레버리지자재\nLeverage Materials\n\n범용 부품\n표준 원자재, MRO"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    # Axes labels
    # Y-axis label (Supply Risk)
    y_label = slide.shapes.add_textbox(
        Inches(0.3), matrix_top + Inches(2), Inches(1), Inches(0.5)
    )
    tf = y_label.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "공급 리스크 ↑"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = primary_color

    # X-axis label (Purchase Impact)
    x_label = slide.shapes.add_textbox(
        matrix_left + Inches(2.7), Inches(6.8), Inches(4), Inches(0.5)
    )
    tf = x_label.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "구매 임팩트 (Purchase Impact) →"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = accent_color

    add_footer(slide, "전략적 재고운영 교육", 7)


def add_material_categories_slide(prs, primary_color, accent_color, text_color):
    """Slide 8: 4 Material Categories Detail"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "4대 자재군별 관리 전략", primary_color, "1회차")

    # Left column
    left_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(4.6), Inches(5)
    )
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    run = p.add_run()
    run.text = "전략자재 & 병목자재"
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = primary_color

    p = left_frame.add_paragraph()
    p.space_before = Pt(12)
    run = p.add_run()
    run.text = "전략자재 (고리스크+고금액)"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = RGBColor(142, 68, 173)

    strategic_items = [
        "장기 파트너십 구축",
        "Hybrid 계획 방법론",
        "공급 안정성 최우선"
    ]

    for item in strategic_items:
        p = left_frame.add_paragraph()
        p.level = 1
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(16)
        run.font.color.rgb = text_color

    p = left_frame.add_paragraph()
    p.space_before = Pt(16)
    run = p.add_run()
    run.text = "병목자재 (고리스크+저금액)"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = RGBColor(230, 126, 34)

    bottleneck_items = [
        "공급 연속성 확보",
        "ROP 기반 안전재고",
        "대체재 개발"
    ]

    for item in bottleneck_items:
        p = left_frame.add_paragraph()
        p.level = 1
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(16)
        run.font.color.rgb = text_color

    # Right column
    right_box = slide.shapes.add_textbox(
        Inches(5.53), Inches(1.5), Inches(4.6), Inches(5)
    )
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    run = p.add_run()
    run.text = "레버리지자재 & 일상자재"
    run.font.name = "Arial"
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = accent_color

    p = right_frame.add_paragraph()
    p.space_before = Pt(12)
    run = p.add_run()
    run.text = "레버리지자재 (저리스크+고금액)"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = RGBColor(39, 174, 96)

    leverage_items = [
        "경쟁 입찰 적극 활용",
        "MRP 기반 정밀 계획",
        "물량 레버리지 활용"
    ]

    for item in leverage_items:
        p = right_frame.add_paragraph()
        p.level = 1
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(16)
        run.font.color.rgb = text_color

    p = right_frame.add_paragraph()
    p.space_before = Pt(16)
    run = p.add_run()
    run.text = "일상자재 (저리스크+저금액)"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = RGBColor(149, 165, 166)

    routine_items = [
        "프로세스 효율화",
        "자동 발주 시스템",
        "관리 비용 최소화"
    ]

    for item in routine_items:
        p = right_frame.add_paragraph()
        p.level = 1
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(16)
        run.font.color.rgb = text_color

    add_footer(slide, "전략적 재고운영 교육", 8)


def add_summary_slide(prs, primary_color, accent_color, text_color):
    """Slide 9: Summary and Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "핵심 요약 및 다음 단계", primary_color, "1회차")

    # Content
    content_box = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(9.43), Inches(5.1)
    )
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    summary_items = [
        "JIT → JIC: 재고는 낭비가 아닌 전략적 자산으로 전환",
        "Kraljic Matrix: 공급 리스크 × 구매 임팩트로 4대 자재군 분류",
        "차별화 전략: 각 자재군별로 다른 관리 철학과 계획 방법론 적용",
        "2회차 예고: 자재군별 소싱 전략과 공급업체 성과 평가 방법",
        "실습 준비: 여러분의 조직 자재를 Kraljic Matrix로 분류해보세요"
    ]

    for item in summary_items:
        p = content_frame.add_paragraph()
        p.space_before = Pt(14)
        run = p.add_run()
        run.text = "• " + item
        run.font.name = "맑은 고딕"
        run.font.size = Pt(19)
        run.font.color.rgb = text_color

    # Action item box
    action_box = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(1.5), Inches(6), Inches(7.83), Inches(0.8)
    )
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(39, 174, 96)
    action_box.line.color.rgb = RGBColor(39, 174, 96)

    tf = action_box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "Action Item: 다음 세션까지 자사의 주요 자재 10개를 Kraljic Matrix에 배치해 오세요"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(17)
    run.font.bold = True
    run.font.color.rgb = RGBColor(255, 255, 255)

    add_footer(slide, "전략적 재고운영 교육", 9)


def add_slide_title(slide, title_text, color, badge_text=None):
    """Add standard slide title with optional badge"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(8.5), Inches(0.6)
    )
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = "Arial"
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = color

    if badge_text:
        badge_box = slide.shapes.add_textbox(
            Inches(9), Inches(0.4), Inches(1.3), Inches(0.5)
        )
        badge_frame = badge_box.text_frame
        badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = badge_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = badge_text
        run.font.name = "맑은 고딕"
        run.font.size = Pt(16)
        run.font.bold = True
        run.font.color.rgb = RGBColor(255, 255, 255)

        shape = badge_box
        shape.fill.solid()
        shape.fill.fore_color.rgb = color


def add_footer(slide, footer_text, slide_number):
    """Add standard footer with page number"""
    # Footer text
    footer_left = slide.shapes.add_textbox(
        Inches(0.5), Inches(7), Inches(5), Inches(0.3)
    )
    footer_frame = footer_left.text_frame
    p = footer_frame.paragraphs[0]
    run = p.add_run()
    run.text = footer_text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(108, 117, 125)

    # Slide number
    footer_right = slide.shapes.add_textbox(
        Inches(9.5), Inches(7), Inches(0.8), Inches(0.3)
    )
    footer_frame = footer_right.text_frame
    p = footer_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(slide_number)
    run.font.name = "Arial"
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(108, 117, 125)


def main():
    """Main function"""
    print("Generating Session 1 PPTX...")
    prs = create_presentation()

    # Save to PPTX_SAMPLE folder
    output_path = "/home/user/Kraljic_Course/PPTX_SAMPLE/Session1_Kraljic_Foundation.pptx"
    prs.save(output_path)

    print(f"✓ PPTX saved: {output_path}")
    print(f"  - 9 slides created")
    print(f"  - Size: 10.83\" x 7.5\" (matching reference)")
    print(f"  - Fonts: Arial (titles), 맑은 고딕 (body)")
    print(f"  - Style: Professional business document")


if __name__ == "__main__":
    main()
