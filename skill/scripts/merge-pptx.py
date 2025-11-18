#!/usr/bin/env python3
"""
PPTX Merger
Merges multiple PowerPoint files into a single presentation
"""

import sys
import os
from pathlib import Path
from pptx import Presentation

def merge_pptx_files(input_files, output_file):
    """
    Merge multiple PPTX files into one

    Args:
        input_files: List of PPTX file paths
        output_file: Output PPTX file path
    """
    print("🔗 PPTX Merger")
    print("=" * 40)
    print(f"Input files: {len(input_files)}")
    print(f"Output: {output_file}\n")

    # Create new presentation
    print("📄 Creating new presentation...")
    merged_prs = Presentation()

    # Set presentation properties
    merged_prs.core_properties.title = "Part 2. 소싱 전략 및 유형별 재고관리 프로세스"
    merged_prs.core_properties.author = "Strategic Inventory Management Course"

    # Process each input file
    slide_count = 0
    for i, input_file in enumerate(input_files):
        print(f"  Processing {i+1}/{len(input_files)}: {Path(input_file).name}")

        if not os.path.exists(input_file):
            print(f"    ⚠️  Warning: File not found, skipping")
            continue

        try:
            # Load source presentation
            source_prs = Presentation(input_file)

            # Copy each slide
            for slide in source_prs.slides:
                # Get slide layout (use blank layout)
                blank_layout = merged_prs.slide_layouts[6]  # Blank layout
                new_slide = merged_prs.slides.add_slide(blank_layout)

                # Copy slide dimensions
                new_slide.shapes._spTree.remove_all()

                # Copy all shapes from source slide
                for shape in slide.shapes:
                    el = shape.element
                    new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')

                slide_count += 1

        except Exception as e:
            print(f"    ❌ Error processing file: {e}")
            continue

    # Save merged presentation
    print(f"\n💾 Saving merged presentation...")
    print(f"   Total slides: {slide_count}")

    # Ensure output directory exists
    output_dir = os.path.dirname(output_file)
    if output_dir and not os.path.exists(output_dir):
        os.makedirs(output_dir, exist_ok=True)

    merged_prs.save(output_file)

    # Get file size
    file_size = os.path.getsize(output_file)
    file_size_mb = file_size / (1024 * 1024)

    print(f"\n✅ Merge complete!")
    print(f"📊 Output: {output_file}")
    print(f"📈 Slides: {slide_count}")
    print(f"💾 Size: {file_size_mb:.2f} MB")

def main():
    if len(sys.argv) < 3:
        print("Usage: python merge-pptx.py <output-file> <input-file-1> <input-file-2> ...")
        print("   OR: python merge-pptx.py <output-file> <input-dir>")
        sys.exit(1)

    output_file = sys.argv[1]

    # Check if second argument is a directory
    if len(sys.argv) == 3 and os.path.isdir(sys.argv[2]):
        # Read all .pptx files from directory
        input_dir = sys.argv[2]
        input_files = sorted([
            os.path.join(input_dir, f)
            for f in os.listdir(input_dir)
            if f.endswith('.pptx') and not f.startswith('~')  # Exclude temp files
        ])

        if not input_files:
            print(f"❌ Error: No PPTX files found in {input_dir}")
            sys.exit(1)

    else:
        # Multiple file arguments
        input_files = sys.argv[2:]

    merge_pptx_files(input_files, output_file)

if __name__ == '__main__':
    main()
