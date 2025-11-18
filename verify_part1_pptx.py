#!/usr/bin/env python3
"""
Verify Part 1 PPTX compliance with S4HANA design system
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def verify_pptx(filepath):
    """Verify PPTX meets S4HANA requirements"""
    print(f"\n{'='*70}")
    print(f"PPTX Compliance Verification")
    print(f"{'='*70}\n")

    prs = Presentation(filepath)

    # 1. Check dimensions
    print(f"1. Slide Dimensions:")
    width_in = prs.slide_width.inches
    height_in = prs.slide_height.inches
    print(f"   Width: {width_in:.2f}\" (Expected: 10.83\")")
    print(f"   Height: {height_in:.2f}\" (Expected: 7.50\")")
    dim_ok = abs(width_in - 10.83) < 0.01 and abs(height_in - 7.50) < 0.01
    print(f"   Status: {'✓ PASS' if dim_ok else '✗ FAIL'}\n")

    # 2. Check total slides
    print(f"2. Slide Count:")
    print(f"   Total: {len(prs.slides)} slides")
    print(f"   Status: {'✓ PASS' if len(prs.slides) == 20 else '✗ FAIL'}\n")

    # 3. Check fonts and governing messages
    print(f"3. Font and Governing Message Check:")
    gov_count = 0
    font_issues = []

    for idx, slide in enumerate(prs.slides, 1):
        has_gov = False
        slide_fonts = set()

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name:
                            slide_fonts.add(run.font.name)

                        # Check for governing message (16pt Bold)
                        if run.font.size and run.font.bold:
                            if run.font.size == Pt(16):
                                has_gov = True

        # Check fonts are Arial or 맑은 고딕
        for font in slide_fonts:
            if font not in ['Arial', '맑은 고딕', 'Malgun Gothic']:
                font_issues.append(f"Slide {idx}: Unexpected font '{font}'")

        if idx > 1:  # Skip cover slide
            if has_gov:
                gov_count += 1

    print(f"   Content slides with governing messages: {gov_count}/19")
    print(f"   Status: {'✓ PASS' if gov_count >= 19 else '✗ FAIL'}")

    if font_issues:
        print(f"   Font Issues:")
        for issue in font_issues[:5]:  # Show first 5
            print(f"     - {issue}")
    else:
        print(f"   Fonts: ✓ All slides use Arial or 맑은 고딕")
    print()

    # 4. Sample slide details
    print(f"4. Sample Slide Analysis:")
    print(f"\n   Slide 1 (Cover):")
    for shape in prs.slides[0].shapes:
        if hasattr(shape, "text") and shape.text.strip():
            print(f"     - Text: \"{shape.text[:50]}...\"")
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name and run.font.size:
                            print(f"       Font: {run.font.name}, Size: {run.font.size.pt}pt, Bold: {run.font.bold}")
                            break
                    break

    if len(prs.slides) > 1:
        print(f"\n   Slide 2 (First Content):")
        for shape in prs.slides[1].shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text[:50]
                print(f"     - Text: \"{text}...\"")
                if hasattr(shape, "text_frame"):
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.name and run.font.size:
                                print(f"       Font: {run.font.name}, Size: {run.font.size.pt}pt, Bold: {run.font.bold}")
                                break
                        break

    print(f"\n{'='*70}")
    print(f"Verification Complete")
    print(f"{'='*70}\n")

if __name__ == "__main__":
    verify_pptx("/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Foundation_Sessions1-3.pptx")
