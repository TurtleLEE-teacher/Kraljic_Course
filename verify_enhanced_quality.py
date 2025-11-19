#!/usr/bin/env python3
"""
Verify Enhanced Part 1 PPTX - Focus on slides 1-8 quality
Check: shape counts, font sizes, content density, text amount
"""

from pptx import Presentation
from pptx.util import Pt

def verify_enhanced_pptx(filepath):
    """Detailed verification of enhanced PPTX"""
    print(f"\n{'='*80}")
    print(f"Enhanced Part 1 PPTX - Quality Verification")
    print(f"{'='*80}\n")

    prs = Presentation(filepath)

    # Basic checks
    print(f"1. Basic Properties:")
    print(f"   Dimensions: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Status: {'✓ PASS' if len(prs.slides) == 48 else f'✗ FAIL (expected 48)'}\n")

    # Detailed analysis of slides 1-8
    print(f"2. Detailed Analysis (Slides 1-8):")
    print(f"{'='*80}")

    for idx in range(1, 9):  # Slides 1-8
        if idx > len(prs.slides):
            break

        slide = prs.slides[idx - 1]

        print(f"\nSlide {idx}:")

        # Count shapes by type
        shape_count = 0
        text_box_count = 0
        total_chars = 0
        font_sizes = []

        for shape in slide.shapes:
            shape_count += 1

            if hasattr(shape, "text_frame"):
                if shape.shape_type == 17:  # TEXT_BOX
                    text_box_count += 1

                # Count text
                text = shape.text
                total_chars += len(text)

                # Collect font sizes
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)

        print(f"  Shapes: {shape_count} total (Text boxes: {text_box_count})")
        print(f"  Text: {total_chars} characters")

        if font_sizes:
            unique_sizes = sorted(set([int(f) for f in font_sizes]))
            print(f"  Font sizes: {unique_sizes}")

            # Count small fonts (6-11pt)
            small_fonts = [f for f in font_sizes if 6 <= f <= 11]
            if font_sizes:
                pct = len(small_fonts) / len(font_sizes) * 100
                print(f"  Small fonts (6-11pt): {len(small_fonts)}/{len(font_sizes)} ({pct:.1f}%)")

        # Quality assessment
        quality = "✓ GOOD" if shape_count >= 15 else "⚠ LOW"
        if shape_count >= 30:
            quality = "✓✓ EXCELLENT"
        print(f"  Quality: {quality} (target: 30-40 shapes)")

    # Overall statistics
    print(f"\n{'='*80}")
    print(f"3. Overall Statistics (All 48 slides):")

    total_shapes = 0
    all_fonts = []

    for slide in prs.slides:
        slide_shapes = 0
        for shape in slide.shapes:
            slide_shapes += 1
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            all_fonts.append(run.font.size.pt)
        total_shapes += slide_shapes

    avg_shapes = total_shapes / len(prs.slides) if len(prs.slides) > 0 else 0
    print(f"   Average shapes per slide: {avg_shapes:.1f}")
    print(f"   Target: 30-40 shapes")
    print(f"   Status: {'✓ PASS' if avg_shapes >= 20 else '⚠ NEEDS IMPROVEMENT'}")

    if all_fonts:
        small_fonts = [f for f in all_fonts if 6 <= f <= 11]
        pct = len(small_fonts) / len(all_fonts) * 100
        print(f"   Small fonts (6-11pt): {pct:.1f}%")
        print(f"   Target: 60%+")

    print(f"\n{'='*80}")
    print(f"Verification Complete")
    print(f"{'='*80}\n")

if __name__ == "__main__":
    verify_enhanced_pptx("/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Session1_Enhanced_v2.pptx")
