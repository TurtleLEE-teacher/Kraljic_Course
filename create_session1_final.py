#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
1회차 강의자료 - 참고 파일 스타일 그대로
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """메인 프레젠테이션 - 참고 파일과 동일한 스타일"""
    prs = Presentation()
    # 참고 파일과 정확히 동일한 크기
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.5)

    # 슬라이드 생성
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_jit_to_jic_slide(prs)
    create_jit_crisis_slide(prs)
    create_kraljic_intro_slide(prs)
    create_kraljic_matrix_slide(prs)
    create_four_materials_slide(prs)
    create_planning_methods_slide(prs)
    create_summary_slide(prs)

    return prs

def create_title_slide(prs):
    """슬라이드 1: 타이틀"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 흰색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 메인 타이틀
    title = slide.shapes.add_textbox(
        Inches(0.45), Inches(1.5),
        Inches(6.77), Inches(2.5)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "전략적 재고운영 및\n자재계획수립"
    p.font.name = "Arial"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # 서브타이틀
    subtitle = slide.shapes.add_textbox(
        Inches(0.45), Inches(4.15),
        Inches(5.42), Inches(0.4)
    )
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = " - [1회차] Kraljic Matrix와 자재계획 방법론 -"
    p.font.name = "Arial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # 날짜/정보
    date = slide.shapes.add_textbox(
        Inches(0.45), Inches(4.75),
        Inches(6.77), Inches(0.65)
    )
    tf = date.text_frame
    p = tf.paragraphs[0]
    p.text = "2025년 11월"
    p.font.name = "Arial"
    p.font.bold = True
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT

def create_agenda_slide(prs):
    """슬라이드 2: Agenda"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.32),
        Inches(5.73), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "금일 Agenda"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)

    # Agenda 내용
    content = slide.shapes.add_textbox(
        Inches(0.71), Inches(1.15),
        Inches(7.15), Inches(4.29)
    )
    tf = content.text_frame

    agenda_items = [
        "패러다임의 전환: JIT에서 JIC로",
        "Kraljic Matrix 프레임워크",
        "4대 자재군 특성 및 관리 철학",
        "자재계획 방법론 맵",
        "통합 KPI 프레임워크"
    ]

    for i, item in enumerate(agenda_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.name = "맑은 고딕"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 0, 0)
        p.space_before = Pt(12)

def create_jit_to_jic_slide(prs):
    """슬라이드 3: JIT to JIC"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "패러다임의 전환: JIT에서 JIC로"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 설명
    desc = slide.shapes.add_textbox(
        Inches(0.3), Inches(1.01),
        Inches(10.32), Inches(0.63)
    )
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "2020년 코로나19 팬데믹 이후, 글로벌 기업들이 JIT(Just-In-Time)에서 JIC(Just-In-Case)로 재고 전략을 전환하고 있습니다."
    p.font.name = "맑은 고딕"
    p.font.size = Pt(16)
    p.font.bold = True

    # 비교표
    comparisons = [
        ["구분", "JIT (과거)", "JIC (현재/미래)"],
        ["재고 철학", "재고 = 낭비", "재고 = 전략적 자산"],
        ["목표", "재고 최소화 (Zero)", "최적 재고 (Optimal)"],
        ["우선순위", "효율성 (Efficiency)", "회복력 (Resilience)"],
        ["리스크 관점", "리스크 무시", "리스크 관리"],
        ["안전재고", "최소 (1-2주)", "차별화 (1주-6개월)"]
    ]

    start_y = 2.0
    for row_idx, row_data in enumerate(comparisons):
        for col_idx, cell_text in enumerate(row_data):
            # 셀 박스
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0.5 + col_idx * 3.3),
                Inches(start_y + row_idx * 0.6),
                Inches(3.2),
                Inches(0.55)
            )

            # 헤더 행 배경색
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(217, 217, 217)
            else:
                cell.fill.background()

            cell.line.color.rgb = RGBColor(150, 150, 150)
            cell.line.width = Pt(0.5)

            # 텍스트
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.word_wrap = True
            tf.vertical_anchor = 1  # MIDDLE

            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.name = "맑은 고딕"
            p.font.size = Pt(13) if row_idx == 0 else Pt(12)
            p.font.bold = row_idx == 0
            p.alignment = PP_ALIGN.CENTER if row_idx == 0 else PP_ALIGN.LEFT

def create_jit_crisis_slide(prs):
    """슬라이드 4: JIT의 위기"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "JIT의 위기: 2021년 글로벌 반도체 대란"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 설명
    desc = slide.shapes.add_textbox(
        Inches(0.3), Inches(1.01),
        Inches(10.32), Inches(0.63)
    )
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "차량용 반도체 부족으로 자동차 생산 1,000만 대 감소, 산업 전체 손실 $210억 (약 25조 원)"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(16)
    p.font.bold = True

    # 주요 내용
    content = slide.shapes.add_textbox(
        Inches(0.71), Inches(2.0),
        Inches(9.5), Inches(4.5)
    )
    tf = content.text_frame

    points = [
        "피해 사례",
        "  • GM, 포드, 폭스바겐 등 수개월간 생산 중단",
        "  • 2021년 한 해 동안 전 세계 자동차 생산량 약 1,000만 대 감소",
        "  • 안전재고 Zero → 공급 충격에 즉시 노출",
        "",
        "JIT가 실패한 이유",
        "  • 팬데믹으로 글로벌 공급망 동시 중단",
        "  • 리드타임 2-4주 → 3-6개월로 증가",
        "  • 컨테이너 운임 10배 폭등",
        "",
        "기업들의 대응",
        "  • Apple: 핵심 부품 안전재고 2주 → 6-8주로 확대",
        "  • Intel: 핵심 원자재 안전재고 4주 → 12주로 확대",
        "  • Toyota: 반도체 안전재고 4-6개월치 확보",
        "  • GM: 전략적 재고에 $70억 투자"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = point
        p.font.name = "맑은 고딕"

        if point and not point.startswith("  •"):
            p.font.size = Pt(16)
            p.font.bold = True
        else:
            p.font.size = Pt(14)
            p.font.bold = False

        p.space_before = Pt(6) if i > 0 else Pt(0)

def create_kraljic_intro_slide(prs):
    """슬라이드 5: Kraljic Matrix 소개"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kraljic Matrix: 자재 포트폴리오 전략의 핵심"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 설명
    desc = slide.shapes.add_textbox(
        Inches(0.3), Inches(1.01),
        Inches(10.32), Inches(1.0)
    )
    tf = desc.text_frame
    p = tf.paragraphs[0]
    p.text = "1983년 Peter Kraljic이 Harvard Business Review에 발표한 프레임워크로,\n자재를 공급 리스크와 구매 임팩트 2개 축으로 분류하여 차별화된 전략을 수립합니다."
    p.font.name = "맑은 고딕"
    p.font.size = Pt(16)
    p.font.bold = True

    # 핵심 통찰
    content = slide.shapes.add_textbox(
        Inches(0.71), Inches(2.5),
        Inches(9.5), Inches(4.0)
    )
    tf = content.text_frame

    points = [
        "핵심 통찰",
        '"Not all materials are created equal"',
        "모든 자재가 동등하게 만들어지지 않았다. 자재의 특성에 따라 차별화된 전략이 필요하다.",
        "",
        "2개 축",
        "  • Y축: 공급 리스크 (Supply Risk)",
        "     - 공급업체 수, 대체 가능성, 시장 구조, 지리적 집중도, 리드타임",
        "",
        "  • X축: 구매 임팩트 (Profit Impact)",
        "     - 구매 금액, 원가 비중, 사업 영향도, 부가가치, 품질 중요성",
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = point
        p.font.name = "맑은 고딕"

        if point == "핵심 통찰" or point == "2개 축":
            p.font.size = Pt(18)
            p.font.bold = True
        elif point.startswith('"'):
            p.font.size = Pt(16)
            p.font.bold = True
        elif point.startswith("  • "):
            p.font.size = Pt(14)
            p.font.bold = True
        else:
            p.font.size = Pt(14)
            p.font.bold = False

        p.space_before = Pt(8) if i > 0 else Pt(0)

def create_kraljic_matrix_slide(prs):
    """슬라이드 6: Kraljic Matrix 2x2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "Kraljic Matrix: 2×2 자재 분류"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 매트릭스 그리기
    center_x = 5.4
    center_y = 4.0
    box_size = 2.5

    quadrants = [
        # (x_offset, y_offset, title, desc)
        (-box_size, -box_size, "병목자재\n(Bottleneck)", "높은 공급 리스크\n낮은 구매 금액"),
        (0, -box_size, "전략자재\n(Strategic)", "높은 공급 리스크\n높은 구매 금액"),
        (-box_size, 0, "일상자재\n(Routine)", "낮은 공급 리스크\n낮은 구매 금액"),
        (0, 0, "레버리지자재\n(Leverage)", "낮은 공급 리스크\n높은 구매 금액")
    ]

    for x_off, y_off, name, desc in quadrants:
        # 박스
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(center_x + x_off),
            Inches(center_y + y_off),
            Inches(box_size - 0.05),
            Inches(box_size - 0.05)
        )
        box.fill.background()
        box.line.color.rgb = RGBColor(0, 0, 0)
        box.line.width = Pt(1.5)

        # 텍스트
        tf = box.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "맑은 고딕"
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = "\n" + desc
        p.font.name = "맑은 고딕"
        p.font.size = Pt(12)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)

    # Y축 레이블
    y_label = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.5),
        Inches(1.2), Inches(1.0)
    )
    tf = y_label.text_frame
    p = tf.paragraphs[0]
    p.text = "구매 금액\n↑"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # X축 레이블
    x_label = slide.shapes.add_textbox(
        Inches(5), Inches(6.7),
        Inches(2.0), Inches(0.5)
    )
    tf = x_label.text_frame
    p = tf.paragraphs[0]
    p.text = "공급 리스크 →"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def create_four_materials_slide(prs):
    """슬라이드 7: 4대 자재군 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "4대 자재군: 차별화된 관리 전략"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 표 형식으로 정리
    materials = [
        ["자재군", "목표", "계획 방법", "안전재고", "주요 전략"],
        ["병목자재", "공급 확보", "ROP", "4-8주 (높음)", "Dual Sourcing, 높은 서비스 수준"],
        ["레버리지자재", "원가 절감", "MRP", "1-2주 (낮음)", "경쟁 입찰, 볼륨 레버리지"],
        ["전략자재", "파트너십", "하이브리드", "3-6주 (중상)", "장기 계약, Win-Win 협력"],
        ["일상자재", "효율화", "자동화/VMI", "1주 (최소)", "프로세스 간소화, 통합 관리"]
    ]

    start_y = 1.5
    col_widths = [1.8, 1.5, 1.5, 1.8, 3.5]

    for row_idx, row_data in enumerate(materials):
        x_pos = 0.3
        for col_idx, (cell_text, width) in enumerate(zip(row_data, col_widths)):
            # 셀 박스
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x_pos),
                Inches(start_y + row_idx * 0.8),
                Inches(width),
                Inches(0.75)
            )

            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(217, 217, 217)
            else:
                cell.fill.background()

            cell.line.color.rgb = RGBColor(150, 150, 150)
            cell.line.width = Pt(0.5)

            # 텍스트
            tf = cell.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.word_wrap = True
            tf.vertical_anchor = 1  # MIDDLE

            p = tf.paragraphs[0]
            p.text = cell_text
            p.font.name = "맑은 고딕"
            p.font.size = Pt(12) if row_idx == 0 else Pt(11)
            p.font.bold = row_idx == 0 or col_idx == 0
            p.alignment = PP_ALIGN.CENTER if row_idx == 0 else PP_ALIGN.LEFT

            x_pos += width

def create_planning_methods_slide(prs):
    """슬라이드 8: 자재계획 방법론"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "자재계획 방법론 맵"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 내용
    content = slide.shapes.add_textbox(
        Inches(0.71), Inches(1.3),
        Inches(9.5), Inches(5.5)
    )
    tf = content.text_frame

    points = [
        "1. ROP (Re-Order Point) - 병목자재",
        "   • 재고 수준이 기준점 도달 시 자동 발주",
        "   • 지속 모니터링, 높은 서비스 수준 (95-99%)",
        "",
        "2. MRP (Material Requirements Planning) - 레버리지자재",
        "   • MPS + BOM 기반 소요량 계산",
        "   • 계획적 발주, 재고 최적화, Lot Sizing",
        "",
        "3. 하이브리드 계획 (LTP + MRP + ROP) - 전략자재",
        "   • 장기 계획 (18-24개월) + 중기 계획 (3-6개월) + 실시간 모니터링",
        "   • Framework Agreement, Capacity 확보, 분기별 총량 계획",
        "",
        "4. 자동화 (Min-Max / VMI) - 일상자재",
        "   • 최소/최대값 기준 자동 보충",
        "   • 공급업체 주도 재고 관리 (VMI)",
        "   • 단순 자동화, 관리 최소화, Zero Touch",
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = point
        p.font.name = "맑은 고딕"

        if point and not point.startswith("   • "):
            p.font.size = Pt(16)
            p.font.bold = True
        else:
            p.font.size = Pt(14)
            p.font.bold = False

        p.space_before = Pt(8) if i > 0 else Pt(0)

def create_summary_slide(prs):
    """슬라이드 9: 핵심 요약"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 제목
    title = slide.shapes.add_textbox(
        Inches(0.3), Inches(0.31),
        Inches(7.56), Inches(0.43)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "핵심 요약 (Key Takeaways)"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(20)
    p.font.bold = True

    # 내용
    content = slide.shapes.add_textbox(
        Inches(0.71), Inches(1.3),
        Inches(9.5), Inches(5.5)
    )
    tf = content.text_frame

    points = [
        "1. 패러다임의 전환",
        "   • JIT (재고 = 낭비) → JIC (재고 = 전략적 자산)",
        "   • 효율성 추구 → 회복력 확보",
        "   • 획일적 관리 → 차별화된 전략",
        "",
        "2. Kraljic Matrix",
        "   • 2개 축: 공급 리스크 × 구매 임팩트",
        "   • 4개 자재군: 병목 / 레버리지 / 전략 / 일상",
        "   • 각 자재군별 차별화된 관리 전략 필요",
        "",
        "3. 자재계획 방법론",
        "   • 병목자재 → ROP (재주문점 발주, 높은 안전재고)",
        "   • 레버리지자재 → MRP (계획 기반, 낮은 재고)",
        "   • 전략자재 → 하이브리드 (장기 + 중기 + 단기 통합)",
        "   • 일상자재 → 자동화 (Min-Max, VMI)",
        "",
        "다음 회차 예고: [2회차] 자재군별 소싱 전략 및 공급업체 관계 관리"
    ]

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = point
        p.font.name = "맑은 고딕"

        if point.startswith("1.") or point.startswith("2.") or point.startswith("3.") or point.startswith("다음"):
            p.font.size = Pt(18)
            p.font.bold = True
        elif point.startswith("   • "):
            p.font.size = Pt(14)
            p.font.bold = False
        else:
            p.font.size = Pt(14)

        p.space_before = Pt(10) if i > 0 else Pt(0)

def main():
    """메인 실행 함수"""
    print("\n" + "="*70)
    print("1회차 강의자료 생성 중 (참고 파일 스타일)")
    print("="*70 + "\n")

    prs = create_presentation()

    output_file = "/home/user/Kraljic_Course/PPTX_SAMPLE/Session1_KraljicMatrix_Foundation.pptx"
    prs.save(output_file)

    print(f"✅ 완료!")
    print(f"📁 파일: {output_file}")
    print(f"📊 슬라이드: {len(prs.slides)}장")
    print(f"📐 크기: 10.83\" x 7.5\" (참고 파일과 동일)")
    print(f"🎨 스타일: 흰색 배경, Arial/맑은 고딕, 심플한 비즈니스 문서")
    print("\n" + "="*70 + "\n")

if __name__ == "__main__":
    main()
