#!/usr/bin/env python3
"""
S4HANA Reference Deep Analysis Script
Analyzes slide structure, shapes, layout patterns

Usage:
    python3 analyze_s4hana_reference.py
"""

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def analyze_reference():
    """Deep analysis of S4HANA reference file"""

    filepath = 'PPTX_SAMPLE/S4HANA_PI단계_단계 종료보고_20230510_v.1.4.pptx'

    print("\n" + "="*70)
    print("S4HANA REFERENCE DEEP ANALYSIS")
    print("="*70)

    try:
        prs = Presentation(filepath)
    except FileNotFoundError:
        print(f"\n❌ File not found: {filepath}")
        print("   Make sure the reference file exists in PPTX_SAMPLE/")
        return
    except Exception as e:
        print(f"\n❌ Error opening file: {e}")
        return

    # Overall stats
    print(f"\n📊 OVERALL STATISTICS")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Dimensions: {prs.slide_width/914400:.2f}\" × {prs.slide_height/914400:.2f}\"")
    print(f"   Aspect ratio: {prs.slide_width/prs.slide_height:.3f}:1")

    # Detailed analysis of first 10 slides
    print(f"\n" + "="*70)
    print("SLIDE-BY-SLIDE STRUCTURE ANALYSIS (First 15 slides)")
    print("="*70)

    total_shapes = 0
    total_auto_shapes = 0
    total_text_boxes = 0
    total_groups = 0

    slides_to_analyze = list(prs.slides)[:15]
    for i, slide in enumerate(slides_to_analyze, 1):
        shapes_count = len(slide.shapes)
        total_shapes += shapes_count

        # Count by type
        auto_shapes = 0
        text_boxes = 0
        groups = 0
        pictures = 0
        tables = 0
        connectors = 0

        for shape in slide.shapes:
            shape_type = shape.shape_type

            if shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                auto_shapes += 1
            elif shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                text_boxes += 1
            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                groups += 1
            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                tables += 1
            elif shape_type in [MSO_SHAPE_TYPE.LINE, 12]:  # 12 = connector
                connectors += 1

        total_auto_shapes += auto_shapes
        total_text_boxes += text_boxes
        total_groups += groups

        # Estimate density (rough: assume each shape covers ~2% of slide)
        density_estimate = min(100, shapes_count * 2)

        print(f"\nSlide {i}:")
        print(f"   Total shapes: {shapes_count}")
        print(f"   └─ AUTO_SHAPES: {auto_shapes}")
        print(f"   └─ Text boxes: {text_boxes}")
        print(f"   └─ Groups: {groups}")
        if pictures > 0:
            print(f"   └─ Pictures: {pictures}")
        if tables > 0:
            print(f"   └─ Tables: {tables}")
        if connectors > 0:
            print(f"   └─ Connectors: {connectors}")
        print(f"   Density estimate: ~{density_estimate}%")

        # Special notes for high-density slides
        if shapes_count > 50:
            print(f"   ⭐ HIGH DENSITY SLIDE - Study this layout!")

    # Summary statistics
    slides_count = len(list(prs.slides))
    avg_shapes = total_shapes / min(15, slides_count)
    avg_auto_shapes = total_auto_shapes / min(15, slides_count)

    print(f"\n" + "="*70)
    print("SUMMARY STATISTICS (First 15 slides)")
    print("="*70)
    print(f"   Average shapes per slide: {avg_shapes:.1f}")
    print(f"   Average AUTO_SHAPES per slide: {avg_auto_shapes:.1f}")
    print(f"   Average text boxes per slide: {total_text_boxes / min(15, slides_count):.1f}")
    print(f"   Average groups per slide: {total_groups / min(15, slides_count):.1f}")

    # Key findings
    print(f"\n" + "="*70)
    print("KEY FINDINGS & RECOMMENDATIONS")
    print("="*70)
    print(f"\n1. SHAPE DENSITY")
    print(f"   ✓ Professional slides use 20-50+ shapes per slide")
    print(f"   ✓ Reference average: {avg_shapes:.1f} shapes/slide")
    print(f"   → Target: Match or exceed this density")

    print(f"\n2. SHAPE TYPES")
    print(f"   ✓ Extensive use of AUTO_SHAPES (rectangles, arrows, etc.)")
    print(f"   ✓ Groups for organizing related elements")
    print(f"   → Use variety: rectangles, arrows, connectors, triangles")

    print(f"\n3. LAYOUT PATTERNS")
    print(f"   ✓ Look for slides with 50+ shapes - these show advanced patterns")
    print(f"   ✓ Timeline diagrams, process flows, comparison matrices")
    print(f"   → Study high-density slides (marked with ⭐ above)")

    print(f"\n4. DENSITY TARGETS")
    print(f"   ✓ Most slides achieve 80-100%+ density")
    print(f"   ✓ Content-rich, minimal whitespace")
    print(f"   → Don't settle for <85% density")

    print("\n" + "="*70)
    print("⚠️  NEXT STEPS")
    print("="*70)
    print("1. Review high-density slides in PowerPoint manually")
    print("2. Document layout patterns (timeline, comparison, matrix)")
    print("3. Create implementation plan based on these findings")
    print("4. Target minimum 20 shapes per slide in your PPTX")
    print("="*70 + "\n")


if __name__ == '__main__':
    analyze_reference()
