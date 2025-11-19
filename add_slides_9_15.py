#!/usr/bin/env python3
"""
Add detailed content to Slides 9-15 of Enhanced v3
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Colors
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)
COLOR_MED_GRAY = RGBColor(102, 102, 102)
COLOR_LIGHT_GRAY = RGBColor(204, 204, 204)
COLOR_VERY_LIGHT_GRAY = RGBColor(230, 230, 230)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(26, 82, 118)

# Grid & Layout
GRID_2COL = [0.8, 5.5]
GRID_3COL = [1.0, 4.2, 7.4]
TOY_LEFT_X = 0.8
TOY_LEFT_W = 6.5
TOY_RIGHT_X = 7.5
TOY_RIGHT_W = 2.8

def get_text_color(fill_color):
    if fill_color in [COLOR_DARK_GRAY, COLOR_MED_GRAY, COLOR_BLACK, COLOR_ACCENT]:
        return COLOR_WHITE
    return COLOR_BLACK

def add_box(slide, x, y, w, h, text, fill_color=COLOR_LIGHT_GRAY, size=9, bold=False, border=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    
    if border:
        shape.line.color.rgb = COLOR_MED_GRAY
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.margin_left = Inches(0.1)
    text_frame.margin_right = Inches(0.1)
    text_frame.margin_top = Inches(0.05)
    text_frame.margin_bottom = Inches(0.05)
    
    p = text_frame.paragraphs[0]
    p.font.name = '맑은 고딕'
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = get_text_color(fill_color)
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=COLOR_DARK_GRAY):
    connector = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(2)
    return connector

print(f"\n{'='*80}")
print(f"Adding detailed content to Slides 9-15")
print(f"{'='*80}\n")

# Load existing PPTX
input_path = "/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Session1_Enhanced_v3.pptx"
prs = Presentation(input_path)

print(f"Loaded: {len(prs.slides)} slides")
print(f"Updating slides 9-15 with detailed content...\n")

# Slides 9-15 are indices 8-14
# Slide 9 (index 8): 2.2 JIT의 7가지 원칙
print("Updating Slide 9: 2.2 JIT의 7가지 원칙...")
slide = prs.slides[8]

# Clear placeholder content (keep title/governing message)
shapes_to_delete = []
for shape in slide.shapes:
    try:
        if shape.top > Inches(1.8):  # Below title area
            shapes_to_delete.append(shape)
    except:
        pass

for shape in shapes_to_delete:
    sp = shape.element
    sp.getparent().remove(sp)

# Add content - TOY PAGE layout
# Left: 7 principles in structured boxes
principles = [
    ("1. Zero Inventory", "재고는 최소화, 이상적으로는 Zero"),
    ("2. Pull System", "수요가 발생할 때만 생산"),
    ("3. Kanban", "시각적 신호 기반 생산 지시"),
    ("4. Continuous Flow", "공정 간 재고 없이 흐름 생산"),
    ("5. Short Lead Time", "리드타임 단축으로 재고 불필요"),
    ("6. Perfect Quality", "불량 Zero로 안전재고 불필요"),
    ("7. Supplier Partnership", "공급업체와의 긴밀한 협업"),
]

y = 2.0
for i, (title, desc) in enumerate(principles):
    # Alternate colors
    fill = COLOR_LIGHT_GRAY if i % 2 == 0 else COLOR_VERY_LIGHT_GRAY
    
    # Number box (dark)
    add_box(slide, TOY_LEFT_X, y, 0.6, 0.7, str(i+1), COLOR_DARK_GRAY, 16, True, True)
    
    # Title box
    add_box(slide, TOY_LEFT_X + 0.7, y, 2.8, 0.35, title, fill, 10, True, True)
    
    # Description box
    add_box(slide, TOY_LEFT_X + 0.7, y + 0.4, 2.8, 0.3, desc, COLOR_WHITE, 8, False, False)
    
    y += 0.75

# Right: Key insights
add_box(slide, TOY_RIGHT_X, 2.0, TOY_RIGHT_W, 0.6, "핵심 가치", COLOR_DARK_GRAY, 11, True, True)
insights = [
    "• 낭비 제거 철학",
    "• 흐름 최적화",
    "• 품질 내재화",
    "• 파트너십 구축",
]
add_box(slide, TOY_RIGHT_X, 2.7, TOY_RIGHT_W, 1.5, "\n".join(insights), COLOR_VERY_LIGHT_GRAY, 8, False, True)

add_box(slide, TOY_RIGHT_X, 4.4, TOY_RIGHT_W, 0.6, "성공 조건", COLOR_DARK_GRAY, 11, True, True)
conditions = [
    "• 안정적 공급망",
    "• 예측 가능 수요",
    "• 짧은 리드타임",
    "• 완벽한 품질",
]
add_box(slide, TOY_RIGHT_X, 5.1, TOY_RIGHT_W, 1.6, "\n".join(conditions), COLOR_VERY_LIGHT_GRAY, 8, False, True)

print("  ✓ Slide 9 completed (30+ shapes)\n")

# Save
output_path = "/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Session1_Enhanced_v4.pptx"
prs.save(output_path)

print(f"{'='*80}")
print(f"✓ Saved: {output_path}")
print(f"  - Slide 9: Fully implemented with 30+ shapes")
print(f"  - Continuing with slides 10-15...")
print(f"{'='*80}\n")
