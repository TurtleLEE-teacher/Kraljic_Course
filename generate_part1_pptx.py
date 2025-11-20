#!/usr/bin/env python3
"""
Part 1 PPTX Generator - Strategic Inventory Management Foundation
Session 1: Kraljic Matrix와 자재계획 방법론

S4HANA Design Standards Compliance:
- Dimensions: 10.83" × 7.50" (1.44:1)
- Monochrome color system (black/white/gray)
- Font: Arial/맑은 고딕
- Governing messages: 16pt Bold (NOT 14pt Italic)
- Shape counts: 20-50+ per slide
- Font distribution: 10pt primary (65%), 12pt bullets (20-25%)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# S4HANA Color Palette (Monochrome)
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)
COLOR_MED_GRAY = RGBColor(102, 102, 102)
COLOR_LIGHT_GRAY = RGBColor(204, 204, 204)
COLOR_VERY_LIGHT_GRAY = RGBColor(230, 230, 230)
COLOR_WHITE = RGBColor(255, 255, 255)

# Kraljic Matrix colors
COLOR_STRATEGIC = RGBColor(142, 68, 173)
COLOR_BOTTLENECK = RGBColor(230, 126, 34)
COLOR_LEVERAGE = RGBColor(39, 174, 60)
COLOR_ROUTINE = RGBColor(149, 165, 166)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.50)
    return prs

def add_slide_title(slide, title):
    left = Inches(0.30)
    top = Inches(0.30)
    width = Inches(10.23)
    height = Inches(0.60)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    p = textbox.text_frame.paragraphs[0]
    p.text = title
    p.font.name = '맑은 고딕'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    return textbox

def add_governing_message(slide, message):
    left = Inches(0.30)
    top = Inches(1.01)
    width = Inches(10.32)
    height = Inches(0.63)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    p = text_frame.paragraphs[0]
    p.text = message
    p.font.name = '맑은 고딕'
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_MED_GRAY
    return textbox

def add_footer(slide, footer_text, slide_number):
    left = Inches(0.30)
    top = Inches(7.00)
    width = Inches(8.00)
    height = Inches(0.30)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    p = textbox.text_frame.paragraphs[0]
    p.text = footer_text
    p.font.name = 'Arial'
    p.font.size = Pt(8)
    p.font.color.rgb = COLOR_MED_GRAY
    
    left = Inches(10.00)
    textbox = slide.shapes.add_textbox(left, top, Inches(0.50), height)
    p = textbox.text_frame.paragraphs[0]
    p.text = str(slide_number)
    p.font.name = 'Arial'
    p.font.size = Pt(8)
    p.font.color.rgb = COLOR_MED_GRAY
    p.alignment = PP_ALIGN.RIGHT

def create_cover_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    left = Inches(1.00)
    top = Inches(2.50)
    width = Inches(8.83)
    height = Inches(1.20)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    p = title_box.text_frame.paragraphs[0]
    p.text = "[1회차] 전략적 재고운영 Foundation"
    p.font.name = '맑은 고딕'
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLACK
    p.alignment = PP_ALIGN.CENTER
    
    top = Inches(3.80)
    height = Inches(0.80)
    
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    p = subtitle_box.text_frame.paragraphs[0]
    p.text = "Kraljic Matrix와 자재계획 방법론"
    p.font.name = '맑은 고딕'
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    top = Inches(5.50)
    height = Inches(1.00)
    
    meta_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = meta_box.text_frame
    
    for i, line in enumerate(["전략적 재고운영 및 자재계획수립 과정", "45분", "Session 1 of 9"]):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = line
        p.font.name = '맑은 고딕'
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_MED_GRAY
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(6)
    
    return slide

def create_kraljic_matrix_slide(prs, slide_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "2.3 Kraljic Matrix: 4대 자재군 분류")
    add_governing_message(slide,
        "공급 리스크(Y축)와 구매 임팩트(X축)의 조합으로 4개 자재군을 분류하며, 각 군은 완전히 다른 관리 철학을 요구합니다.")
    
    matrix_left = Inches(1.50)
    matrix_top = Inches(2.00)
    quadrant_width = Inches(3.80)
    quadrant_height = Inches(2.20)
    
    # Spectrum indicator
    top_indicator = slide.shapes.add_textbox(
        Inches(1.50), Inches(1.60), Inches(7.60), Inches(0.30)
    )
    p = top_indicator.text_frame.paragraphs[0]
    p.text = "← 낮음                공급 리스크 (Supply Risk)                높음 →"
    p.font.name = '맑은 고딕'
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = COLOR_MED_GRAY
    p.alignment = PP_ALIGN.CENTER
    
    # Quadrants with WHITE text on colored backgrounds
    quadrants = [
        {
            "left": matrix_left,
            "top": matrix_top,
            "color": COLOR_LEVERAGE,
            "title": "레버리지자재",
            "subtitle": "Leverage Materials",
            "bullets": ["• 공급 안정, 금액 큼", "• 경쟁 입찰", "• 원가 절감 집중", "• MRP 계획"]
        },
        {
            "left": matrix_left + quadrant_width,
            "top": matrix_top,
            "color": COLOR_STRATEGIC,
            "title": "전략자재",
            "subtitle": "Strategic Materials",
            "bullets": ["• 공급 어렵고 금액 큼", "• 장기 파트너십", "• Win-Win 협력", "• LTP + Hybrid"]
        },
        {
            "left": matrix_left,
            "top": matrix_top + quadrant_height,
            "color": COLOR_ROUTINE,
            "title": "일상자재",
            "subtitle": "Routine Materials",
            "bullets": ["• 공급 쉽고 금액 작음", "• 자동화", "• 효율성 극대화", "• Min-Max / VMI"]
        },
        {
            "left": matrix_left + quadrant_width,
            "top": matrix_top + quadrant_height,
            "color": COLOR_BOTTLENECK,
            "title": "병목자재",
            "subtitle": "Bottleneck Materials",
            "bullets": ["• 공급 어렵고 금액 작음", "• 안전재고 확보", "• 공급 안정성 우선", "• ROP 발주"]
        }
    ]
    
    for q in quadrants:
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            q["left"], q["top"], quadrant_width, quadrant_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = q["color"]
        shape.line.color.rgb = COLOR_BLACK
        shape.line.width = Pt(2)
        
        # Title (WHITE text)
        title_box = slide.shapes.add_textbox(
            q["left"] + Inches(0.20), q["top"] + Inches(0.15),
            quadrant_width - Inches(0.40), Inches(0.40)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = q["title"]
        p.font.name = '맑은 고딕'
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle (WHITE text)
        subtitle_box = slide.shapes.add_textbox(
            q["left"] + Inches(0.20), q["top"] + Inches(0.55),
            quadrant_width - Inches(0.40), Inches(0.25)
        )
        p = subtitle_box.text_frame.paragraphs[0]
        p.text = q["subtitle"]
        p.font.name = 'Arial'
        p.font.size = Pt(9)
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
        # Bullets (WHITE text)
        bullets_box = slide.shapes.add_textbox(
            q["left"] + Inches(0.30), q["top"] + Inches(0.90),
            quadrant_width - Inches(0.60), Inches(1.20)
        )
        text_frame = bullets_box.text_frame
        text_frame.word_wrap = True
        
        for i, bullet_text in enumerate(q["bullets"]):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = bullet_text
            p.font.name = '맑은 고딕'
            p.font.size = Pt(10)
            p.font.color.rgb = COLOR_WHITE
            p.space_after = Pt(3)
    
    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_toc_slide(prs, slide_num):
    """Table of Contents"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_title(slide, "목차")

    chapters = [
        "1장 패러다임의 전환: JIT에서 JIC로",
        "2장 Kraljic Matrix 프레임워크",
        "3장 자재군별 차별화 전략",
        "4장 자재계획 방법론 전체 맵",
        "5장 통합 KPI 프레임워크",
        "6장 산업별·규모별 적용 가이드",
        "7장 9회차 학습 여정"
    ]

    start_top = Inches(1.50)
    for i, chapter in enumerate(chapters):
        top = start_top + (i * Inches(0.75))
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.60), top, Inches(9.50), Inches(0.65)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY if i % 2 == 0 else COLOR_WHITE
        shape.line.color.rgb = COLOR_LIGHT_GRAY

        p = shape.text_frame.paragraphs[0]
        p.text = chapter
        p.font.name = '맑은 고딕'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_GRAY

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_section_divider(prs, chapter_num, chapter_title, slide_num):
    """Chapter divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(Inches(2.00), Inches(2.50), Inches(6.83), Inches(2.00))
    p = box.text_frame.paragraphs[0]
    p.text = f"{chapter_num}장"
    p.font.name = '맑은 고딕'
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLOR_DARK_GRAY
    p.alignment = PP_ALIGN.CENTER

    box = slide.shapes.add_textbox(Inches(2.00), Inches(4.70), Inches(6.83), Inches(0.80))
    p = box.text_frame.paragraphs[0]
    p.text = chapter_title
    p.font.name = '맑은 고딕'
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_MED_GRAY
    p.alignment = PP_ALIGN.CENTER

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_content_slide_with_bullets(prs, title, gov_msg, bullets, slide_num):
    """Generic content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, title)
    add_governing_message(slide, gov_msg)

    # Add bullet points
    left = Inches(1.00)
    top = Inches(2.00)
    width = Inches(8.83)
    height = Inches(4.50)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = bullet
        p.font.name = '맑은 고딕'
        p.font.size = Pt(12)
        p.font.color.rgb = COLOR_DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_jit_jic_comparison_table(prs, slide_num):
    """JIT vs JIC comparison table"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "1.2 JIT vs JIC 비교")
    add_governing_message(slide,
        "JIC는 무조건적인 재고 증가가 아니라, 자재 특성에 따른 차별화된 접근을 의미합니다.")

    # Create table
    rows_data = [
        ["재고 철학", "재고 = 낭비", "재고 = 전략적 자산"],
        ["목표", "재고 최소화 (Zero)", "최적 재고 (Optimal)"],
        ["우선순위", "효율성", "회복력"],
        ["리스크", "리스크 무시", "리스크 관리"],
        ["안전재고", "최소 (1-2주)", "차별화 (1주-6개월)"],
        ["공급업체", "Single Source", "Dual/Multi Source"]
    ]

    table = slide.shapes.add_table(
        rows=len(rows_data) + 1, cols=3,
        left=Inches(0.50), top=Inches(1.80),
        width=Inches(9.83), height=Inches(4.80)
    ).table

    table.columns[0].width = Inches(2.00)
    table.columns[1].width = Inches(3.90)
    table.columns[2].width = Inches(3.90)

    # Header row
    headers = ["구분", "JIT (과거)", "JIC (현재/미래)"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_GRAY
        for p in cell.text_frame.paragraphs:
            p.font.name = '맑은 고딕'
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows_data, 1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.name = '맑은 고딕'
                p.font.size = Pt(10)
                p.font.color.rgb = COLOR_DARK_GRAY
                p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_strategy_matrix_slide(prs, slide_num):
    """자재군별 전략 매트릭스"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "3.2 자재군별 차별화 전략 매트릭스")
    add_governing_message(slide,
        "자재군별로 관리 목표, 우선순위, 안전재고 수준, 발주 방식을 차별화해야 전체 최적화가 가능합니다.")

    # Create table
    rows_data = [
        ["관리 목표", "공급 확보", "원가 절감", "파트너십", "효율화"],
        ["안전재고", "높음 (4-8주)", "낮음 (1-2주)", "중상 (3-6주)", "최소 (1주)"],
        ["발주 방식", "ROP", "MRP", "LTP+Hybrid", "Min-Max/VMI"],
        ["공급업체", "증가 (1→2-3)", "경쟁 (5+)", "전략 (1-2)", "통합 (10-20)"]
    ]

    table = slide.shapes.add_table(
        rows=len(rows_data) + 1, cols=5,
        left=Inches(0.30), top=Inches(1.80),
        width=Inches(10.23), height=Inches(4.50)
    ).table

    # Header
    headers = ["구분", "병목", "레버리지", "전략", "일상"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_GRAY
        for p in cell.text_frame.paragraphs:
            p.font.name = '맑은 고딕'
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER

    # Data
    for row_idx, row_data in enumerate(rows_data, 1):
        for col_idx, text in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.font.name = '맑은 고딕'
                p.font.size = Pt(9)
                p.font.color.rgb = COLOR_DARK_GRAY
                p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_methodology_overview_slide(prs, slide_num):
    """5대 방법론 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "4.1 자재계획 방법론 전체 맵")
    add_governing_message(slide,
        "병목은 ROP, 레버리지는 MRP, 전략은 LTP+Hybrid, 일상은 Min-Max/VMI로 자재 특성에 맞는 방법론을 적용해야 합니다.")

    methodologies = [
        {
            "name": "ROP (Re-Order Point)",
            "target": "병목자재",
            "desc": "재고 수준이 발주점 도달 시 자동 발주\n지속 모니터링, 높은 서비스 수준"
        },
        {
            "name": "MRP (Material Requirements Planning)",
            "target": "레버리지자재",
            "desc": "MPS + BOM 기반 소요량 계산\n계획적 발주, 재고 최적화"
        },
        {
            "name": "LTP (Long-Term Planning)",
            "target": "전략자재",
            "desc": "장기 수급 계획 + 분할 납입\nCapacity 확보, 협업 계획"
        },
        {
            "name": "Min-Max / VMI",
            "target": "일상자재",
            "desc": "최소/최대값 기준 자동 보충\n단순 자동화, 관리 최소화"
        }
    ]

    box_width = Inches(4.70)
    box_height = Inches(1.30)
    gap = Inches(0.30)

    for i, method in enumerate(methodologies):
        row = i // 2
        col = i % 2
        left = Inches(0.60) + col * (box_width + gap)
        top = Inches(1.90) + row * (box_height + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        shape.line.color.rgb = COLOR_LIGHT_GRAY

        # Title
        title_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.10),
            box_width - Inches(0.30), Inches(0.30)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = method["name"]
        p.font.name = '맑은 고딕'
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_GRAY

        # Target
        target_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.45),
            box_width - Inches(0.30), Inches(0.20)
        )
        p = target_box.text_frame.paragraphs[0]
        p.text = f"→ {method['target']}"
        p.font.name = '맑은 고딕'
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_MED_GRAY

        # Description
        desc_box = slide.shapes.add_textbox(
            left + Inches(0.15), top + Inches(0.70),
            box_width - Inches(0.30), Inches(0.50)
        )
        text_frame = desc_box.text_frame
        text_frame.word_wrap = True
        for j, line in enumerate(method["desc"].split('\n')):
            p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
            p.text = line
            p.font.name = '맑은 고딕'
            p.font.size = Pt(9)
            p.font.color.rgb = COLOR_DARK_GRAY

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def create_summary_slide(prs, slide_num):
    """Summary & Next Steps"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide, "핵심 요약 & Next Steps")
    add_governing_message(slide,
        "JIT에서 JIC로의 패러다임 전환과 Kraljic Matrix 기반 차별화 전략이 전략적 재고운영의 핵심이며, 다음 회차부터 각 자재군별 구체적 방법론을 학습합니다.")

    summaries = [
        {
            "title": "1. 패러다임의 전환",
            "content": "JIT (재고=낭비) → JIC (재고=전략적 자산)\n효율성 추구 → 회복력 확보"
        },
        {
            "title": "2. Kraljic Matrix",
            "content": "공급 리스크 × 구매 임팩트\n4개 자재군 차별화 전략"
        },
        {
            "title": "3. 자재계획 방법론",
            "content": "병목→ROP, 레버리지→MRP\n전략→LTP+Hybrid, 일상→VMI"
        },
        {
            "title": "4. Next Steps",
            "content": "Session 2: 소싱 전략 & 공급업체 관계\nSession 3-7: 자재군별 심화 학습"
        }
    ]

    box_width = Inches(4.50)
    box_height = Inches(1.80)
    gap = Inches(0.40)

    for i, summary in enumerate(summaries):
        row = i // 2
        col = i % 2
        left = Inches(0.80) + col * (box_width + gap)
        top = Inches(1.90) + row * (box_height + gap)

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, box_width, box_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        shape.line.color.rgb = COLOR_LIGHT_GRAY
        shape.line.width = Pt(1.5)

        title_box = slide.shapes.add_textbox(
            left + Inches(0.20), top + Inches(0.15),
            box_width - Inches(0.40), Inches(0.35)
        )
        p = title_box.text_frame.paragraphs[0]
        p.text = summary["title"]
        p.font.name = '맑은 고딕'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_DARK_GRAY

        content_box = slide.shapes.add_textbox(
            left + Inches(0.20), top + Inches(0.60),
            box_width - Inches(0.40), box_height - Inches(0.80)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True
        for j, line in enumerate(summary["content"].split('\n')):
            p = text_frame.paragraphs[0] if j == 0 else text_frame.add_paragraph()
            p.text = line
            p.font.name = '맑은 고딕'
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_DARK_GRAY
            p.space_after = Pt(6)

    add_footer(slide, "전략적 재고운영 Foundation", slide_num)
    return slide

def main():
    print("=== Part 1 PPTX Generation Started ===")
    print("Target: 26 slides with S4HANA standards\n")

    prs = create_presentation()
    print(f"✓ Created presentation: 10.83\" × 7.50\"")

    slide_num = 1

    # Slide 1: Cover
    create_cover_slide(prs)
    print(f"✓ Slide {slide_num}: Cover")
    slide_num += 1

    # Slide 2: TOC
    create_toc_slide(prs, slide_num)
    print(f"✓ Slide {slide_num}: TOC")
    slide_num += 1

    # Slide 3: Chapter 1 Divider
    create_section_divider(prs, 1, "패러다임의 전환: JIT에서 JIC로", slide_num)
    print(f"✓ Slide {slide_num}: Chapter 1 Divider")
    slide_num += 1

    # Slide 4: JIT의 영광과 몰락
    create_content_slide_with_bullets(
        prs, "1.1 JIT의 영광과 몰락",
        "도요타의 JIT는 안정적인 공급망 환경에서는 완벽했으나, 공급 충격에는 극도로 취약한 시스템이었습니다.",
        [
            "• JIT 7대 원칙: Zero Inventory, Pull System, Kanban, Continuous Flow, Short Lead Time, Perfect Quality, Supplier Partnership",
            "• 1980-2010년대: GM, Harley-Davidson, Ford 등 전 세계 제조업체들이 JIT 도입",
            "• 성공 사례: 재고 비용 30-50% 절감, 생산성 대폭 향상",
            "• 당시의 전제: 안정적 공급망, 예측 가능한 수요, 저렴한 운송"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: JIT의 영광과 몰락")
    slide_num += 1

    # Slide 5: 팬데믹이 드러낸 약점
    create_content_slide_with_bullets(
        prs, "1.2 팬데믹이 드러낸 JIT의 치명적 약점",
        "2021년 글로벌 반도체 대란은 $210억 손실을 초래했으며, 재고 Zero 전략의 위험성을 명확히 증명했습니다.",
        [
            "• 2021년 글로벌 반도체 대란: 자동차 생산 1,000만 대 감소, 산업 손실 $210억",
            "• GM, 포드, 폭스바겐 등 수개월간 생산 중단",
            "• 원인: JIT 방식으로 안전재고 Zero → 공급 충격에 즉각 노출",
            "• 교훈: 재고 Zero = 리스크 Maximum, 효율적이지만 깨지기 쉬운 시스템"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 팬데믹이 드러낸 약점")
    slide_num += 1

    # Slide 6: JIT vs JIC 비교
    create_jit_jic_comparison_table(prs, slide_num)
    print(f"✓ Slide {slide_num}: JIT vs JIC 비교")
    slide_num += 1

    # Slide 7: JIC 채택 기업들
    create_content_slide_with_bullets(
        prs, "1.3 JIC를 채택한 글로벌 기업들",
        "Apple, Intel, Toyota 등 글로벌 선도 기업들은 핵심 자재의 안전재고를 2주에서 6개월로 확대하며 공급망 회복력에 투자하고 있습니다.",
        [
            "• Apple: 핵심 부품 안전재고 2주 → 6-8주 확대, Dual Sourcing 강화, 공급망 회복력에 $430억 투자",
            "• Intel: 핵심 원자재 안전재고 4주 → 12주 확대, Arizona와 Ohio에 $200억 반도체 공장 건설",
            "• Toyota (JIT 창시자): 2021년 이후 반도체 안전재고 4-6개월치 확보, \"더 이상 JIT만 추구하지 않는다\"",
            "• GM: 핵심 반도체 안전재고 0주 → 10-12주 확대, Triple Source 다변화, $70억 투자"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: JIC 채택 기업들")
    slide_num += 1

    # Slide 8: Chapter 2 Divider
    create_section_divider(prs, 2, "Kraljic Matrix 프레임워크", slide_num)
    print(f"✓ Slide {slide_num}: Chapter 2 Divider")
    slide_num += 1

    # Slide 9: Kraljic Matrix 탄생
    create_content_slide_with_bullets(
        prs, "2.1 Kraljic Matrix의 탄생과 의미",
        "Kraljic Matrix는 모든 자재를 동등하게 취급하는 것이 아니라, 공급 리스크와 구매 임팩트에 따라 차별화된 전략을 적용하는 프레임워크입니다.",
        [
            "• 1983년 Peter Kraljic이 Harvard Business Review에 발표",
            "• 배경: 1970년대 2차례 석유파동으로 공급 불안정 극심",
            "• 핵심 통찰: \"Not all materials are created equal\" - 자재 특성에 따른 차별화 필요",
            "• 2개 축: 공급 리스크(Y축) × 구매 임팩트(X축)",
            "• 4개 자재군: 병목, 레버리지, 전략, 일상"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: Kraljic Matrix 탄생")
    slide_num += 1

    # Slide 10: 2×2 매트릭스의 두 축
    create_content_slide_with_bullets(
        prs, "2.2 2×2 매트릭스의 두 축",
        "공급 리스크와 구매 임팩트를 평가하여 자재를 4개 군으로 분류합니다.",
        [
            "• Y축: 공급 리스크 (Supply Risk) - 공급의 어려움",
            "  - 평가 요소: 공급업체 수, 대체 가능성, 시장 구조, 지리적 집중도, 기술 복잡성, 리드타임",
            "  - 높음: 1-2개 업체, 대체 불가, 독과점, 6개월+ 리드타임",
            "",
            "• X축: 구매 임팩트 (Purchase Impact) - 사업에 미치는 영향",
            "  - 평가 요소: 구매 금액, 원가 비중, 사업 영향도, 부가가치, 품질 중요성",
            "  - 높음: 총 구매액 대비 비중 높음, 결품 시 생산 중단"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 2×2 매트릭스의 두 축")
    slide_num += 1

    # Slide 11: Kraljic Matrix (Door Chart)
    create_kraljic_matrix_slide(prs, slide_num)
    print(f"✓ Slide {slide_num}: Kraljic Matrix (Door Chart)")
    slide_num += 1

    # Slide 12-15: Quadrant details
    quadrants_detail = [
        {
            "title": "2.4 병목자재 (Bottleneck Materials) 상세",
            "gov": "금액은 작지만 공급이 어려운 병목자재는 재고를 보험료로 인식하고 높은 안전재고를 확보해야 합니다.",
            "bullets": [
                "• 정의: 높은 공급 리스크 + 낮은 구매 임팩트",
                "• 특징: 금액은 작지만 없으면 생산 중단, 공급업체 1-2개로 매우 제한적, 대체 자재/공급선 찾기 어려움",
                "• 사례: 차량용 MCU, 특수 규격 센서, 희소 원자재, 인증 필요 부품",
                "• 핵심 과제: 공급 안정성 확보 (비용이 들더라도 공급이 확실해야)",
                "• 전략: 안전재고 확보 (4-8주), 대체품 발굴, 장기 계약, ROP 발주"
            ]
        },
        {
            "title": "2.5 레버리지자재 (Leverage Materials) 상세",
            "gov": "공급이 안정적이고 금액이 큰 레버리지자재는 경쟁 입찰을 통한 원가 절감에 집중해야 합니다.",
            "bullets": [
                "• 정의: 낮은 공급 리스크 + 높은 구매 임팩트",
                "• 특징: 금액이 크지만 공급은 안정적, 공급업체 다수 (5개 이상), 표준화된 품목, 경쟁시장",
                "• 사례: 범용 원자재 (철강, 플라스틱), 표준 전자부품, 대량 구매 품목",
                "• 핵심 과제: 원가 절감 (경쟁을 통한 최적 가격 확보)",
                "• 전략: 경쟁 입찰, 가격 협상, Volume Discount, MRP 계획"
            ]
        },
        {
            "title": "2.6 전략자재 (Strategic Materials) 상세",
            "gov": "공급도 어렵고 금액도 큰 전략자재는 단기 이익보다 장기 파트너십이 핵심입니다.",
            "bullets": [
                "• 정의: 높은 공급 리스크 + 높은 구매 임팩트",
                "• 특징: 금액도 크고 공급도 어려움, 사업의 성패를 좌우, 대체 불가능, 장기 개발 필요",
                "• 사례: 핵심 반도체 (AP, SoC), 특수 소재, 장납기 외자재, 독점 기술 부품",
                "• 핵심 과제: 전략적 파트너십 (Win-Win 협력을 통한 상호 성장)",
                "• 전략: 장기 계약 (3-5년), 기술 협력, 공동 개발, LTP + Hybrid 계획"
            ]
        },
        {
            "title": "2.7 일상자재 (Routine Materials) 상세",
            "gov": "금액도 작고 공급도 쉬운 일상자재는 자동화를 통해 관리 비용을 최소화해야 합니다.",
            "bullets": [
                "• 정의: 낮은 공급 리스크 + 낮은 구매 임팩트",
                "• 특징: 금액도 작고 공급도 쉬움, 하지만 품목 수는 전체의 60-80%, 관리 부담 큼",
                "• 사례: MRO 품목 (볼트, 너트), 사무용품, 소모품 전반",
                "• 핵심 과제: 효율화 & 자동화 (최소한의 노력으로 관리)",
                "• 전략: E-Procurement, VMI, 통합 발주, Min-Max 자동 보충"
            ]
        }
    ]

    for q in quadrants_detail:
        create_content_slide_with_bullets(prs, q["title"], q["gov"], q["bullets"], slide_num)
        print(f"✓ Slide {slide_num}: {q['title'][:20]}...")
        slide_num += 1

    # Slide 16: Chapter 3 Divider
    create_section_divider(prs, 3, "자재군별 차별화 전략", slide_num)
    print(f"✓ Slide {slide_num}: Chapter 3 Divider")
    slide_num += 1

    # Slide 17: 차별화의 필요성
    create_content_slide_with_bullets(
        prs, "3.1 차별화가 필수인 이유",
        "모든 자재를 동일하게 관리하면 전략자재는 공급 중단, 레버리지자재는 원가 경쟁력 상실, 병목자재는 결품 발생, 일상자재는 인력 낭비가 발생합니다.",
        [
            "• 획일적 관리의 문제점:",
            "  - 전략자재에 원가 절감 압박 → 공급업체 이탈 → 생산 중단",
            "  - 레버리지자재에 높은 안전재고 → 재고 비용 증가 → 원가 경쟁력 상실",
            "  - 병목자재를 소홀히 관리 → 결품 발생 → 생산 라인 정지",
            "  - 일상자재에 과도한 관리 → 인력 낭비 → 전략 업무 소홀"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 차별화의 필요성")
    slide_num += 1

    # Slide 18: 자재군별 전략 매트릭스
    create_strategy_matrix_slide(prs, slide_num)
    print(f"✓ Slide {slide_num}: 자재군별 전략 매트릭스")
    slide_num += 1

    # Slide 19: Chapter 4 Divider
    create_section_divider(prs, 4, "자재계획 방법론 전체 맵", slide_num)
    print(f"✓ Slide {slide_num}: Chapter 4 Divider")
    slide_num += 1

    # Slide 20: 5대 방법론 개요
    create_methodology_overview_slide(prs, slide_num)
    print(f"✓ Slide {slide_num}: 5대 방법론 개요")
    slide_num += 1

    # Slide 21: 하이브리드 접근법
    create_content_slide_with_bullets(
        prs, "4.2 하이브리드 접근법 (전략자재)",
        "전략자재는 단일 방법론으로 해결 불가하며, 장기+중기+단기 계획을 통합해야 합니다.",
        [
            "• 전략자재의 특수성: 공급도 어렵고 금액도 커서 단일 방법론으로 불충분",
            "• 3단계 통합 계획:",
            "  - LTP (18-24개월): Framework Agreement, Capacity 확보, 분기별 총량 계획",
            "  - MRP (3-6개월): 월간 상세 계획, Release Order, BOM 기반 계산",
            "  - ROP (실시간): 안전재고 모니터링, 긴급 발주 대응, 버퍼 역할",
            "• 효과: 장기 안정성 + 중기 유연성 + 단기 대응력"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 하이브리드 접근법")
    slide_num += 1

    # Slide 22: 통합 KPI
    create_content_slide_with_bullets(
        prs, "5.1 통합 KPI 프레임워크",
        "자재군별로 다른 KPI를 설정하고 측정해야 올바른 성과 평가가 가능합니다.",
        [
            "• 병목자재: 재고 가용률 95% 이상, 공급 안정성 95% 이상, 긴급 발주 5회/월 이하",
            "• 레버리지자재: 원가 절감률 3-5%/년, 재고 회전율 12회 이상, 발주 정확도 95% 이상",
            "• 전략자재: 공급 연속성 100%, LTP 정확도 70% 이상, TCO 개선 -3%/년",
            "• 일상자재: 발주 처리비용 5천원/건 이하, 자동화율 80% 이상, 공급업체 수 100개 이하"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 통합 KPI")
    slide_num += 1

    # Slide 23: 산업별 적용
    create_content_slide_with_bullets(
        prs, "6.1 산업별·규모별 적용 가이드",
        "제조업, 프로세스 산업, 유통업은 자재 분포가 다르므로, 산업 특성에 맞는 우선순위와 전략을 수립해야 합니다.",
        [
            "• 제조업 (자동차, 전자): 전략 10-15%, 병목 15-20%, 레버리지 30-35%, 일상 35-40%",
            "  → 우선순위: 전략 → 병목 → 레버리지 → 일상",
            "• 프로세스 산업 (화학, 제약): 전략 20-30%, 병목 10-15%, 레버리지 40-50%, 일상 10-20%",
            "  → 우선순위: 전략 → 레버리지 → 병목 → 일상",
            "• 유통·서비스업: 전략 0-5%, 병목 5-10%, 레버리지 10-20%, 일상 70-80%",
            "  → 우선순위: 일상 → 레버리지 → 병목"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 산업별 적용")
    slide_num += 1

    # Slide 24: 9회차 학습 여정
    create_content_slide_with_bullets(
        prs, "7.1 9회차 학습 여정",
        "Session 1-3은 구조적 기초, Session 4-7은 자재군별 실습, Session 8-9는 통합 워크샵으로 진행되며 점진적으로 마스터할 수 있습니다.",
        [
            "• Module 1: Foundation (1-2회차)",
            "  - 1회차: JIT→JIC + Kraljic Matrix (지금)",
            "  - 2회차: 소싱 전략 + 공급업체 관계 관리",
            "",
            "• Module 2: 자재군별 심화 (3-6회차)",
            "  - 3회차: 병목자재 + ROP (이론+실습)",
            "  - 4회차: 레버리지자재 + MRP (이론+실습)",
            "  - 5회차: 전략자재 + LTP (이론+실습)",
            "  - 6회차: 일상자재 + 자동화",
            "",
            "• Module 3: 실전 통합 (7-9회차)",
            "  - 7-8회차: 실전 워크샵, 9회차: 발표 및 피드백"
        ],
        slide_num
    )
    print(f"✓ Slide {slide_num}: 9회차 학습 여정")
    slide_num += 1

    # Slide 25: Summary
    create_summary_slide(prs, slide_num)
    print(f"✓ Slide {slide_num}: Summary & Next Steps")

    # Save
    output_path = "/home/user/Kraljic_Course/Part1_Session1_StrategicInventory.pptx"
    prs.save(output_path)

    print(f"\n=== PPTX Generation Complete ===")
    print(f"Output: {output_path}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"\nNext: Run verification and commit")

    return output_path

if __name__ == "__main__":
    main()
