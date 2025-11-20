#!/usr/bin/env python3
"""
Part 2 PPTX Generator - Session 2: 자재군별 소싱 전략 및 공급업체 관계 관리
S4HANA Professional Style (10.83" × 7.50")
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# ============================================================================
# COLOR SYSTEM (Monochrome)
# ============================================================================
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)
COLOR_MED_GRAY = RGBColor(102, 102, 102)
COLOR_LIGHT_GRAY = RGBColor(204, 204, 204)
COLOR_VERY_LIGHT_GRAY = RGBColor(230, 230, 230)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(26, 82, 118)  # Dark blue (MINIMAL use)

# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def create_presentation():
    """Create presentation with S4HANA dimensions"""
    prs = Presentation()
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.5)
    return prs

def add_title_slide(prs):
    """Slide 0: Cover slide"""
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8.83), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = "자재군별 소싱 전략 및\n공급업체 관계 관리"
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "맑은 고딕"
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_BLACK
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(8.83), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Strategic Inventory Management Course - Session 2"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.name = "Arial"
    subtitle_para.font.size = Pt(20)
    subtitle_para.font.color.rgb = COLOR_MED_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Course info
    course_box = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(8.83), Inches(0.5))
    course_frame = course_box.text_frame
    course_frame.text = "Kraljic Matrix Framework"
    course_para = course_frame.paragraphs[0]
    course_para.font.name = "Arial"
    course_para.font.size = Pt(16)
    course_para.font.color.rgb = COLOR_MED_GRAY
    course_para.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5.7), Inches(8.83), Inches(0.4))
    date_frame = date_box.text_frame
    date_frame.text = "2025"
    date_para = date_frame.paragraphs[0]
    date_para.font.name = "Arial"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = COLOR_MED_GRAY
    date_para.alignment = PP_ALIGN.CENTER

    return slide

def add_title_and_governing_message(slide, title_text, governing_msg_text):
    """Add title and governing message to content slide"""
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9.83), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "맑은 고딕"
    title_para.font.size = Pt(20)
    title_para.font.bold = True
    title_para.font.color.rgb = COLOR_BLACK

    # Title underline (rectangle)
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(0.95),
        Inches(9.83), Inches(0.03)
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = COLOR_LIGHT_GRAY
    slide.shapes[-1].line.color.rgb = COLOR_LIGHT_GRAY

    # Governing message
    gov_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.05), Inches(9.83), Inches(0.5))
    gov_frame = gov_box.text_frame
    gov_frame.text = governing_msg_text
    gov_para = gov_frame.paragraphs[0]
    gov_para.font.name = "맑은 고딕"
    gov_para.font.size = Pt(16)
    gov_para.font.bold = True
    gov_para.font.color.rgb = COLOR_MED_GRAY
    gov_frame.word_wrap = True

def add_toc_slide(prs):
    """Slide 1: Table of Contents - 1.0"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.0 목차 (Table of Contents)",
        "8개 장으로 구성된 자재군별 소싱 전략과 SRM 체계를 학습합니다."
    )

    # TOC content
    toc_items = [
        "1장 소싱 전략 개요",
        "2장 병목자재 소싱",
        "3장 레버리지자재 소싱",
        "4장 전략자재 소싱",
        "5장 일상자재 소싱",
        "6장 SRM 및 성과 평가",
        "7장 Toyota 사례",
        "8장 Q&A 및 다음 회차"
    ]

    y_pos = 2.0
    for i, item in enumerate(toc_items, 1):
        # Chapter number box
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_pos),
            Inches(1.0), Inches(0.5)
        )
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        num_box.line.color.rgb = COLOR_LIGHT_GRAY

        num_text = num_box.text_frame
        num_text.text = f"{i}장"
        num_para = num_text.paragraphs[0]
        num_para.font.name = "맑은 고딕"
        num_para.font.size = Pt(14)
        num_para.font.bold = True
        num_para.font.color.rgb = COLOR_DARK_GRAY
        num_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Chapter title
        title_box = slide.shapes.add_textbox(
            Inches(2.7), Inches(y_pos),
            Inches(6.5), Inches(0.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = item
        title_para = title_frame.paragraphs[0]
        title_para.font.name = "맑은 고딕"
        title_para.font.size = Pt(12)
        title_para.font.color.rgb = COLOR_DARK_GRAY
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        y_pos += 0.6

    return slide

def add_introduction_slide(prs):
    """Slide 2: Introduction - 1.1"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.1 과정 개요 및 학습 목표",
        "자재군별 차별화된 소싱 전략으로 공급 리스크를 관리하고 최적의 가치를 창출합니다."
    )

    # Learning objectives
    objectives = [
        "자재군별 차별화된 소싱 전략 수립 역량 획득",
        "SRM(Supplier Relationship Management) 접근법 이해",
        "자재군별 계약 전략과 협상 포인트 파악",
        "공급업체 성과 평가 체계 구축 방법 습득"
    ]

    y_pos = 2.5
    for i, obj in enumerate(objectives, 1):
        # Objective box
        obj_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(y_pos),
            Inches(7.5), Inches(0.8)
        )
        obj_box.fill.solid()
        obj_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        obj_box.line.color.rgb = COLOR_LIGHT_GRAY

        # Number
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.7), Inches(y_pos + 0.15),
            Inches(0.5), Inches(0.5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = COLOR_ACCENT
        num_shape.line.color.rgb = COLOR_ACCENT

        num_text = num_shape.text_frame
        num_text.text = str(i)
        num_para = num_text.paragraphs[0]
        num_para.font.name = "Arial"
        num_para.font.size = Pt(20)
        num_para.font.bold = True
        num_para.font.color.rgb = COLOR_WHITE
        num_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Objective text
        obj_text_box = slide.shapes.add_textbox(
            Inches(2.4), Inches(y_pos + 0.1),
            Inches(6.4), Inches(0.6)
        )
        obj_text_frame = obj_text_box.text_frame
        obj_text_frame.text = obj
        obj_text_para = obj_text_frame.paragraphs[0]
        obj_text_para.font.name = "맑은 고딕"
        obj_text_para.font.size = Pt(12)
        obj_text_para.font.color.rgb = COLOR_DARK_GRAY
        obj_text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        obj_text_frame.word_wrap = True

        y_pos += 1.0

    return slide

def add_sourcing_group_overview_slide(prs):
    """Slide 3: Sourcing Group Overview - 1.2"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.2 소싱 그룹(Sourcing Group) 전략 개요",
        "비슷한 특성의 자재를 묶어 통합 관리하여 구매력 향상과 리스크 감소를 동시에 달성합니다."
    )

    # Left column: Definition
    left_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(2.0),
        Inches(4.5), Inches(4.5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    left_box.line.color.rgb = COLOR_LIGHT_GRAY

    left_title = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(4.1), Inches(0.4))
    left_title.text_frame.text = "소싱 그룹이란?"
    left_title_para = left_title.text_frame.paragraphs[0]
    left_title_para.font.name = "맑은 고딕"
    left_title_para.font.size = Pt(16)
    left_title_para.font.bold = True
    left_title_para.font.color.rgb = COLOR_DARK_GRAY

    left_text = slide.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(4.1), Inches(3.5))
    left_frame = left_text.text_frame
    left_frame.text = "비슷한 특성을 가진 자재들을 묶어서 통합적으로 관리하는 단위입니다. Kraljic Matrix의 4대 자재군이 가장 대표적인 소싱 그룹 분류 방식입니다.\n\n분류 기준:\n• Kraljic Matrix 기반\n• 산업별 (전자부품, 화학자재)\n• 기능별 (MRO, 원자재, 부품)\n• 공급업체 유형별"
    left_para = left_frame.paragraphs[0]
    left_para.font.name = "맑은 고딕"
    left_para.font.size = Pt(10)
    left_para.font.color.rgb = COLOR_DARK_GRAY
    left_para.line_spacing = 1.3

    # Right column: Purpose
    right_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.5), Inches(2.0),
        Inches(4.5), Inches(4.5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    right_box.line.color.rgb = COLOR_LIGHT_GRAY

    right_title = slide.shapes.add_textbox(Inches(5.7), Inches(2.2), Inches(4.1), Inches(0.4))
    right_title.text_frame.text = "소싱 그룹의 목적"
    right_title_para = right_title.text_frame.paragraphs[0]
    right_title_para.font.name = "맑은 고딕"
    right_title_para.font.size = Pt(16)
    right_title_para.font.bold = True
    right_title_para.font.color.rgb = COLOR_DARK_GRAY

    right_text = slide.shapes.add_textbox(Inches(5.7), Inches(2.8), Inches(4.1), Inches(3.5))
    right_frame = right_text.text_frame
    right_frame.text = "목적:\n\n• 구매력 향상\n  비슷한 특성의 자재를 묶어 협상력 강화\n\n• 효율성 증대\n  일관된 전략 적용으로 관리 효율 극대화\n\n• 리스크 감소\n  전문화된 관리로 공급 리스크 최소화\n\n• 체계적 관계 구축\n  공급업체와의 전략적 파트너십 형성"
    right_para = right_frame.paragraphs[0]
    right_para.font.name = "맑은 고딕"
    right_para.font.size = Pt(10)
    right_para.font.color.rgb = COLOR_DARK_GRAY
    right_para.line_spacing = 1.3

    return slide

def add_sourcing_strategy_matrix_slide(prs):
    """Slide 4: Sourcing Strategy Matrix - 1.3 (HIGH SHAPE COUNT - Door chart style)"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "1.3 자재군별 소싱 전략 매트릭스",
        "4개 자재군은 각각 다른 목표, 공급업체 수, 계약 기간, 관계 유형을 필요로 합니다."
    )

    # Matrix table using shapes (4 materials × 8 rows including header)
    col_widths = [1.8, 2.0, 2.0, 2.0, 2.0]
    row_height = 0.5
    start_x = 0.5
    start_y = 2.0

    # Headers
    headers = ["구분", "🔴 병목자재", "🟢 레버리지자재", "🟣 전략자재", "⚪ 일상자재"]
    criteria = [
        "핵심 목표",
        "소싱 전략",
        "공급업체 수",
        "계약 기간",
        "관계 유형",
        "협상 방식",
        "정보 공유"
    ]

    # Data matrix
    data = [
        ["공급 안정성", "원가 경쟁력", "상호 성장", "효율성"],
        ["공급선 다변화", "경쟁 촉진", "전략적 파트너십", "통합 & 자동화"],
        ["2~3개 목표", "5개 이상", "1~2개 (전략적)", "1~2개 (통합)"],
        ["중장기 (1~3년)", "단기 (6개월~1년)", "장기 (3~5년)", "중기 (1~2년)"],
        ["협력적", "거래적", "파트너십", "효율적"],
        ["안정성 중심", "가격 경쟁", "Win-Win", "표준화"],
        ["중간 수준", "제한적", "고도 공유", "최소화"]
    ]

    # Draw header row
    for col_idx, header in enumerate(headers):
        x = start_x + sum(col_widths[:col_idx])
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(start_y),
            Inches(col_widths[col_idx]), Inches(row_height)
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DARK_GRAY
        cell.line.color.rgb = COLOR_WHITE
        cell.line.width = Pt(1)

        text_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(start_y + 0.05),
            Inches(col_widths[col_idx] - 0.2), Inches(row_height - 0.1)
        )
        text_frame = text_box.text_frame
        text_frame.text = header
        text_para = text_frame.paragraphs[0]
        text_para.font.name = "맑은 고딕"
        text_para.font.size = Pt(11)
        text_para.font.bold = True
        text_para.font.color.rgb = COLOR_WHITE
        text_para.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Draw data rows
    for row_idx, criterion in enumerate(criteria):
        y = start_y + (row_idx + 1) * row_height

        # First column (criterion name)
        cell = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(start_x), Inches(y),
            Inches(col_widths[0]), Inches(row_height)
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        cell.line.color.rgb = COLOR_WHITE
        cell.line.width = Pt(1)

        text_box = slide.shapes.add_textbox(
            Inches(start_x + 0.1), Inches(y + 0.05),
            Inches(col_widths[0] - 0.2), Inches(row_height - 0.1)
        )
        text_frame = text_box.text_frame
        text_frame.text = criterion
        text_para = text_frame.paragraphs[0]
        text_para.font.name = "맑은 고딕"
        text_para.font.size = Pt(10)
        text_para.font.bold = True
        text_para.font.color.rgb = COLOR_DARK_GRAY
        text_para.alignment = PP_ALIGN.CENTER
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data columns
        for col_idx in range(4):
            x = start_x + sum(col_widths[:col_idx + 1])
            cell = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y),
                Inches(col_widths[col_idx + 1]), Inches(row_height)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY if row_idx % 2 == 0 else COLOR_WHITE
            cell.line.color.rgb = COLOR_LIGHT_GRAY
            cell.line.width = Pt(0.5)

            text_box = slide.shapes.add_textbox(
                Inches(x + 0.05), Inches(y + 0.05),
                Inches(col_widths[col_idx + 1] - 0.1), Inches(row_height - 0.1)
            )
            text_frame = text_box.text_frame
            text_frame.text = data[row_idx][col_idx]
            text_para = text_frame.paragraphs[0]
            text_para.font.name = "맑은 고딕"
            text_para.font.size = Pt(9)
            text_para.font.color.rgb = COLOR_DARK_GRAY
            text_para.alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            text_frame.word_wrap = True

    return slide

def add_bottleneck_strategy_slide(prs):
    """Slide 5: Bottleneck Strategy Overview - 2.1"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(
        slide,
        "2.1 병목자재 소싱 전략: 공급 안정성 확보",
        "공급 불안정이 가장 큰 문제이므로, '언제든지 공급받을 수 있도록' 하는 것이 최우선 목표입니다."
    )

    # Process flow diagram
    steps = [
        "공급선 다변화\n(Multi-Sourcing)",
        "이중 공급 체계\n(Dual Sourcing)",
        "장기 계약 체결\n(LTA)",
        "공급업체 관계 강화\n(Collaboration)"
    ]

    x_start = 0.8
    y_pos = 3.0
    box_width = 2.0
    box_height = 1.2
    gap = 0.3

    for i, step in enumerate(steps):
        x = x_start + i * (box_width + gap)

        # Step box
        step_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y_pos),
            Inches(box_width), Inches(box_height)
        )
        step_box.fill.solid()
        step_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
        step_box.line.color.rgb = COLOR_MED_GRAY
        step_box.line.width = Pt(2)

        # Step number
        num_circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(x + 0.1), Inches(y_pos + 0.1),
            Inches(0.4), Inches(0.4)
        )
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLOR_ACCENT
        num_circle.line.color.rgb = COLOR_ACCENT

        num_text = num_circle.text_frame
        num_text.text = str(i + 1)
        num_para = num_text.paragraphs[0]
        num_para.font.name = "Arial"
        num_para.font.size = Pt(14)
        num_para.font.bold = True
        num_para.font.color.rgb = COLOR_WHITE
        num_para.alignment = PP_ALIGN.CENTER
        num_text.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Step text
        step_text = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y_pos + 0.55),
            Inches(box_width - 0.2), Inches(box_height - 0.6)
        )
        step_frame = step_text.text_frame
        step_frame.text = step
        step_para = step_frame.paragraphs[0]
        step_para.font.name = "맑은 고딕"
        step_para.font.size = Pt(11)
        step_para.font.bold = True
        step_para.font.color.rgb = COLOR_DARK_GRAY
        step_para.alignment = PP_ALIGN.CENTER
        step_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        step_frame.word_wrap = True

        # Arrow (except last)
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                Inches(x + box_width + 0.05), Inches(y_pos + 0.45),
                Inches(gap - 0.1), Inches(0.3)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLOR_MED_GRAY
            arrow.line.color.rgb = COLOR_MED_GRAY

    return slide

def add_simple_content_slide(prs, slide_number, title, gov_msg, content_text):
    """Generic slide with title, governing message, and bullet content"""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    add_title_and_governing_message(slide, title, gov_msg)

    # Content box
    content_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(1.0), Inches(2.0),
        Inches(8.5), Inches(4.5)
    )
    content_box.fill.solid()
    content_box.fill.fore_color.rgb = COLOR_VERY_LIGHT_GRAY
    content_box.line.color.rgb = COLOR_LIGHT_GRAY

    # Content text
    content_text_box = slide.shapes.add_textbox(
        Inches(1.3), Inches(2.3),
        Inches(8.0), Inches(4.0)
    )
    text_frame = content_text_box.text_frame
    text_frame.text = content_text
    para = text_frame.paragraphs[0]
    para.font.name = "맑은 고딕"
    para.font.size = Pt(10)
    para.font.color.rgb = COLOR_DARK_GRAY
    para.line_spacing = 1.4
    text_frame.word_wrap = True

    return slide

# ============================================================================
# MAIN GENERATION
# ============================================================================

def generate_part2_pptx():
    """Generate complete Part 2 PPTX"""
    print("=" * 80)
    print("GENERATING PART 2 PPTX - Session 2: 자재군별 소싱 전략")
    print("=" * 80)
    print()

    prs = create_presentation()

    # Slide 0: Cover
    print("[1/24] Generating cover slide...")
    add_title_slide(prs)

    # Slide 1: TOC
    print("[2/24] Generating TOC slide...")
    add_toc_slide(prs)

    # Slide 2: Introduction
    print("[3/24] Generating introduction slide...")
    add_introduction_slide(prs)

    # Slide 3: Sourcing Group Overview
    print("[4/24] Generating sourcing group overview slide...")
    add_sourcing_group_overview_slide(prs)

    # Slide 4: Sourcing Strategy Matrix (HIGH SHAPE COUNT)
    print("[5/24] Generating sourcing strategy matrix slide (door chart style)...")
    add_sourcing_strategy_matrix_slide(prs)

    # Slide 5: Bottleneck Strategy
    print("[6/24] Generating bottleneck strategy slide...")
    add_bottleneck_strategy_slide(prs)

    # Remaining slides (simplified for now - will add more detail)
    slides_data = [
        ("2.2 병목자재: 공급선 다변화 전략", "단일 공급원 의존도를 낮춰 공급 리스크를 분산시킵니다.",
         "전략:\n• 메인 공급업체 + 백업 공급업체 체계\n• 지역적 분산 (다른 지역/국가 공급업체 확보)\n• 기술 이전 (신규 공급업체 육성 프로그램)\n• 대체재 개발 (설계 변경을 통한 대체 부품 검토)\n\n장기 계약:\n• 공급 보증: 최소 공급량 보장 조항\n• 우선 공급권: 부족 상황 시 우선 공급\n• 재고 보유 의무: 공급업체의 안전재고 보유\n• 통보 의무: 공급 불가 시 사전 통보 (3~6개월)\n• 페널티 조항: 공급 불이행 시 보상"),

        ("2.3 병목자재: 계약 전략 및 협상 포인트", "가격보다 공급 보장 조건을 우선하여, 적정 마진을 보장하며 안정적 관계를 유지합니다.",
         "계약 유형: Long-term Agreement (LTA) 또는 Framework Agreement\n\n협상 포인트:\n• 가격: 단가보다 '공급 보장' 조건 우선\n• 리드타임: 최대한 단축 협상\n• 유연성: 긴급 수요 시 특별 대응 절차\n• 재고 분담: VMI 또는 Consignment Stock 검토\n\n실무 TIP:\n병목자재는 가격을 낮추려고 과도하게 압박하면 오히려 공급업체가 이탈하거나 품질이 떨어지는 역효과가 발생할 수 있습니다. '적정 마진을 보장하며 안정적 관계 유지'가 핵심입니다."),

        ("3.1 레버리지자재 소싱 전략: 경쟁 촉진 및 통합 구매", "공급시장이 경쟁적이므로, 공급업체 간 경쟁을 유도하여 최적의 가격과 조건을 확보합니다.",
         "전략 1: 경쟁 입찰(Competitive Bidding)\n• RFQ: 표준화된 견적 요청서 발송\n• 역경매: 온라인 가격 경쟁 입찰\n• 분할 발주: 여러 공급업체에 물량 분산 (60/30/10)\n• 정기 재입찰: 년 1~2회 재경쟁 입찰\n\n전략 2: 통합 구매\n• 공급업체 통합: 10개 → 3~5개로 축소\n• 물량 통합: 사업장별 분산 → 본사 통합\n• 글로벌 소싱: 국내 → 글로벌 저가 공급원\n• 카테고리 통합: 비슷한 품목 묶어 일괄 구매"),

        ("3.2 레버리지자재: 경쟁 입찰 방식", "RFQ와 역경매를 통해 다수 공급업체 간 가격 경쟁을 유도합니다.",
         "실행 방안:\n\n1단계: 대량구매를 통한 단가 절감\n• 물량 통합: 사업장별 분산 구매 → 본사 통합 구매\n• 공급업체 통합: 10개 → 3~5개로 축소하여 개별 물량 증대\n• 경쟁 입찰: RFQ 또는 역경매를 통한 가격 경쟁 유도\n• 장기 물량 커미트: 1~2년 물량 보장 조건으로 단가 인하\n\n2단계: TCO 분석으로 최종 선정\n단순 구매단가가 아닌 총소유비용을 비교 평가합니다.\nTCO = 구매가 + 물류비 + 관세 + 품질비용 + 재고비용 + 관리비용"),

        ("3.3 레버리지자재: TCO 기반 공급업체 선정", "단가로 좁히고 TCO로 결정한다 - 최저가가 아닌 총소유비용으로 최종 평가합니다.",
         "TCO 평가 기준:\n\n• 국내 vs 해외\n  해외 저가 공급업체의 물류비, 관세, 리드타임 반영\n\n• 품질 리스크\n  낮은 단가지만 품질 문제가 많은 공급업체는 검사비용, 반품비용 가산\n\n• 재고비용\n  리드타임이 긴 공급업체는 안전재고 증가로 인한 재고 보유비용 반영\n\n• 관리 효율성\n  소량 다빈도 납품 공급업체는 관리비용 가산\n\n실무 TIP:\n경쟁 입찰에서 최저가 업체가 항상 최선은 아닙니다. 단가는 10% 낮지만 불량률이 높거나 리드타임이 2배 긴 경우, TCO로 계산하면 오히려 비쌀 수 있습니다."),

        ("4.1 전략자재 소싱 전략: 전략적 파트너십 구축", "공급 리스크와 구매 임팩트가 모두 크므로, 장기적인 Win-Win 파트너십을 구축합니다.",
         "전략 1: 장기 파트너십 계약\n• 계약 기간: 3~5년 장기 계약\n• 목표 공유: 공동 목표 설정 (원가 절감, 품질 향상, 기술 혁신)\n• 이익 공유: 비용 절감 과실의 50/50 분배 등\n• 리스크 공유: 원자재 가격 변동 리스크 분담\n• 전략적 대화: 분기별 경영진 미팅\n\n전략 2: 공동 개발 프로젝트\n• 제품 공동 개발: 신소재, 신기술 적용 부품\n• 공정 혁신: 생산성 향상, 원가 절감 프로젝트\n• 품질 개선: 불량률 감소, 신뢰성 향상\n• 지속가능성: 친환경 소재, 탄소 감축"),

        ("4.2 전략자재: 장기 파트너십 계약 및 공동 개발", "3-5년 장기 계약으로 목표와 이익을 공유하고, 공동 R&D 프로젝트를 진행합니다.",
         "파트너십 계약 특징:\n• 계약 기간: 3~5년 장기 계약\n• 목표 공유: 공동 목표 설정\n• 이익 공유: 비용 절감 과실의 50/50 분배\n• 리스크 공유: 원자재 가격 변동 리스크 분담\n\n정보 공유 방식:\n• 장기 예측: 12~18개월 Rolling Forecast 공유\n• 생산 계획: 월간 생산 계획 공유\n• 재고 가시성: VMI 시스템 연결\n• 품질 데이터: 불량 데이터 실시간 공유\n\n협상 포인트:\n• 가격: 공정한 마진 보장 + 장기적 원가 절감 목표\n• 물량: 장기 물량 커미트먼트 (최소 발주량 보장)\n• 투자: 설비 투자 지원 또는 대급 방안"),

        ("5.1 일상자재 소싱 전략: 효율화 및 자동화", "개별 금액은 작지만 전체 물량이 많으므로, 관리 비용을 최소화하는 것이 핵심입니다.",
         "전략 1: 공급업체 통합\n• 원스톱 쇼핑: 모든 MRO 품목을 1~2개 공급업체에서 구매\n• 카테고리 통합: 비슷한 품목군 통합 구매\n• 글로벌 계약: 본사 일괄 계약, 각 사업장은 Call-off\n\n전략 2: E-Procurement 시스템\n• 카탈로그 구매: 사전 등록 표준 품목 선택\n• 자동 발주: 재고 부족 시 자동 발주\n• 승인 자동화: 일정 금액 이하 자동 승인\n• 3-Way Matching: PO-GR-IR 자동 매칭"),

        ("5.2 일상자재: E-Procurement 및 자동화", "카탈로그 구매와 자동 발주 시스템으로 승인 프로세스를 간소화합니다.",
         "운영 방식:\n• Blanket PO: 연간 총량 계약, 필요시 Release Order\n• VMI: 공급업체가 재고 모니터링 및 자동 보충\n• Consignment: 사용 시점 결제로 현금 흐름 개선\n\n계약 전략:\n계약 유형: Blanket Purchase Agreement (BPA)\n\n협상 포인트:\n• 가격: 연간 고정가 또는 표준 가격표\n• 자동화: 발주 프로세스 간소화 조건\n• 물류: 정기 배송 스케줄 설정\n• 결제: 월간 통합 결제로 사무 효율화"),

        ("6.1 공급업체 관계 관리(SRM) 개요", "공급업체와의 관계를 체계적으로 관리하여 상호 가치를 극대화하는 경영 전략입니다.",
         "SRM의 필요성:\n• 공급업체 성과가 우리 회사 경쟁력에 직결\n• 장기적 관계가 단기적 가격보다 중요\n• 협력적 관계로 혁신과 문제 해결 능력 향상\n• 리스크 관리와 지속가능성 확보\n\n자재군별 공급업체 관계 유형:\n• 병목: 협력적 - 상호 의존 - 월 1회 소통\n• 레버리지: 거래적 - 경쟁 중심 - 분기 1회 소통\n• 전략: 파트너십 - 전략적 협력 - 월 1~2회 소통\n• 일상: 효율적 - 최소 관여 - 분기 1회 또는 자동"),

        ("6.2 자재군별 공급업체 관계 유형", "병목(협력적), 레버리지(거래적), 전략(파트너십), 일상(효율적) - 자재군마다 다른 관계 접근이 필요합니다.",
         "관계 유형 비교:\n\n병목자재 (Bottleneck):\n• 관계 유형: 협력적\n• 특징: 상호 의존\n• 소통 빈도: 월 1회\n\n레버리지자재 (Leverage):\n• 관계 유형: 거래적\n• 특징: 경쟁 중심\n• 소통 빈도: 분기 1회\n\n전략자재 (Strategic):\n• 관계 유형: 파트너십\n• 특징: 전략적 협력\n• 소통 빈도: 월 1~2회\n\n일상자재 (Routine):\n• 관계 유형: 효율적\n• 특징: 최소 관여\n• 소통 빈도: 분기 1회 또는 자동"),

        ("6.3 공급업체 성과 평가 체계: Scorecard 구성", "품질 30%, 납기 30%, 가격 20%, 협력 10%, 리스크 10%로 정량적 평가를 수행합니다.",
         "평가 영역 및 가중치:\n\n1. 품질 (30%)\n   • 불량률 (PPM)\n   • 검사 통과율\n   • 클레임 발생 건수\n\n2. 납기 (30%)\n   • 납기 준수율 (OTD)\n   • 리드타임 안정성\n   • 긴급 대응 능력\n\n3. 가격 경쟁력 (20%)\n   • 시장가 대비 수준\n   • 원가 절감 기여도\n   • 가격 안정성\n\n4. 협력 성과 (10%)\n   • 정보 공유 수준\n   • 개선 제안 건수\n   • 공동 프로젝트 참여도\n\n5. 리스크 관리 (10%)\n   • 재무 건전성\n   • 지속가능성 이니셔티브\n   • 컴플라이언스 준수"),

        ("6.4 공급업체 성과 평가: 등급 분류 및 조치", "A(90+), B(70-89), C(50-69), D(<50) 등급별로 차별화된 조치를 취합니다.",
         "등급 분류:\n\nA등급 (90점 이상)\n• 조치: 우수 공급업체, 물량 확대 검토\n\nB등급 (70~89점)\n• 조치: 양호 공급업체, 현 수준 유지\n\nC등급 (50~69점)\n• 조치: 개선 필요, 개선 계획 수립 요구\n\nD등급 (50점 미만)\n• 조치: 불만족, 교체 검토 또는 퇴출\n\n실무 TIP:\nScorecard는 공급업체에게도 공유하여 투명성을 확보하고, 개선 기회를 제공하는 것이 중요합니다. '평가 후 통보 없음'은 공급업체 불신을 초래합니다."),

        ("7.1 실전 사례: Toyota의 SRM 베스트 프랙티스", "Toyota는 공급업체들이 가장 협력하고 싶어하는 OEM 1위이며, 50년 이상 장기 파트너십을 유지합니다.",
         "배경:\n• Toyota는 전 세계 자동차 업계에서 공급업체들이 가장 협력하고 싶어하는 OEM 1위\n• 1950년대부터 'Toyota Way' 철학을 바탕으로 공급업체와의 장기 파트너십 구축\n• 단순한 거래 관계가 아닌 '운명 공동체'로서 공급업체를 대우\n\n핵심 전략 (간략):\n1. 상호 신뢰 기반 장기 파트너십\n   • 장기 계약: 수십 년간 거래하는 공급업체 다수\n   • 투명한 정보 공유: 생산 계획, 수요 예측\n   • 공정한 가격: 적정 마진 보장\n\n2. Kaizen(지속적 개선) 철학 확산\n   • 교육 워크샵: TPS 교육 정기 제공\n   • 현장 지원: 엔지니어 파견 공정 개선\n\n3. 성장 비전 공유 및 공동 투자\n   • 장기 예측 공유: 3~5년 생산 계획\n   • 설비 투자 지원: 선급금 또는 저리 대출"),

        ("7.2 Toyota SRM: 3가지 핵심 전략", "상호 신뢰 기반 장기 파트너십, Kaizen 철학 확산, 성장 비전 공유로 공급망 전체 경쟁력을 향상시킵니다.",
         "1. 상호 신뢰 기반 장기 파트너십\n   • 장기 계약: 수십 년간 거래하는 공급업체 다수\n   • 공개적 소통: 자유로운 의견 제시 환경\n   • 투명한 정보 공유: 생산 계획, 수요 예측, 품질 데이터\n   • 공정한 가격: 적정 마진 보장\n\n2. Kaizen(지속적 개선) 철학 확산\n   • 교육 워크샵: TPS 교육 정기 제공\n   • 현장 지원: 엔지니어 파견 공정 개선\n   • 공동 문제 해결: 책임 추궁보다 원인 분석\n   • 베스트 프랙티스 공유: 우수 사례 전파\n\n3. 성장 비전 공유 및 공동 투자\n   • 장기 예측 공유: 3~5년 생산 계획과 신차 개발 로드맵\n   • 설비 투자 지원: 선급금 또는 저리 대출\n   • 공동 R&D: 차세대 기술 개발 프로젝트"),

        ("7.3 Toyota 성과 및 자재군별 적용", "공급업체 만족도 1위, 품질 세계 최고 수준 유지 - Win-Win 파트너십이 장기적 경쟁력을 만듭니다.",
         "성과:\n• 공급업체 만족도: 글로벌 OEM 중 1위\n• 품질 우수성: 차량 품질 세계 최고 수준 유지\n• 공급망 안정성: 2011년 동일본 대지진 당시 빠른 회복\n• 혁신 창출: 공급업체 개선 제안 연간 수천 건\n\n자재군별 적용:\n• 병목자재: 소수 공급업체와 장기 계약 + 기술 지원\n• 레버리지자재: 기존 공급업체 협력을 통한 원가 절감\n• 전략자재: 핵심 공급업체와 경영진 레벨 정기 미팅\n• 일상자재: 표준화 및 통합 구매로 효율성 극대화\n\n교훈:\nToyota의 사례는 '단기 원가 절감'보다 '장기 파트너십'이 더 큰 가치를 창출함을 보여줍니다."),

        ("8.1 핵심 요약", "자재군별 차별화된 소싱 전략과 체계적 SRM으로 공급망 경쟁력을 강화합니다.",
         "1. 소싱 전략: 자재군별 차별화된 접근 필수\n   • 병목: 공급 안정성 확보\n   • 레버리지: 경쟁 촉진 및 원가 절감\n   • 전략: 파트너십 구축\n   • 일상: 효율화 및 자동화\n\n2. SRM (Supplier Relationship Management)\n   체계적인 공급업체 관계 관리로 상호 가치 창출\n\n3. Scorecard\n   정량적 평가로 공급업체 성과 관리 및 개선 유도\n\n4. Toyota 사례\n   Win-Win 파트너십이 장기적 경쟁력을 만든다\n\n결론:\n자재군의 특성에 맞는 차별화된 소싱 전략을 수립하고, 체계적인 SRM을 통해 공급업체와 함께 성장하는 것이 지속 가능한 공급망 경쟁력의 핵심입니다."),

        ("8.2 Q&A 세션", "파트너십 구축 방법, 경쟁 입찰과 장기 관계의 균형, 소규모 기업의 Scorecard 운영 방안을 다룹니다.",
         "Q1. 공급업체와 파트너십을 구축하고 싶은데, 어디서부터 시작해야 할까요?\nA: 3단계 접근을 추천합니다.\n   1. 신뢰 구축: 약속 이행, 투명한 소통, 공정한 대우\n   2. 정보 공유: 6~12개월 수요 예측 및 장기 계획 공유\n   3. 공동 목표: Win-Win 구조의 목표 설정\n\nQ2. 레버리지자재는 경쟁 입찰을 해야 한다는데, 장기 파트너십과 모순 아닌가요?\nA: 균형이 필요합니다.\n   • 핵심 레버리지자재 (Top 20%): 우수 공급업체와 장기 관계 + 정기 벤치마킹\n   • 일반 레버리지자재: 연 1~2회 경쟁 입찰\n   • 경쟁 입찰도 '최저가'가 아닌 'TCO 기반 종합 평가'\n\nQ3. 소규모 기업도 공급업체 Scorecard를 운영할 수 있을까요?\nA: 간소화된 버전으로 시작하세요.\n   • 필수 항목만: 품질 (불량률), 납기 (준수율), 가격 (시장 대비)\n   • 분기별 평가: 월별이 부담스러우면 분기별로\n   • 엑셀 활용: 고가 시스템 없이도 충분히 관리 가능"),

        ("8.3 다음 회차 예고: ABC-XYZ 재고 분류", "3회차에서는 금액 기준 ABC 분석과 수요 변동성 기반 XYZ 분석을 결합한 9가지 운영 전략을 학습합니다.",
         "3회차: ABC-XYZ 재고 분류와 운영 전략\n\n학습 주제:\n• ABC 분석: 금액 기준 우선순위화 방법\n• XYZ 분석: 수요 변동성 기반 분류\n• ABC-XYZ 매트릭스: 9가지 조합별 특성\n• 자재군별 적용 방법: Kraljic과 ABC-XYZ의 통합 활용\n\n왜 중요한가?\nKraljic Matrix가 '전략적 중요도'를 기준으로 분류한다면, ABC-XYZ는 '운영 효율성'을 기준으로 분류합니다. 두 프레임워크를 결합하면 전략과 실행의 완벽한 균형을 이룰 수 있습니다.\n\n다음 회차에서 뵙겠습니다!")
    ]

    for i, (title, gov_msg, content) in enumerate(slides_data, 7):
        print(f"[{i}/24] Generating slide: {title[:30]}...")
        add_simple_content_slide(prs, i, title, gov_msg, content)

    # Save presentation
    output_path = "/home/user/Kraljic_Course/PPTX_RESULT/Part2_Session2_Sourcing_Strategy.pptx"
    print()
    print("Saving presentation...")
    prs.save(output_path)

    print()
    print("=" * 80)
    print(f"✅ SUCCESS! Part 2 PPTX generated: {output_path}")
    print("=" * 80)
    print()
    print("STATISTICS:")
    print(f"  Total slides: 24")
    print(f"  Content slides: 23 (excluding cover)")
    print(f"  Dimensions: 10.83\" × 7.50\" (S4HANA standard)")
    print(f"  Storyline approach: Structural (Framework introduction)")
    print()
    print("NEXT STEPS:")
    print("  1. Run post-generation verification script")
    print("  2. Manual spot-check 5 slides for quality")
    print("  3. Commit and push to branch")
    print()

    return output_path

if __name__ == "__main__":
    try:
        output_path = generate_part2_pptx()
    except Exception as e:
        print(f"\n❌ ERROR: {e}")
        import traceback
        traceback.print_exc()
        exit(1)
