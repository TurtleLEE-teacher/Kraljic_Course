#!/usr/bin/env python3
"""
Verify NEW Part 1 PPTX (Session 1 only)
"""

from pptx import Presentation
from pptx.util import Inches, Pt

def verify_pptx(filepath):
    """Verify PPTX meets S4HANA requirements"""
    print(f"\n{'='*70}")
    print(f"NEW Part 1 PPTX Compliance Verification")
    print(f"{'='*70}\n")

    prs = Presentation(filepath)

    # 1. Dimensions
    print(f"1. Slide Dimensions:")
    width_in = prs.slide_width.inches
    height_in = prs.slide_height.inches
    print(f"   Width: {width_in:.2f}\" (Expected: 10.83\")")
    print(f"   Height: {height_in:.2f}\" (Expected: 7.50\")")
    dim_ok = abs(width_in - 10.83) < 0.01 and abs(height_in - 7.50) < 0.01
    print(f"   Status: {'✓ PASS' if dim_ok else '✗ FAIL'}\n")

    # 2. Slide count
    print(f"2. Slide Count:")
    print(f"   Total: {len(prs.slides)} slides (Expected: 20)")
    print(f"   Status: {'✓ PASS' if len(prs.slides) == 20 else '✗ FAIL'}\n")

    # 3. Font sizes and governing messages
    print(f"3. Font Sizes and Governing Messages:")
    font_sizes = []
    gov_count = 0
    shape_counts = []

    for idx, slide in enumerate(prs.slides, 1):
        has_gov = False
        slide_fonts = []
        shape_count = 0

        for shape in slide.shapes:
            shape_count += 1

            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)
                            slide_fonts.append(run.font.size.pt)

                        # Check for governing message (16pt Bold)
                        if run.font.size and run.font.bold and run.font.size == Pt(16):
                            has_gov = True

        shape_counts.append(shape_count)

        if idx > 1 and has_gov:  # Skip cover slide
            gov_count += 1

        # Sample slide details (slides 2, 5, 11)
        if idx in [2, 5, 11]:
            print(f"\n   Slide {idx}:")
            print(f"     - Shapes: {shape_count}")
            print(f"     - Font sizes: {sorted(set([int(f) for f in slide_fonts]))}")
            print(f"     - Governing message: {'✓ YES' if has_gov else '✗ NO'}")

    print(f"\n   Content slides with governing messages: {gov_count}/19")
    print(f"   Status: {'✓ PASS' if gov_count >= 19 else '✗ FAIL'}")

    # 4. Font size statistics
    print(f"\n4. Font Size Statistics:")
    if font_sizes:
        unique_sizes = sorted(set([int(f) for f in font_sizes]))
        print(f"   Unique sizes: {unique_sizes}")
        print(f"   Smallest: {min(font_sizes):.0f}pt (S4HANA uses 6-11pt for body)")
        print(f"   Largest: {max(font_sizes):.0f}pt")

        # Count body text fonts (8-11pt)
        body_fonts = [f for f in font_sizes if 8 <= f <= 11]
        print(f"   Body fonts (8-11pt): {len(body_fonts)}/{len(font_sizes)} ({len(body_fonts)/len(font_sizes)*100:.1f}%)")

    # 5. Shape density
    print(f"\n5. Shape Density (Content Slides):")
    avg_shapes = sum(shape_counts[1:]) / (len(shape_counts) - 1) if len(shape_counts) > 1 else 0
    print(f"   Average shapes per slide: {avg_shapes:.1f}")
    print(f"   Target: 10-50+ shapes for high density")
    print(f"   Status: {'✓ PASS' if avg_shapes >= 10 else '⚠ LOW (need more diagrams)'}")

    print(f"\n{'='*70}")
    print(f"Verification Complete")
    print(f"{'='*70}\n")

if __name__ == "__main__":
    verify_pptx("/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Session1_Kraljic_Foundation.pptx")
