#!/usr/bin/env python3
"""
Part 1 PPTX Generator - Session 1 ONLY (NOT Sessions 1-3!)
S4HANA Design System Compliance:
- Fonts: 8-11pt body (SMALL fonts!), 16pt Bold governing messages
- Content Density: 85%+ (minimize whitespace)
- Visual Elements: 10-50+ shapes per slide (flowcharts, arrows, diagrams)
- Dimensions: 10.83" × 7.5"
- Monochrome: Black/White/Gray + 1 accent color
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# S4HANA Monochrome Color System
COLOR_BLACK = RGBColor(0, 0, 0)
COLOR_DARK_GRAY = RGBColor(51, 51, 51)
COLOR_MED_GRAY = RGBColor(102, 102, 102)
COLOR_LIGHT_GRAY = RGBColor(204, 204, 204)
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_ACCENT = RGBColor(26, 82, 118)  # Dark blue accent

# Kraljic colors (ONLY for matrix diagram)
COLOR_STRATEGIC = RGBColor(142, 68, 173)
COLOR_BOTTLENECK = RGBColor(230, 126, 34)
COLOR_LEVERAGE = RGBColor(39, 174, 96)
COLOR_ROUTINE = RGBColor(149, 165, 166)

def add_title_and_governing(slide, title_text, governing_text, slide_num):
    """Add title + governing message (S4HANA style)"""
    # Title: 20pt Bold
    title_box = slide.shapes.add_textbox(Inches(0.3), Inches(0.31), Inches(7.56), Inches(0.43))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title_text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = COLOR_BLACK

    # Governing message: 16pt Bold (NOT 14pt Italic!)
    gov_box = slide.shapes.add_textbox(Inches(0.3), Inches(1.01), Inches(10.32), Inches(0.63))
    gtf = gov_box.text_frame
    gtf.word_wrap = True
    gp = gtf.paragraphs[0]
    grun = gp.add_run()
    grun.text = governing_text
    grun.font.name = "맑은 고딕"
    grun.font.size = Pt(16)
    grun.font.bold = True
    grun.font.color.rgb = COLOR_DARK_GRAY

    # Slide number: bottom right
    num_box = slide.shapes.add_textbox(Inches(5.29), Inches(7.15), Inches(0.3), Inches(0.25))
    ntf = num_box.text_frame
    np = ntf.paragraphs[0]
    np.alignment = PP_ALIGN.RIGHT
    nrun = np.add_run()
    nrun.text = str(slide_num)
    nrun.font.name = "Arial"
    nrun.font.size = Pt(10)
    nrun.font.color.rgb = COLOR_MED_GRAY

def add_small_text(slide, left, top, width, height, text, font_size=9, bold=False):
    """Add text with SMALL font (8-11pt for high density)"""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = COLOR_BLACK
    return box

def add_shape_box(slide, left, top, width, height, text, fill_color, font_size=9):
    """Add shape with text for flowcharts/diagrams"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = COLOR_DARK_GRAY
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "맑은 고딕"
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE if fill_color != COLOR_WHITE else COLOR_BLACK
    return shape

def add_arrow(slide, x1, y1, x2, y2):
    """Add arrow connector"""
    connector = slide.shapes.add_connector(
        1,  # Straight connector
        Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = COLOR_DARK_GRAY
    connector.line.width = Pt(2)
    return connector

def create_presentation():
    """Create Part 1 PPTX (Session 1 ONLY - 20 slides)"""
    print("Generating Part 1 PPTX (Session 1 - 20 slides)...")
    print("S4HANA Compliance: Small fonts (8-11pt), 85%+ density, 10-50+ shapes")

    prs = Presentation()
    prs.slide_width = Inches(10.83)
    prs.slide_height = Inches(7.50)

    # Slide 1: Cover
    add_cover_slide(prs)

    # Slides 2-3: 학습 목표 & Agenda
    add_slide_2_objectives(prs)
    add_slide_3_agenda(prs)

    # Slides 4-6: JIT to JIC
    add_slide_4_jit_glory(prs)
    add_slide_5_jit_failure(prs)
    add_slide_6_jic_shift(prs)

    # Slides 7-8: Crisis cases
    add_slide_7_semiconductor_crisis(prs)
    add_slide_8_company_responses(prs)

    # Slides 9-11: Kraljic Matrix
    add_slide_9_kraljic_intro(prs)
    add_slide_10_two_axes(prs)
    add_slide_11_four_quadrants(prs)

    # Slides 12-16: 4 Material categories
    add_slide_12_bottleneck(prs)
    add_slide_13_leverage(prs)
    add_slide_14_strategic(prs)
    add_slide_15_routine(prs)
    add_slide_16_comparison_matrix(prs)

    # Slides 17-19: Planning methodologies
    add_slide_17_planning_overview(prs)
    add_slide_18_decision_tree(prs)
    add_slide_19_kpi_framework(prs)

    # Slide 20: Summary
    add_slide_20_summary(prs)

    output_path = "/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Session1_Kraljic_Foundation.pptx"
    prs.save(output_path)

    print(f"\n✓ PPTX saved: {output_path}")
    print(f"  - 20 slides (Session 1 ONLY)")
    print(f"  - S4HANA Monochrome Design")
    print(f"  - 16pt Bold governing messages")
    print(f"  - 8-11pt body fonts (high density)")
    print(f"  - 10-50+ shapes per slide (flowcharts, diagrams)")
    print(f"  - 10.83\" × 7.50\"")

    return output_path

def add_cover_slide(prs):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Background: solid accent color (NO gradient!)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_ACCENT

    # Main title: 48pt Bold
    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(2.0), Inches(9.5), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "전략적 재고운영 및 자재계획수립\n\nPart 1: Kraljic Matrix Foundation"
    run.font.name = "맑은 고딕"
    run.font.size = Pt(48)
    run.font.bold = True
    run.font.color.rgb = COLOR_WHITE

    # Subtitle: 20pt
    sub_box = slide.shapes.add_textbox(Inches(0.45), Inches(4.5), Inches(9.5), Inches(1.0))
    stf = sub_box.text_frame
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srun = sp.add_run()
    srun.text = "Session 1: JIT에서 JIC로의 전환과 Kraljic Matrix"
    srun.font.name = "맑은 고딕"
    srun.font.size = Pt(20)
    srun.font.color.rgb = COLOR_WHITE

    # Date
    date_box = slide.shapes.add_textbox(Inches(0.45), Inches(6.5), Inches(9.5), Inches(0.5))
    dtf = date_box.text_frame
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    drun = dp.add_run()
    drun.text = "2025년"
    drun.font.name = "Arial"
    drun.font.size = Pt(16)
    drun.font.color.rgb = COLOR_WHITE

def add_slide_2_objectives(prs):
    """Slide 2: 학습 목표"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "Session 1: 학습 목표",
        "JIT에서 JIC로의 패러다임 전환과 Kraljic Matrix를 활용한 자재 분류 방법을 학습합니다.",
        2
    )

    # Content area with HIGH DENSITY
    content_y = 2.0

    # 4 objectives in boxes with shapes
    objectives = [
        "JIT에서 JIC로의 패러다임 전환 배경과 필요성 이해",
        "전략적 재고운영의 핵심 개념과 차별화된 접근법 습득",
        "Kraljic Matrix를 활용한 자재 포트폴리오 분류 역량 확보",
        "자재군별 관리 철학과 자재계획 방법론의 전체 맵 이해"
    ]

    for i, obj in enumerate(objectives):
        y_pos = content_y + (i * 1.0)
        # Numbered box
        add_shape_box(slide, 0.5, y_pos, 0.6, 0.6, str(i+1), COLOR_ACCENT, 14)
        # Text with small font (9pt)
        add_small_text(slide, 1.3, y_pos + 0.15, 8.5, 0.6, obj, 9, False)

        # Add decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(y_pos + 0.7), Inches(9.5), Inches(0.05)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = COLOR_LIGHT_GRAY
        line.line.fill.background()

def add_slide_3_agenda(prs):
    """Slide 3: Agenda with flowchart"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "Session 1 Agenda",
        "5개 주제로 구성된 Session 1에서 JIT→JIC 전환과 Kraljic Matrix 전체를 다룹니다.",
        3
    )

    # Flowchart with 5 sections (LOTS of shapes for density)
    sections = [
        ("1. JIT→JIC\n전환", COLOR_DARK_GRAY),
        ("2. Kraljic\nMatrix", COLOR_ACCENT),
        ("3. 4대\n자재군", COLOR_MED_GRAY),
        ("4. 계획\n방법론", COLOR_DARK_GRAY),
        ("5. KPI\n프레임워크", COLOR_ACCENT)
    ]

    start_x = 0.8
    box_width = 1.6
    box_height = 1.2
    gap = 0.3
    y_pos = 2.5

    for i, (text, color) in enumerate(sections):
        x_pos = start_x + (i * (box_width + gap))
        add_shape_box(slide, x_pos, y_pos, box_width, box_height, text, color, 11)

        # Add arrow between boxes
        if i < len(sections) - 1:
            add_arrow(slide, x_pos + box_width, y_pos + 0.6, x_pos + box_width + gap, y_pos + 0.6)

    # Add detail boxes below (8pt small text for density)
    details = [
        "JIT의 몰락\nJIC의 부상",
        "2축 분석\n4사분면",
        "병목/레버리지\n전략/일상",
        "ROP/MRP\nLTP/VMI",
        "자재군별\n핵심 지표"
    ]

    for i, detail in enumerate(details):
        x_pos = start_x + (i * (box_width + gap))
        add_small_text(slide, x_pos, y_pos + 1.5, box_width, 0.8, detail, 8, False)

def add_slide_4_jit_glory(prs):
    """Slide 4: JIT의 영광 (with timeline diagram)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "JIT (Just-In-Time)의 영광: 1970-2010년",
        "도요타의 JIT 방식은 40년간 전 세계 제조업의 표준이었으나, 재고=낭비라는 극단적 해석이 문제였습니다.",
        4
    )

    # Timeline with shapes (HIGH DENSITY - 20+ shapes)
    timeline_y = 2.2

    # Timeline base line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0), Inches(timeline_y + 0.5), Inches(8.5), Inches(0.05)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_DARK_GRAY

    # Timeline events
    events = [
        (1.0, "1970\n도요타 JIT\n개발", COLOR_ACCENT),
        (3.5, "1980\nGE 도입\n30% 절감", COLOR_DARK_GRAY),
        (6.0, "1990\n전세계\n벤치마킹", COLOR_MED_GRAY),
        (8.5, "2010\nJIT 전성기\n\"재고=악\"", COLOR_DARK_GRAY)
    ]

    for x, text, color in events:
        # Event box
        add_shape_box(slide, x, timeline_y - 0.5, 1.3, 0.9, text, color, 8)
        # Connector to timeline
        connector = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x + 0.65), Inches(timeline_y + 0.4), Inches(0.03), Inches(0.5)
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = color
        connector.line.fill.background()

    # JIT 7 principles in small boxes (8pt - HIGH DENSITY)
    principles = [
        "Zero Inventory", "Pull System", "Kanban",
        "Continuous Flow", "Short Lead Time", "Perfect Quality", "Supplier Partnership"
    ]

    princ_y = 4.0
    for i, princ in enumerate(principles):
        x_pos = 1.0 + ((i % 4) * 2.2)
        y_pos = princ_y + ((i // 4) * 0.7)
        add_small_text(slide, x_pos, y_pos, 2.0, 0.5, f"• {princ}", 8, False)

    # Warning box at bottom
    add_shape_box(slide, 1.0, 5.8, 8.5, 1.0,
                   "⚠️ 문제: '재고=낭비'라는 극단적 해석 → 안전재고 Zero → 공급 충격에 취약",
                   COLOR_BOTTLENECK, 9)

def add_slide_5_jit_failure(prs):
    """Slide 5: JIT 실패 사례 (2020 pandemic)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "2020년 팬데믹: JIT의 치명적 약점 노출",
        "코로나19로 전 세계 공급망이 동시 중단되며 재고 Zero 전략의 위험성이 극명하게 드러났습니다.",
        5
    )

    # 3 crisis cases with detailed boxes (MANY shapes for density)
    cases = [
        {
            "title": "반도체 대란",
            "impact": "자동차 생산\n1,000만대 감소",
            "loss": "$210억\n손실",
            "cause": "안전재고\nZero",
            "color": COLOR_BOTTLENECK
        },
        {
            "title": "마스크 대란",
            "impact": "중국 생산\n50% 의존",
            "loss": "글로벌\n공급 중단",
            "cause": "단일 공급원\n집중",
            "color": COLOR_DARK_GRAY
        },
        {
            "title": "컨테이너 대란",
            "impact": "운임\n10배 폭등",
            "loss": "리드타임\n2배 증가",
            "cause": "예측 불가\n운송",
            "color": COLOR_MED_GRAY
        }
    ]

    case_y = 2.2
    for i, case in enumerate(cases):
        x_base = 0.8 + (i * 3.3)

        # Title box
        add_shape_box(slide, x_base, case_y, 3.0, 0.6, case["title"], case["color"], 11)

        # Detail boxes (8-9pt small fonts)
        add_shape_box(slide, x_base, case_y + 0.8, 1.4, 0.9, case["impact"], COLOR_LIGHT_GRAY, 8)
        add_shape_box(slide, x_base + 1.6, case_y + 0.8, 1.4, 0.9, case["loss"], COLOR_LIGHT_GRAY, 8)
        add_shape_box(slide, x_base, case_y + 1.9, 3.0, 0.9, case["cause"], COLOR_WHITE, 9)

    # Comparison table (JIT assumptions vs 2020 reality)
    table_y = 5.2
    table_data = [
        ("안정적 공급", "✅ 안정", "❌ 중단"),
        ("예측 가능 수요", "✅ 가능", "❌ 불가"),
        ("짧은 리드타임", "✅ 2-4주", "❌ 3-6개월")
    ]

    add_small_text(slide, 1.0, table_y, 2.5, 0.3, "JIT 전제 조건", 9, True)
    add_small_text(slide, 3.7, table_y, 2.5, 0.3, "2020 이전", 9, True)
    add_small_text(slide, 6.4, table_y, 2.5, 0.3, "2020 이후", 9, True)

    for i, (condition, before, after) in enumerate(table_data):
        y = table_y + 0.4 + (i * 0.4)
        add_small_text(slide, 1.0, y, 2.5, 0.3, condition, 8, False)
        add_small_text(slide, 3.7, y, 2.5, 0.3, before, 8, False)
        add_small_text(slide, 6.4, y, 2.5, 0.3, after, 8, False)

def add_slide_6_jic_shift(prs):
    """Slide 6: JIC 전환 (comparison diagram)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "JIC (Just-In-Case)로의 패러다임 전환",
        "2021년 이후 글로벌 기업들은 재고를 '전략적 자산'으로 재정의하고 차별화된 재고 전략을 도입했습니다.",
        6
    )

    # Before/After comparison with arrows (MANY shapes)
    # Left side: JIT
    add_shape_box(slide, 0.8, 2.0, 4.2, 0.8, "JIT (과거)", COLOR_DARK_GRAY, 14)

    jit_items = [
        "재고 = 낭비",
        "목표 = Zero 재고",
        "우선순위 = 효율성",
        "안전재고 = 1-2주",
        "공급업체 = Single Source"
    ]

    for i, item in enumerate(jit_items):
        y = 3.0 + (i * 0.65)
        add_small_text(slide, 1.0, y, 3.8, 0.5, f"• {item}", 9, False)

    # Center: Big arrow
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(5.2), Inches(3.5), Inches(0.8), Inches(1.5)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = COLOR_ACCENT
    arrow.line.fill.background()

    # Right side: JIC
    add_shape_box(slide, 6.2, 2.0, 4.2, 0.8, "JIC (현재/미래)", COLOR_ACCENT, 14)

    jic_items = [
        "재고 = 전략적 자산",
        "목표 = 최적 재고",
        "우선순위 = 회복력",
        "안전재고 = 차별화 (1주-6개월)",
        "공급업체 = Dual/Multi Source"
    ]

    for i, item in enumerate(jic_items):
        y = 3.0 + (i * 0.65)
        add_small_text(slide, 6.4, y, 3.8, 0.5, f"• {item}", 9, False)

def add_slide_7_semiconductor_crisis(prs):
    """Slide 7: 반도체 대란 상세 (case study with data)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "실전 사례: 2021년 글로벌 반도체 대란",
        "차량용 반도체 부족으로 1,000만대 생산 차질이 발생하며 JIT의 한계가 극명하게 드러났습니다.",
        7
    )

    # Timeline of crisis (horizontal flowchart with 15+ shapes)
    phases = [
        ("2020.03\n수요 급감\n예상", "주문 취소", COLOR_DARK_GRAY),
        ("2020.06\n수요 회복", "재주문\n시도", COLOR_BOTTLENECK),
        ("2020.09\n공급 부족", "생산 중단\n시작", COLOR_BOTTLENECK),
        ("2021.06\n최악의 상황", "1,000만대\n차질", COLOR_BOTTLENECK)
    ]

    phase_y = 2.2
    for i, (time, event, color) in enumerate(phases):
        x = 1.0 + (i * 2.2)
        # Time box
        add_shape_box(slide, x, phase_y, 1.8, 0.7, time, color, 9)
        # Event box
        add_small_text(slide, x, phase_y + 0.9, 1.8, 0.6, event, 8, True)

        # Arrow
        if i < len(phases) - 1:
            add_arrow(slide, x + 1.8, phase_y + 0.35, x + 2.2, phase_y + 0.35)

    # Impact data (numbers with small fonts - HIGH DENSITY)
    impacts = [
        ("생산 감소", "1,000만대", COLOR_BOTTLENECK),
        ("산업 손실", "$210억", COLOR_BOTTLENECK),
        ("영향 기업", "GM, Ford,\nVW 등", COLOR_DARK_GRAY),
        ("리드타임", "6-12개월", COLOR_MED_GRAY)
    ]

    impact_y = 4.5
    for i, (label, value, color) in enumerate(impacts):
        x = 1.5 + (i * 2.1)
        add_small_text(slide, x, impact_y, 1.8, 0.3, label, 8, True)
        add_shape_box(slide, x, impact_y + 0.4, 1.8, 0.7, value, color, 11)

    # Lesson learned box
    add_shape_box(slide, 1.0, 6.0, 8.5, 0.9,
                   "교훈: 재고 Zero = 리스크 Maximum | 안전재고 비용 << 생산 중단 손실",
                   COLOR_ACCENT, 9)

def add_slide_8_company_responses(prs):
    """Slide 8: 기업 대응 사례 (Apple, Intel, Toyota, GM)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "글로벌 기업의 JIC 전환 사례",
        "Apple, Intel, Toyota, GM 등은 안전재고를 2-12주에서 6-12개월로 확대하고 공급망을 다변화했습니다.",
        8
    )

    # 4 companies in grid (16 boxes total)
    companies = [
        {
            "name": "Apple",
            "before": "안전재고\n2주",
            "after": "안전재고\n6-8주",
            "investment": "$430억\n투자",
            "color": COLOR_DARK_GRAY
        },
        {
            "name": "Intel",
            "before": "안전재고\n4주",
            "after": "안전재고\n12주",
            "investment": "$200억\n공장",
            "color": COLOR_ACCENT
        },
        {
            "name": "Toyota",
            "before": "JIT 철학\n고수",
            "after": "안전재고\n4-6개월",
            "investment": "JIT 포기\n선언",
            "color": COLOR_MED_GRAY
        },
        {
            "name": "GM",
            "before": "Single\nSource",
            "after": "Triple\nSource",
            "investment": "$70억\n투자",
            "color": COLOR_DARK_GRAY
        }
    ]

    row_y = [2.2, 4.5]
    col_x = [0.8, 5.6]

    for i, company in enumerate(companies):
        row = i // 2
        col = i % 2
        x_base = col_x[col]
        y_base = row_y[row]

        # Company name
        add_shape_box(slide, x_base, y_base, 4.3, 0.6, company["name"], company["color"], 12)

        # Before/After/Investment (3 small boxes)
        add_shape_box(slide, x_base, y_base + 0.7, 1.4, 0.8, company["before"], COLOR_LIGHT_GRAY, 8)
        add_arrow(slide, x_base + 1.5, y_base + 1.1, x_base + 1.7, y_base + 1.1)
        add_shape_box(slide, x_base + 1.8, y_base + 0.7, 1.4, 0.8, company["after"], COLOR_ACCENT, 8)
        add_shape_box(slide, x_base + 3.3, y_base + 0.7, 1.0, 0.8, company["investment"], company["color"], 7)

def add_slide_9_kraljic_intro(prs):
    """Slide 9: Kraljic Matrix 소개"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "Kraljic Matrix: 자재 포트폴리오 전략의 핵심",
        "1983년 Peter Kraljic이 제시한 2×2 매트릭스는 자재를 차별화하여 관리하는 혁신적 프레임워크입니다.",
        9
    )

    # History timeline (1970s oil crisis → 1983 HBR paper)
    add_shape_box(slide, 1.0, 2.2, 2.5, 1.0,
                   "1970년대\n석유파동", COLOR_BOTTLENECK, 10)
    add_arrow(slide, 3.5, 2.7, 4.0, 2.7)
    add_shape_box(slide, 4.0, 2.2, 2.5, 1.0,
                   "공급 불안정\n극심", COLOR_DARK_GRAY, 10)
    add_arrow(slide, 6.5, 2.7, 7.0, 2.7)
    add_shape_box(slide, 7.0, 2.2, 2.5, 1.0,
                   "1983년\nKraljic 논문", COLOR_ACCENT, 10)

    # Key insight quote box
    add_shape_box(slide, 1.0, 3.8, 8.5, 1.2,
                   '"Not all materials are created equal"\n모든 자재가 동등하게 만들어지지 않았다.\n자재 특성에 따라 차별화된 전략이 필요하다.',
                   COLOR_WHITE, 10)

    # Problem of uniform management (8pt small text - HIGH DENSITY)
    problems = [
        "전략자재에 원가 절감 압박 → 공급업체 이탈 → 생산 중단",
        "레버리지자재에 높은 안전재고 → 재고 비용 증가 → 경쟁력 상실",
        "병목자재를 소홀히 관리 → 결품 발생 → 라인 정지",
        "일상자재에 과도한 관리 → 인력 낭비 → 전략 업무 소홀"
    ]

    problem_y = 5.5
    for i, prob in enumerate(problems):
        add_small_text(slide, 1.2, problem_y + (i * 0.4), 8.5, 0.35, f"❌ {prob}", 8, False)

def add_slide_10_two_axes(prs):
    """Slide 10: 2개 축 설명 (detailed with examples)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "Kraljic Matrix의 2개 축: 공급 리스크 × 구매 임팩트",
        "Y축은 공급의 어려움을, X축은 사업에 미치는 영향을 나타내며, 각 6개 평가 요소로 측정합니다.",
        10
    )

    # Y-axis: Supply Risk (left side with 6 factors)
    add_shape_box(slide, 0.5, 2.0, 4.5, 0.7, "Y축: 공급 리스크 (Supply Risk)", COLOR_BOTTLENECK, 11)

    risk_factors = [
        "공급업체 수: 1-2개 (고) ↔ 10개+ (저)",
        "대체 가능성: 대체 불가 (고) ↔ 쉽게 대체 (저)",
        "시장 구조: 독과점 (고) ↔ 경쟁시장 (저)",
        "지역 집중도: 특정 지역 (고) ↔ 분산 (저)",
        "기술 복잡성: 특수 기술 (고) ↔ 범용 (저)",
        "리드타임: 6개월+ (고) ↔ 1개월- (저)"
    ]

    for i, factor in enumerate(risk_factors):
        add_small_text(slide, 0.7, 2.9 + (i * 0.55), 4.0, 0.45, f"• {factor}", 8, False)

    # X-axis: Profit Impact (right side with 5 factors)
    add_shape_box(slide, 5.5, 2.0, 4.5, 0.7, "X축: 구매 임팩트 (Profit Impact)", COLOR_ACCENT, 11)

    impact_factors = [
        "구매 금액: 총 구매액 대비 비중",
        "원가 비중: 제품 원가 중 차지 비율",
        "사업 영향도: 결품 시 생산 중단 여부",
        "부가가치: 최종 제품 성능 영향",
        "품질 중요성: 품질 문제 시 파급 효과"
    ]

    for i, factor in enumerate(impact_factors):
        add_small_text(slide, 5.7, 2.9 + (i * 0.65), 4.0, 0.55, f"• {factor}", 8, False)

    # Scoring formula (bottom - 9pt)
    add_shape_box(slide, 1.0, 6.0, 8.5, 0.9,
                   "점수화: 공급 리스크 = (업체수×0.3) + (대체성×0.3) + (리드타임×0.2) + (지역×0.2)\n          구매 임팩트 = (구매금액×0.4) + (원가비중×0.3) + (사업영향×0.3)",
                   COLOR_LIGHT_GRAY, 8)

def add_slide_11_four_quadrants(prs):
    """Slide 11: 2×2 Matrix with 4 quadrants"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "Kraljic Matrix 2×2: 4대 자재군",
        "공급 리스크와 구매 임팩트의 조합으로 병목, 레버리지, 전략, 일상 4개 자재군으로 분류합니다.",
        11
    )

    # Draw 2x2 matrix (centered, large)
    matrix_x = 2.0
    matrix_y = 2.5
    box_w = 3.5
    box_h = 2.0

    # Quadrant 1: Bottleneck (top-left)
    add_shape_box(slide, matrix_x, matrix_y, box_w, box_h,
                   "병목자재\n(Bottleneck)\n\n높은 리스크\n낮은 임팩트",
                   COLOR_BOTTLENECK, 11)

    # Quadrant 2: Strategic (top-right)
    add_shape_box(slide, matrix_x + box_w, matrix_y, box_w, box_h,
                   "전략자재\n(Strategic)\n\n높은 리스크\n높은 임팩트",
                   COLOR_STRATEGIC, 11)

    # Quadrant 3: Routine (bottom-left)
    add_shape_box(slide, matrix_x, matrix_y + box_h, box_w, box_h,
                   "일상자재\n(Routine)\n\n낮은 리스크\n낮은 임팩트",
                   COLOR_ROUTINE, 11)

    # Quadrant 4: Leverage (bottom-right)
    add_shape_box(slide, matrix_x + box_w, matrix_y + box_h, box_w, box_h,
                   "레버리지자재\n(Leverage)\n\n낮은 리스크\n높은 임팩트",
                   COLOR_LEVERAGE, 11)

    # Axis labels (8pt)
    add_small_text(slide, 0.8, 3.5, 1.0, 0.5, "공급\n리스크\n↑", 9, True)
    add_small_text(slide, 4.0, 7.0, 3.0, 0.3, "구매 임팩트 →", 9, True)

def add_slide_12_bottleneck(prs):
    """Slide 12: 병목자재 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "병목자재 (Bottleneck Items): 공급 확보가 최우선",
        "금액은 작지만 없으면 생산 중단되는 자재로, 안전재고 4-8주와 Dual Sourcing이 핵심 전략입니다.",
        12
    )

    # Characteristics (4 boxes)
    chars = [
        "높은 공급 리스크\n+\n낮은 구매 임팩트",
        "공급업체\n1-2개로 제한",
        "결품 시\n생산 중단",
        "리드타임 길거나\n불안정"
    ]

    char_y = 2.0
    for i, char in enumerate(chars):
        add_shape_box(slide, 1.0 + (i * 2.2), char_y, 2.0, 1.0, char, COLOR_BOTTLENECK, 9)

    # Examples (8pt small text)
    add_small_text(slide, 1.0, 3.3, 8.5, 0.3, "사례: 차량용 MCU, 특수 센서, 희토류, FDA 인증 부품", 8, True)

    # Strategy table (HIGH DENSITY)
    strategy_items = [
        ("목표", "공급 안정성 확보", "비용보다 공급이 우선"),
        ("안전재고", "4-8주 (높음)", "95-99% 가용률 목표"),
        ("공급업체", "Dual Sourcing 시도", "1개 → 2-3개로 증가"),
        ("발주 방식", "ROP (재주문점)", "재고 수준 지속 모니터링"),
        ("KPI", "재고 가용률 95%+", "긴급 발주 5회/월 이하")
    ]

    table_y = 4.0
    for i, (category, strategy, detail) in enumerate(strategy_items):
        y = table_y + (i * 0.5)
        add_small_text(slide, 1.0, y, 2.0, 0.4, category, 9, True)
        add_small_text(slide, 3.2, y, 3.0, 0.4, strategy, 8, False)
        add_small_text(slide, 6.4, y, 3.2, 0.4, detail, 8, False)

def add_slide_13_leverage(prs):
    """Slide 13: 레버리지자재 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "레버리지자재 (Leverage Items): 원가 절감이 핵심",
        "금액이 크지만 공급 안정적인 자재로, 경쟁 입찰과 MRP 최적화를 통해 연 3-5% 절감이 목표입니다.",
        13
    )

    # Characteristics
    chars = [
        "낮은 공급 리스크\n+\n높은 구매 임팩트",
        "공급업체\n5개 이상 다수",
        "표준화되어\n대체 가능",
        "경쟁시장\n형성"
    ]

    char_y = 2.0
    for i, char in enumerate(chars):
        add_shape_box(slide, 1.0 + (i * 2.2), char_y, 2.0, 1.0, char, COLOR_LEVERAGE, 9)

    # Examples
    add_small_text(slide, 1.0, 3.3, 8.5, 0.3, "사례: 범용 원자재 (철강, 플라스틱), 표준 전자부품, 포장재, 케이블", 8, True)

    # Strategy table
    strategy_items = [
        ("목표", "원가 절감", "연 3-5% 절감 목표"),
        ("안전재고", "1-2주 (낮음)", "재고 최소화"),
        ("공급업체", "경쟁 입찰", "5개 이상 경쟁 유지"),
        ("발주 방식", "MRP + EOQ", "계획 기반 최적화"),
        ("KPI", "원가 절감률 3-5%", "재고 회전율 12회+")
    ]

    table_y = 4.0
    for i, (category, strategy, detail) in enumerate(strategy_items):
        y = table_y + (i * 0.5)
        add_small_text(slide, 1.0, y, 2.0, 0.4, category, 9, True)
        add_small_text(slide, 3.2, y, 3.0, 0.4, strategy, 8, False)
        add_small_text(slide, 6.4, y, 3.2, 0.4, detail, 8, False)

def add_slide_14_strategic(prs):
    """Slide 14: 전략자재 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "전략자재 (Strategic Items): 장기 파트너십 구축",
        "금액도 크고 공급도 어려운 핵심 자재로, 3-5년 장기 계약과 Win-Win 협력이 필수입니다.",
        14
    )

    # Characteristics
    chars = [
        "높은 공급 리스크\n+\n높은 구매 임팩트",
        "사업 성패를\n좌우",
        "대체 불가능\n전환 비용 높음",
        "장기 개발\n필요"
    ]

    char_y = 2.0
    for i, char in enumerate(chars):
        add_shape_box(slide, 1.0 + (i * 2.2), char_y, 2.0, 1.0, char, COLOR_STRATEGIC, 9)

    # Examples
    add_small_text(slide, 1.0, 3.3, 8.5, 0.3, "사례: 핵심 반도체 (AP, SoC), OLED 발광재료, 항공기 엔진, 독점 기술 부품", 8, True)

    # Strategy table
    strategy_items = [
        ("목표", "전략적 파트너십", "Win-Win 협력"),
        ("안전재고", "3-6주 (중상)", "공급 연속성 100%"),
        ("공급업체", "선택적 협력", "1-2개 전략 파트너"),
        ("발주 방식", "LTP + Hybrid", "장기 + 중기 + 단기 통합"),
        ("KPI", "공급 연속성 100%", "TCO 개선 -3%/년")
    ]

    table_y = 4.0
    for i, (category, strategy, detail) in enumerate(strategy_items):
        y = table_y + (i * 0.5)
        add_small_text(slide, 1.0, y, 2.0, 0.4, category, 9, True)
        add_small_text(slide, 3.2, y, 3.0, 0.4, strategy, 8, False)
        add_small_text(slide, 6.4, y, 3.2, 0.4, detail, 8, False)

def add_slide_15_routine(prs):
    """Slide 15: 일상자재 상세"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "일상자재 (Routine Items): 자동화와 효율화",
        "금액도 작고 공급도 쉬운 자재로, VMI와 자동 발주를 통해 관리 비용을 최소화합니다.",
        15
    )

    # Characteristics
    chars = [
        "낮은 공급 리스크\n+\n낮은 구매 임팩트",
        "품목 수는\n전체의 60-80%",
        "관리 부담\n큼",
        "표준화\n가능"
    ]

    char_y = 2.0
    for i, char in enumerate(chars):
        add_shape_box(slide, 1.0 + (i * 2.2), char_y, 2.0, 1.0, char, COLOR_ROUTINE, 9)

    # Examples
    add_small_text(slide, 1.0, 3.3, 8.5, 0.3, "사례: MRO 품목 (볼트, 너트, 공구), 사무용품, 청소용품, 소모품 전반", 8, True)

    # Strategy table
    strategy_items = [
        ("목표", "효율화 & 자동화", "최소 노력으로 관리"),
        ("안전재고", "1주 (최소)", "자동 보충"),
        ("공급업체", "통합", "10-20개 → 3-5개"),
        ("발주 방식", "Min-Max / VMI", "완전 자동화"),
        ("KPI", "발주 비용 5천원/건↓", "자동화율 80%+")
    ]

    table_y = 4.0
    for i, (category, strategy, detail) in enumerate(strategy_items):
        y = table_y + (i * 0.5)
        add_small_text(slide, 1.0, y, 2.0, 0.4, category, 9, True)
        add_small_text(slide, 3.2, y, 3.0, 0.4, strategy, 8, False)
        add_small_text(slide, 6.4, y, 3.2, 0.4, detail, 8, False)

def add_slide_16_comparison_matrix(prs):
    """Slide 16: 4대 자재군 비교표 (comprehensive)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "4대 자재군 핵심 전략 비교",
        "자재군별로 목표, 안전재고, 발주 방식, KPI가 명확히 차별화되어야 최적 성과를 달성합니다.",
        16
    )

    # Comparison table (VERY HIGH DENSITY with 8pt fonts)
    headers = ["", "병목", "레버리지", "전략", "일상"]
    rows = [
        ("관리 목표", "공급 확보", "원가 절감", "파트너십", "효율화"),
        ("우선순위", "공급>원가", "원가>품질", "관계>이익", "효율>정확"),
        ("안전재고", "4-8주↑", "1-2주↓", "3-6주→", "1주↓"),
        ("발주 방식", "ROP", "MRP", "LTP+Hybrid", "Min-Max/VMI"),
        ("서비스 수준", "95-99%", "85-90%", "90-95%", "80-85%"),
        ("KPI", "가용률95%+", "절감3-5%", "연속성100%", "자동화80%+")
    ]

    # Header row
    header_y = 2.0
    col_widths = [1.8, 1.9, 1.9, 1.9, 1.9]
    col_x = [0.8, 2.7, 4.7, 6.7, 8.7]

    for i, header in enumerate(headers):
        if i == 0:
            add_small_text(slide, col_x[i], header_y, col_widths[i], 0.4, header, 9, True)
        else:
            colors = [COLOR_BOTTLENECK, COLOR_LEVERAGE, COLOR_STRATEGIC, COLOR_ROUTINE]
            add_shape_box(slide, col_x[i], header_y, col_widths[i], 0.4, header, colors[i-1], 9)

    # Data rows
    for row_idx, (category, *values) in enumerate(rows):
        y = header_y + 0.5 + (row_idx * 0.55)
        add_small_text(slide, col_x[0], y, col_widths[0], 0.5, category, 8, True)
        for col_idx, value in enumerate(values):
            add_small_text(slide, col_x[col_idx + 1], y, col_widths[col_idx + 1], 0.5, value, 7, False)

def add_slide_17_planning_overview(prs):
    """Slide 17: 자재계획 방법론 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "자재계획 방법론 전체 맵",
        "4대 자재군은 각각 ROP, MRP, LTP+Hybrid, Min-Max/VMI 방법론을 적용하여 최적화합니다.",
        17
    )

    # 5 methodologies with flowchart
    methods = [
        ("ROP\n재주문점", "병목자재", "재고 도달 시\n자동 발주", COLOR_BOTTLENECK),
        ("MRP\n소요량 계획", "레버리지", "MPS+BOM\n기반 계산", COLOR_LEVERAGE),
        ("LTP\n장기 계획", "전략자재", "18-24개월\n프레임워크", COLOR_STRATEGIC),
        ("Hybrid\n통합", "전략자재", "장기+중기\n+단기", COLOR_STRATEGIC),
        ("Min-Max\nVMI", "일상자재", "자동 보충\n완전 자동화", COLOR_ROUTINE)
    ]

    method_y = 2.2
    for i, (name, material, desc, color) in enumerate(methods):
        y = method_y + (i * 0.9)
        # Method name
        add_shape_box(slide, 1.0, y, 1.8, 0.7, name, color, 10)
        # Material type
        add_small_text(slide, 3.0, y + 0.1, 1.8, 0.5, material, 9, True)
        # Description
        add_small_text(slide, 5.0, y + 0.1, 4.5, 0.5, desc, 8, False)

def add_slide_18_decision_tree(prs):
    """Slide 18: 방법론 선택 의사결정 트리"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "자재계획 방법론 선택 Decision Tree",
        "공급 리스크와 구매 임팩트를 순차적으로 평가하여 적합한 계획 방법론을 선택합니다.",
        18
    )

    # Decision tree with shapes (15+ boxes and arrows)
    # Start
    add_shape_box(slide, 4.0, 2.0, 2.5, 0.6, "START:\n자재 분석", COLOR_DARK_GRAY, 10)
    add_arrow(slide, 5.25, 2.6, 5.25, 3.0)

    # First question
    add_shape_box(slide, 3.5, 3.0, 3.5, 0.7, "공급 리스크 높은가?", COLOR_ACCENT, 10)

    # YES branch (left)
    add_arrow(slide, 3.5, 3.35, 2.5, 3.35)
    add_small_text(slide, 2.6, 3.2, 0.6, 0.3, "YES", 8, True)
    add_shape_box(slide, 1.0, 3.8, 1.4, 0.6, "구매\n임팩트?", COLOR_MED_GRAY, 9)

    # YES-YES: Strategic
    add_arrow(slide, 1.7, 4.4, 1.7, 4.8)
    add_small_text(slide, 1.8, 4.5, 0.5, 0.2, "높음", 7, False)
    add_shape_box(slide, 0.8, 4.8, 1.8, 0.8, "전략자재\n↓\nLTP+Hybrid", COLOR_STRATEGIC, 9)

    # YES-NO: Bottleneck
    add_arrow(slide, 2.4, 4.1, 3.5, 4.1)
    add_small_text(slide, 2.5, 4.0, 0.5, 0.2, "낮음", 7, False)
    add_shape_box(slide, 3.5, 4.8, 1.8, 0.8, "병목자재\n↓\nROP", COLOR_BOTTLENECK, 9)

    # NO branch (right)
    add_arrow(slide, 7.0, 3.35, 8.0, 3.35)
    add_small_text(slide, 7.2, 3.2, 0.5, 0.3, "NO", 8, True)
    add_shape_box(slide, 8.0, 3.8, 1.4, 0.6, "구매\n임팩트?", COLOR_MED_GRAY, 9)

    # NO-YES: Leverage
    add_arrow(slide, 8.7, 4.4, 8.7, 4.8)
    add_small_text(slide, 8.8, 4.5, 0.5, 0.2, "높음", 7, False)
    add_shape_box(slide, 7.8, 4.8, 1.8, 0.8, "레버리지\n↓\nMRP", COLOR_LEVERAGE, 9)

    # NO-NO: Routine
    add_arrow(slide, 9.4, 4.1, 10.2, 4.1)
    add_small_text(slide, 9.5, 4.0, 0.5, 0.2, "낮음", 7, False)
    add_shape_box(slide, 10.0, 4.8, 1.0, 0.8, "일상\n↓\nVMI", COLOR_ROUTINE, 8)

def add_slide_19_kpi_framework(prs):
    """Slide 19: KPI 프레임워크"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "자재군별 핵심 KPI 프레임워크",
        "병목은 가용률 95%, 레버리지는 절감률 3-5%, 전략은 연속성 100%, 일상은 자동화율 80%가 목표입니다.",
        19
    )

    # KPI table (HIGH DENSITY - 8pt fonts)
    kpi_data = [
        ("병목", "재고 가용률\n95%+", "공급 안정성\n95%+", "긴급 발주\n5회/월↓", COLOR_BOTTLENECK),
        ("레버리지", "원가 절감률\n3-5%/년", "재고 회전율\n12회+", "발주 정확도\n95%+", COLOR_LEVERAGE),
        ("전략", "공급 연속성\n100%", "LTP 정확도\n70%+", "TCO 개선\n-3%/년", COLOR_STRATEGIC),
        ("일상", "발주 비용\n5천원/건↓", "자동화율\n80%+", "업체 수\n100개↓", COLOR_ROUTINE)
    ]

    table_y = 2.2
    for i, (material, kpi1, kpi2, kpi3, color) in enumerate(kpi_data):
        y = table_y + (i * 1.2)
        # Material name
        add_shape_box(slide, 0.8, y, 1.5, 0.9, material, color, 11)
        # 3 KPIs
        add_shape_box(slide, 2.5, y, 2.4, 0.9, kpi1, COLOR_LIGHT_GRAY, 8)
        add_shape_box(slide, 5.1, y, 2.4, 0.9, kpi2, COLOR_LIGHT_GRAY, 8)
        add_shape_box(slide, 7.7, y, 2.4, 0.9, kpi3, COLOR_LIGHT_GRAY, 8)

def add_slide_20_summary(prs):
    """Slide 20: Summary"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_and_governing(
        slide,
        "Part 1 핵심 요약",
        "JIT→JIC 전환, Kraljic 4대 자재군, 차별화 전략, 계획 방법론이 전략적 재고운영의 핵심입니다.",
        20
    )

    # 4 key takeaways in boxes
    takeaways = [
        {
            "title": "1. 패러다임 전환",
            "content": "JIT (재고=낭비)\n→\nJIC (재고=전략자산)\n효율성 → 회복력",
            "color": COLOR_DARK_GRAY
        },
        {
            "title": "2. Kraljic Matrix",
            "content": "2개 축: 공급 리스크\n× 구매 임팩트\n4개 자재군 분류",
            "color": COLOR_ACCENT
        },
        {
            "title": "3. 차별화 전략",
            "content": "병목: 공급 확보\n레버리지: 원가 절감\n전략: 파트너십\n일상: 자동화",
            "color": COLOR_MED_GRAY
        },
        {
            "title": "4. 계획 방법론",
            "content": "ROP (병목)\nMRP (레버리지)\nLTP+Hybrid (전략)\nMin-Max/VMI (일상)",
            "color": COLOR_DARK_GRAY
        }
    ]

    for i, item in enumerate(takeaways):
        row = i // 2
        col = i % 2
        x = 1.0 + (col * 4.5)
        y = 2.2 + (row * 2.5)

        # Title
        add_shape_box(slide, x, y, 4.0, 0.6, item["title"], item["color"], 11)
        # Content
        add_small_text(slide, x + 0.2, y + 0.8, 3.6, 1.4, item["content"], 9, False)

    # Next session preview
    add_shape_box(slide, 2.0, 6.5, 6.5, 0.8,
                   "Next → Session 2: 소싱 전략 및 공급업체 관계 관리 (SRM)",
                   COLOR_ACCENT, 10)

if __name__ == "__main__":
    create_presentation()
