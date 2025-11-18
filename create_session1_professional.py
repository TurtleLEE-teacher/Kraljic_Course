#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
1회차 강의자료 생성기 (전문 컨설팅 스타일)
전략적 재고운영 Foundation: Kraljic Matrix와 자재계획 방법론
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    """메인 프레젠테이션 생성 - 16:9 비율"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 비율
    prs.slide_height = Inches(7.5)

    # 슬라이드 생성
    create_title_slide(prs)
    create_agenda_slide(prs)
    create_section_divider(prs, "1. 패러다임의 전환", "JIT에서 JIC로")
    create_jit_crisis_slide(prs)
    create_jit_vs_jic_comparison(prs)
    create_section_divider(prs, "2. Kraljic Matrix", "자재 포트폴리오 전략의 핵심")
    create_kraljic_matrix_slide(prs)
    create_four_quadrants_detail(prs)
    create_key_takeaways_slide(prs)

    return prs

def add_header_footer(slide, title_text):
    """헤더 및 타이틀 추가 (상단 바 스타일)"""
    # 상단 색상 바
    header_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(0.7)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 32, 96)  # 진한 네이비
    header_bar.line.fill.background()

    # 타이틀 텍스트
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15),
        Inches(12), Inches(0.4)
    )
    text_frame = title_box.text_frame
    text_frame.text = title_text

    p = text_frame.paragraphs[0]
    p.font.name = "맑은 고딕"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.LEFT

    # 하단 페이지 바
    footer_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(7.3),
        Inches(13.333), Inches(0.2)
    )
    footer_bar.fill.solid()
    footer_bar.fill.fore_color.rgb = RGBColor(0, 112, 192)  # 밝은 블루
    footer_bar.line.fill.background()

def create_title_slide(prs):
    """슬라이드 1: 타이틀 페이지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경 그라디언트 효과를 위한 큰 사각형
    bg_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(3.5)
    )
    bg_top.fill.solid()
    bg_top.fill.fore_color.rgb = RGBColor(0, 32, 96)
    bg_top.line.fill.background()

    bg_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(3.5),
        Inches(13.333), Inches(4)
    )
    bg_bottom.fill.solid()
    bg_bottom.fill.fore_color.rgb = RGBColor(245, 248, 252)
    bg_bottom.line.fill.background()

    # 메인 타이틀
    title = slide.shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(11), Inches(1.2)
    )
    tf = title.text_frame
    p = tf.paragraphs[0]
    p.text = "전략적 재고운영 및 자재계획수립"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle = slide.shapes.add_textbox(
        Inches(1), Inches(4.2),
        Inches(11), Inches(0.8)
    )
    tf = subtitle.text_frame
    p = tf.paragraphs[0]
    p.text = "[1회차] Kraljic Matrix와 자재계획 방법론"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 32, 96)
    p.alignment = PP_ALIGN.CENTER

    # 정보 박스
    info_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(4), Inches(5.5),
        Inches(5.333), Inches(1.2)
    )
    info_box.fill.solid()
    info_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
    info_box.line.color.rgb = RGBColor(0, 112, 192)
    info_box.line.width = Pt(2)

    tf = info_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    p = tf.paragraphs[0]
    p.text = "난이도: 중급  |  소요시간: 45분"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0, 32, 96)
    p.alignment = PP_ALIGN.CENTER

def create_agenda_slide(prs):
    """슬라이드 2: 목차 (Agenda)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_header_footer(slide, "학습 목표 및 과정 구성")

    # 학습 목표 박스
    objectives_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.2),
        Inches(5.5), Inches(2.5)
    )
    objectives_box.fill.solid()
    objectives_box.fill.fore_color.rgb = RGBColor(230, 240, 255)
    objectives_box.line.color.rgb = RGBColor(0, 112, 192)
    objectives_box.line.width = Pt(3)

    tf = objectives_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "🎯 학습 목표"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 32, 96)

    objectives = [
        "JIT에서 JIC로의 패러다임 전환 이해",
        "전략적 재고운영의 핵심 개념 습득",
        "Kraljic Matrix 자재 분류 역량 확보",
        "자재군별 관리 철학 이해"
    ]

    for obj in objectives:
        p = tf.add_paragraph()
        p.text = f"• {obj}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(14)
        p.space_before = Pt(6)
        p.level = 0

    # 과정 구성 박스
    agenda_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7), Inches(1.2),
        Inches(5.5), Inches(5.5)
    )
    agenda_box.fill.solid()
    agenda_box.fill.fore_color.rgb = RGBColor(255, 250, 240)
    agenda_box.line.color.rgb = RGBColor(230, 126, 34)
    agenda_box.line.width = Pt(3)

    tf = agenda_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "📋 과정 구성"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 32, 96)

    agenda_items = [
        ("1", "패러다임의 전환: JIT → JIC"),
        ("2", "Kraljic Matrix 프레임워크"),
        ("3", "4대 자재군 특성 및 관리 철학"),
        ("4", "자재계획 방법론 맵"),
        ("5", "통합 KPI 프레임워크")
    ]

    for num, item in agenda_items:
        p = tf.add_paragraph()
        p.text = f"{num}. {item}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(16)
        p.font.bold = True
        p.space_before = Pt(10)
        p.font.color.rgb = RGBColor(230, 126, 34)

def create_section_divider(prs, section_num, section_title):
    """섹션 구분 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 전체 배경
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0, 32, 96)
    bg.line.fill.background()

    # 액센트 바
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(3),
        Inches(0.3), Inches(1.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(0, 176, 240)
    accent.line.fill.background()

    # 섹션 번호
    section_box = slide.shapes.add_textbox(
        Inches(2), Inches(2.5),
        Inches(9), Inches(1)
    )
    tf = section_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_num
    p.font.name = "맑은 고딕"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 섹션 타이틀
    title_box = slide.shapes.add_textbox(
        Inches(2), Inches(3.8),
        Inches(9), Inches(1)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.name = "맑은 고딕"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(0, 176, 240)

def create_jit_crisis_slide(prs):
    """JIT의 위기 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_header_footer(slide, "JIT의 위기: 2021년 글로벌 반도체 대란")

    # 좌측: 문제 상황
    problem_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8), Inches(1.2),
        Inches(5.8), Inches(2.5)
    )
    problem_box.fill.solid()
    problem_box.fill.fore_color.rgb = RGBColor(255, 235, 235)
    problem_box.line.color.rgb = RGBColor(231, 76, 60)
    problem_box.line.width = Pt(3)

    tf = problem_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "⚠️ JIT의 붕괴"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(192, 0, 0)

    crisis_points = [
        "자동차 생산 1,000만 대 감소",
        "산업 전체 손실 $210억",
        "GM, 포드, 폭스바겐 생산 중단",
        "안전재고 Zero → 공급 충격 즉시 노출"
    ]

    for point in crisis_points:
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(14)
        p.space_before = Pt(8)

    # 우측: 교훈
    lesson_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(6.8), Inches(1.2),
        Inches(5.8), Inches(2.5)
    )
    lesson_box.fill.solid()
    lesson_box.fill.fore_color.rgb = RGBColor(235, 255, 245)
    lesson_box.line.color.rgb = RGBColor(46, 204, 113)
    lesson_box.line.width = Pt(3)

    tf = lesson_box.text_frame
    tf.margin_left = Inches(0.3)
    tf.margin_top = Inches(0.2)

    p = tf.paragraphs[0]
    p.text = "💡 핵심 교훈"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)

    lessons = [
        "JIT = Efficient but Fragile",
        "재고 Zero = 리스크 Maximum",
        "차별화된 재고 전략 필요",
        "병목자재 안전재고 확대 필수"
    ]

    for lesson in lessons:
        p = tf.add_paragraph()
        p.text = f"✓ {lesson}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(14)
        p.space_before = Pt(8)

    # 하단: 통계 박스들
    stats = [
        ("1-2주", "Before\nJIT 안전재고"),
        ("6개월", "After\nJIC 안전재고"),
        ("15배", "재고 증가율")
    ]

    for i, (value, label) in enumerate(stats):
        stat_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 4), Inches(4.2),
            Inches(3.5), Inches(2)
        )
        stat_box.fill.solid()
        stat_box.fill.fore_color.rgb = RGBColor(0, 32, 96)
        stat_box.line.fill.background()

        tf = stat_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = tf.paragraphs[0]
        p.text = value
        p.font.name = "맑은 고딕"
        p.font.size = Pt(42)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 176, 240)
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = label
        p.font.name = "맑은 고딕"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

def create_jit_vs_jic_comparison(prs):
    """JIT vs JIC 비교 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_header_footer(slide, "패러다임의 전환: JIT vs JIC")

    # 테이블 헤더
    headers = ["구분", "JIT (과거)", "JIC (현재/미래)"]
    header_colors = [RGBColor(0, 32, 96), RGBColor(231, 76, 60), RGBColor(46, 204, 113)]

    for i, (header, color) in enumerate(zip(headers, header_colors)):
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8 + i * 4.2), Inches(1.2),
            Inches(4), Inches(0.5)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = header
        p.font.name = "맑은 고딕"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER

    # 비교 항목
    comparisons = [
        ("재고 철학", "재고 = 낭비", "재고 = 전략적 자산"),
        ("목표", "재고 최소화 (Zero)", "최적 재고 (Optimal)"),
        ("우선순위", "효율성 (Efficiency)", "회복력 (Resilience)"),
        ("리스크 관점", "리스크 무시", "리스크 관리"),
        ("공급망 구조", "글로벌 최적화", "지역 분산"),
        ("안전재고", "최소 (1-2주)", "차별화 (1주-6개월)")
    ]

    for i, (category, jit, jic) in enumerate(comparisons):
        y_pos = 1.8 + i * 0.75

        # 카테고리
        cat_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(y_pos),
            Inches(4), Inches(0.6)
        )
        cat_box.fill.solid()
        cat_box.fill.fore_color.rgb = RGBColor(240, 240, 240)
        cat_box.line.color.rgb = RGBColor(200, 200, 200)

        tf = cat_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = category
        p.font.name = "맑은 고딕"
        p.font.size = Pt(13)
        p.font.bold = True

        # JIT
        jit_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(5), Inches(y_pos),
            Inches(4), Inches(0.6)
        )
        jit_box.fill.solid()
        jit_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
        jit_box.line.color.rgb = RGBColor(200, 200, 200)

        tf = jit_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = jit
        p.font.name = "맑은 고딕"
        p.font.size = Pt(12)

        # JIC
        jic_box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(9.2), Inches(y_pos),
            Inches(4), Inches(0.6)
        )
        jic_box.fill.solid()
        jic_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
        jic_box.line.color.rgb = RGBColor(200, 200, 200)

        tf = jic_box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        p = tf.paragraphs[0]
        p.text = jic
        p.font.name = "맑은 고딕"
        p.font.size = Pt(12)

def create_kraljic_matrix_slide(prs):
    """Kraljic Matrix 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_header_footer(slide, "Kraljic Matrix: 2×2 자재 포트폴리오")

    # 매트릭스 중심
    center_x = 7
    center_y = 4.2
    box_size = 2.8

    # 4개 사분면
    quadrants = [
        # (x, y, name, desc, bg_color, text_color, icon)
        (center_x - box_size, center_y - box_size,
         "병목자재\nBottleneck", "높은 공급 리스크\n낮은 구매 금액",
         RGBColor(255, 200, 200), RGBColor(192, 0, 0), "🔴"),

        (center_x, center_y - box_size,
         "전략자재\nStrategic", "높은 공급 리스크\n높은 구매 금액",
         RGBColor(230, 200, 255), RGBColor(128, 0, 128), "🟣"),

        (center_x - box_size, center_y,
         "일상자재\nRoutine", "낮은 공급 리스크\n낮은 구매 금액",
         RGBColor(240, 240, 240), RGBColor(96, 96, 96), "⚪"),

        (center_x, center_y,
         "레버리지자재\nLeverage", "낮은 공급 리스크\n높은 구매 금액",
         RGBColor(200, 255, 200), RGBColor(0, 128, 0), "🟢")
    ]

    for x, y, name, desc, bg_color, text_color, icon in quadrants:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(box_size - 0.1), Inches(box_size - 0.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = RGBColor(100, 100, 100)
        box.line.width = Pt(2)

        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = f"{icon}"
        p.font.size = Pt(32)
        p.alignment = PP_ALIGN.CENTER

        p = tf.add_paragraph()
        p.text = name
        p.font.name = "맑은 고딕"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

        p = tf.add_paragraph()
        p.text = desc
        p.font.name = "맑은 고딕"
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(6)

    # Y축 화살표 및 레이블
    y_arrow = slide.shapes.add_shape(
        MSO_SHAPE.UP_ARROW,
        Inches(2), Inches(2.5),
        Inches(0.4), Inches(3)
    )
    y_arrow.fill.solid()
    y_arrow.fill.fore_color.rgb = RGBColor(0, 112, 192)
    y_arrow.line.fill.background()

    y_label = slide.shapes.add_textbox(
        Inches(1.2), Inches(3.5),
        Inches(1.5), Inches(1)
    )
    tf = y_label.text_frame
    p = tf.paragraphs[0]
    p.text = "구매 금액\n(Purchase\nImpact)"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER

    # X축 화살표 및 레이블
    x_arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW,
        Inches(4.5), Inches(6.5),
        Inches(3), Inches(0.4)
    )
    x_arrow.fill.solid()
    x_arrow.fill.fore_color.rgb = RGBColor(0, 112, 192)
    x_arrow.line.fill.background()

    x_label = slide.shapes.add_textbox(
        Inches(5.5), Inches(6.8),
        Inches(2), Inches(0.5)
    )
    tf = x_label.text_frame
    p = tf.paragraphs[0]
    p.text = "공급 리스크 (Supply Risk)"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER

def create_four_quadrants_detail(prs):
    """4대 자재군 상세 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_header_footer(slide, "4대 자재군: 차별화된 관리 전략")

    materials = [
        ("🔴 병목자재", "공급 확보", "ROP (Re-Order Point)", "높은 안전재고 (4-8주)",
         RGBColor(255, 200, 200), RGBColor(192, 0, 0)),
        ("🟢 레버리지자재", "원가 절감", "MRP (계획 기반)", "경쟁 입찰, 낮은 재고",
         RGBColor(200, 255, 200), RGBColor(0, 128, 0)),
        ("🟣 전략자재", "파트너십", "하이브리드 계획", "장기 계약, 협력",
         RGBColor(230, 200, 255), RGBColor(128, 0, 128)),
        ("⚪ 일상자재", "효율화", "자동화 (VMI)", "프로세스 간소화",
         RGBColor(240, 240, 240), RGBColor(96, 96, 96))
    ]

    for i, (name, goal, method, strategy, bg_color, text_color) in enumerate(materials):
        row = i // 2
        col = i % 2

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + col * 6.2), Inches(1.5 + row * 2.8),
            Inches(5.8), Inches(2.4)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
        box.line.color.rgb = text_color
        box.line.width = Pt(3)

        tf = box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "맑은 고딕"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = text_color

        p = tf.add_paragraph()
        p.text = f"\n목표: {goal}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(15)
        p.font.bold = True
        p.space_before = Pt(8)

        p = tf.add_paragraph()
        p.text = f"계획 방법: {method}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(13)
        p.space_before = Pt(6)

        p = tf.add_paragraph()
        p.text = f"전략: {strategy}"
        p.font.name = "맑은 고딕"
        p.font.size = Pt(13)
        p.space_before = Pt(4)

def create_key_takeaways_slide(prs):
    """핵심 요약 슬라이드"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    add_header_footer(slide, "핵심 요약 (Key Takeaways)")

    # 타이틀
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.2),
        Inches(11.333), Inches(0.6)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💎 오늘 배운 핵심 내용"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 32, 96)
    p.alignment = PP_ALIGN.CENTER

    # 3개 핵심 포인트
    takeaways = [
        ("1. 패러다임 전환",
         ["JIT → JIC", "효율성 → 회복력", "획일적 관리 → 차별화 전략"],
         RGBColor(230, 240, 255)),

        ("2. Kraljic Matrix",
         ["2개 축: 공급 리스크 × 구매 금액", "4개 자재군 분류", "차별화된 관리 전략"],
         RGBColor(255, 250, 240)),

        ("3. 실행 방향",
         ["자재 특성 분석", "포트폴리오 분류", "자재군별 맞춤 전략 수립"],
         RGBColor(240, 255, 240))
    ]

    for i, (title, points, color) in enumerate(takeaways):
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.8 + i * 4.2), Inches(2.5),
            Inches(3.8), Inches(3.8)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = RGBColor(0, 112, 192)
        box.line.width = Pt(3)

        tf = box.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)

        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "맑은 고딕"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 32, 96)

        for point in points:
            p = tf.add_paragraph()
            p.text = f"• {point}"
            p.font.name = "맑은 고딕"
            p.font.size = Pt(13)
            p.space_before = Pt(10)
            p.level = 0

    # 하단 메시지
    footer_msg = slide.shapes.add_textbox(
        Inches(1), Inches(6.5),
        Inches(11.333), Inches(0.6)
    )
    tf = footer_msg.text_frame
    p = tf.paragraphs[0]
    p.text = "다음 회차: [2회차] 자재군별 소싱 전략 및 공급업체 관계 관리"
    p.font.name = "맑은 고딕"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 112, 192)
    p.alignment = PP_ALIGN.CENTER

def main():
    """메인 실행 함수"""
    print("\n" + "="*60)
    print("1회차 전문 강의자료 생성 중...")
    print("="*60 + "\n")

    prs = create_presentation()

    output_file = "/home/user/Kraljic_Course/PPTX_SAMPLE/Session1_KraljicMatrix_Foundation.pptx"
    prs.save(output_file)

    print(f"✅ 완료!")
    print(f"📁 파일 저장: {output_file}")
    print(f"📊 총 {len(prs.slides)} 슬라이드 생성")
    print(f"📐 크기: 16:9 (13.333\" x 7.5\")")
    print("\n" + "="*60 + "\n")

if __name__ == "__main__":
    main()
