#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
1회차 강의자료 샘플 생성기
전략적 재고운영 Foundation: Kraljic Matrix와 자재계획 방법론
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """메인 프레젠테이션 생성"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 슬라이드 생성
    create_title_slide(prs)
    create_learning_objectives_slide(prs)
    create_course_structure_slide(prs)
    create_jit_vs_jic_slide(prs)
    create_kraljic_matrix_slide(prs)
    create_material_categories_slide(prs)

    return prs

def add_title_shape(slide, title_text, top=0.5):
    """타이틀 박스 추가"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(top), Inches(9), Inches(0.8)
    )
    text_frame = title_box.text_frame
    text_frame.text = title_text

    # 스타일링
    p = text_frame.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)  # 다크 블루
    p.alignment = PP_ALIGN.CENTER

    return title_box

def add_content_box(slide, left, top, width, height, text, bg_color=None, font_size=14):
    """컨텐츠 박스 추가"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(left), Inches(top), Inches(width), Inches(height)
    )

    # 배경색 설정
    if bg_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
    else:
        shape.fill.background()

    # 테두리
    shape.line.color.rgb = RGBColor(200, 200, 200)
    shape.line.width = Pt(1)

    # 텍스트
    text_frame = shape.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 폰트 스타일
    p = text_frame.paragraphs[0]
    p.font.size = Pt(font_size)
    p.alignment = PP_ALIGN.CENTER

    return shape

def create_title_slide(prs):
    """슬라이드 1: 타이틀"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # 배경색
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 248, 252)  # 연한 블루 그레이

    # 메인 타이틀
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(1.5)
    )
    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "[1회차] 전략적 재고운영 Foundation"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.alignment = PP_ALIGN.CENTER

    # 서브타이틀
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(3.5), Inches(8), Inches(1)
    )
    text_frame = subtitle_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "Kraljic Matrix와 자재계획 방법론"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(51, 102, 153)
    p.alignment = PP_ALIGN.CENTER

    # 하단 정보
    info_box = slide.shapes.add_textbox(
        Inches(1), Inches(5.5), Inches(8), Inches(0.5)
    )
    text_frame = info_box.text_frame
    p = text_frame.paragraphs[0]
    p.text = "난이도: 중급 | 소요시간: 45분"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.CENTER

def create_learning_objectives_slide(prs):
    """슬라이드 2: 학습 목표"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 타이틀
    add_title_shape(slide, "🎯 학습 목표")

    # 4가지 목표 박스
    objectives = [
        "JIT에서 JIC로의\n패러다임 전환 이해",
        "전략적 재고운영의\n핵심 개념 습득",
        "Kraljic Matrix를 활용한\n자재 포트폴리오 분류",
        "자재군별 관리 철학과\n계획 방법론 이해"
    ]

    colors = [
        RGBColor(230, 240, 255),  # 연한 블루
        RGBColor(240, 255, 240),  # 연한 그린
        RGBColor(255, 245, 230),  # 연한 오렌지
        RGBColor(250, 240, 255),  # 연한 퍼플
    ]

    for i, (obj, color) in enumerate(zip(objectives, colors)):
        row = i // 2
        col = i % 2

        add_content_box(
            slide,
            left=0.5 + col * 4.7,
            top=2 + row * 2.2,
            width=4.2,
            height=1.8,
            text=obj,
            bg_color=color,
            font_size=16
        )

def create_course_structure_slide(prs):
    """슬라이드 3: 과정 구성 (MECE)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 타이틀
    add_title_shape(slide, "📋 1회차 구성 (MECE)")

    # 5개 섹션
    sections = [
        ("1", "패러다임의 전환", "JIT → JIC"),
        ("2", "Kraljic Matrix", "프레임워크"),
        ("3", "4대 자재군", "특성 및 관리 철학"),
        ("4", "자재계획 방법론", "전체 맵"),
        ("5", "통합 KPI", "프레임워크"),
    ]

    colors = [
        RGBColor(52, 152, 219),   # Blue
        RGBColor(46, 204, 113),   # Green
        RGBColor(241, 196, 15),   # Yellow
        RGBColor(230, 126, 34),   # Orange
        RGBColor(155, 89, 182),   # Purple
    ]

    for i, ((num, title, subtitle), color) in enumerate(zip(sections, colors)):
        # 메인 박스
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1.5), Inches(2 + i * 1), Inches(7), Inches(0.8)
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = color

        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 번호
        p = text_frame.paragraphs[0]
        p.text = f"{num}. {title}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        # 서브타이틀
        p = text_frame.add_paragraph()
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.level = 1

def create_jit_vs_jic_slide(prs):
    """슬라이드 4: JIT vs JIC 비교"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 타이틀
    add_title_shape(slide, "패러다임의 전환: JIT vs JIC")

    # 헤더
    header_left = add_content_box(slide, 0.5, 1.5, 4.2, 0.6, "JIT (과거)",
                                   RGBColor(231, 76, 60), 18)
    header_left.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header_left.text_frame.paragraphs[0].font.bold = True

    header_right = add_content_box(slide, 5.3, 1.5, 4.2, 0.6, "JIC (현재/미래)",
                                    RGBColor(46, 204, 113), 18)
    header_right.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    header_right.text_frame.paragraphs[0].font.bold = True

    # 비교 항목
    comparisons = [
        ("재고 = 낭비", "재고 = 전략적 자산"),
        ("재고 최소화 (Zero)", "최적 재고 (Optimal)"),
        ("효율성 우선", "회복력 우선"),
        ("리스크 무시", "리스크 관리"),
        ("글로벌 최적화", "지역 분산"),
        ("안전재고 1-2주", "안전재고 차별화"),
    ]

    for i, (jit, jic) in enumerate(comparisons):
        y_pos = 2.3 + i * 0.75

        # JIT
        box_left = add_content_box(slide, 0.5, y_pos, 4.2, 0.6, jit,
                                    RGBColor(255, 235, 235), 13)
        box_left.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        box_left.text_frame.paragraphs[0].level = 1

        # JIC
        box_right = add_content_box(slide, 5.3, y_pos, 4.2, 0.6, jic,
                                     RGBColor(235, 255, 245), 13)
        box_right.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        box_right.text_frame.paragraphs[0].level = 1

def create_kraljic_matrix_slide(prs):
    """슬라이드 5: Kraljic Matrix"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 타이틀
    add_title_shape(slide, "Kraljic Matrix: 2×2 자재 포트폴리오")

    # 매트릭스 중심 위치
    center_x = 5
    center_y = 4
    box_size = 2.5

    # 4개 사분면
    quadrants = [
        # (x_offset, y_offset, title, subtitle, color)
        (-box_size/2, -box_size/2, "🔴 병목자재", "높은 리스크\n낮은 금액", RGBColor(255, 200, 200)),
        (box_size/2, -box_size/2, "🟣 전략자재", "높은 리스크\n높은 금액", RGBColor(230, 200, 255)),
        (-box_size/2, box_size/2, "⚪ 일상자재", "낮은 리스크\n낮은 금액", RGBColor(240, 240, 240)),
        (box_size/2, box_size/2, "🟢 레버리지자재", "낮은 리스크\n높은 금액", RGBColor(200, 255, 200)),
    ]

    for x_off, y_off, title, subtitle, color in quadrants:
        shape = slide.shapes.add_shape(
            1,
            Inches(center_x + x_off - 0.25),
            Inches(center_y + y_off - 0.25),
            Inches(box_size - 0.1),
            Inches(box_size - 0.1)
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(100, 100, 100)
        shape.line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        p2 = text_frame.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.alignment = PP_ALIGN.CENTER

    # Y축 레이블
    y_label = slide.shapes.add_textbox(
        Inches(0.5), Inches(3), Inches(1.5), Inches(2)
    )
    p = y_label.text_frame.paragraphs[0]
    p.text = "구매\n금액\n↑"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # X축 레이블
    x_label = slide.shapes.add_textbox(
        Inches(4), Inches(6.5), Inches(2), Inches(0.5)
    )
    p = x_label.text_frame.paragraphs[0]
    p.text = "공급 리스크 →"
    p.font.size = Pt(14)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

def create_material_categories_slide(prs):
    """슬라이드 6: 4대 자재군 개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 배경
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)

    # 타이틀
    add_title_shape(slide, "4대 자재군: 차별화된 관리 전략")

    # 4개 자재군 카드
    materials = [
        ("🔴 병목자재", "공급 확보", "ROP\n높은 안전재고", RGBColor(255, 200, 200)),
        ("🟢 레버리지자재", "원가 절감", "MRP\n경쟁 입찰", RGBColor(200, 255, 200)),
        ("🟣 전략자재", "파트너십", "하이브리드\n장기 계약", RGBColor(230, 200, 255)),
        ("⚪ 일상자재", "효율화", "자동화\nVMI", RGBColor(240, 240, 240)),
    ]

    for i, (name, goal, strategy, color) in enumerate(materials):
        row = i // 2
        col = i % 2

        shape = slide.shapes.add_shape(
            1,
            Inches(0.7 + col * 4.8),
            Inches(2 + row * 2.5),
            Inches(4.3),
            Inches(2.2)
        )

        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(150, 150, 150)
        shape.line.width = Pt(2)

        text_frame = shape.text_frame
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.margin_top = Inches(0.2)
        text_frame.margin_left = Inches(0.2)

        # 이름
        p = text_frame.paragraphs[0]
        p.text = name
        p.font.size = Pt(18)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # 목표
        p2 = text_frame.add_paragraph()
        p2.text = f"\n목표: {goal}"
        p2.font.size = Pt(14)
        p2.alignment = PP_ALIGN.CENTER

        # 전략
        p3 = text_frame.add_paragraph()
        p3.text = f"\n전략: {strategy}"
        p3.font.size = Pt(12)
        p3.alignment = PP_ALIGN.CENTER

def main():
    """메인 실행 함수"""
    print("1회차 강의자료 샘플 생성 중...")

    prs = create_presentation()

    output_file = "/home/user/Kraljic_Course/Session1_Sample.pptx"
    prs.save(output_file)

    print(f"✅ 완료! 파일 저장: {output_file}")
    print(f"📊 총 {len(prs.slides)} 슬라이드 생성")

if __name__ == "__main__":
    main()
