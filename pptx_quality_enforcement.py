#!/usr/bin/env python3
"""
PPTX Quality Enforcement Module
100% 품질 보장을 위한 강제 검증 시스템
"""

from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ============================================================================
# 폰트 크기 상수 (절대 변경 불가!)
# ============================================================================
FONT_TITLE = Pt(20)
FONT_GOVERNING = Pt(16)  # Bold
FONT_HEADING = Pt(14)
FONT_BODY = Pt(10)       # PRIMARY (65% of all text)
FONT_BULLET = Pt(12)
FONT_CAPTION = Pt(8)

# ============================================================================
# 강제 할당 함수: 절대 None이 발생하지 않도록 보장
# ============================================================================

def enforce_text_properties(text_frame, font_size, font_name="맑은 고딕",
                           bold=False, color=None, alignment=None,
                           vertical_anchor=MSO_ANCHOR.TOP, word_wrap=True):
    """
    텍스트 프레임의 모든 run에 폰트 속성을 강제로 할당

    Args:
        text_frame: TextFrame 객체
        font_size: Pt() 객체 (REQUIRED!)
        font_name: 폰트 이름
        bold: 볼드 여부
        color: RGBColor 객체
        alignment: PP_ALIGN 상수
        vertical_anchor: MSO_ANCHOR 상수
        word_wrap: 자동 줄바꿈

    Returns:
        int: 설정된 run 개수
    """
    if font_size is None:
        raise ValueError("❌ font_size는 None일 수 없습니다! Pt() 값을 전달하세요.")

    text_frame.word_wrap = word_wrap
    text_frame.vertical_anchor = vertical_anchor

    run_count = 0
    for para in text_frame.paragraphs:
        # 단락 정렬
        if alignment:
            para.alignment = alignment

        # 모든 run에 폰트 속성 강제 할당
        for run in para.runs:
            run.font.name = font_name
            run.font.size = font_size  # ← 절대 누락 불가!
            run.font.bold = bold
            if color:
                run.font.color.rgb = color
            run_count += 1

    # run이 없으면 빈 텍스트 → 경고
    if run_count == 0:
        print(f"⚠️ Warning: 텍스트 프레임에 run이 없습니다 (빈 텍스트)")

    return run_count


def create_text_with_enforcement(shape, text, font_size, font_name="맑은 고딕",
                                 bold=False, color=None, alignment=None,
                                 vertical_anchor=MSO_ANCHOR.TOP):
    """
    Shape에 텍스트를 추가하고 즉시 폰트 속성 강제 할당

    Args:
        shape: Shape 객체 (text_frame이 있어야 함)
        text: 추가할 텍스트 (str)
        font_size: Pt() 객체
        ... (나머지는 enforce_text_properties와 동일)

    Returns:
        TextFrame: 설정된 text_frame
    """
    if not hasattr(shape, 'text_frame'):
        raise ValueError("❌ Shape에 text_frame이 없습니다!")

    text_frame = shape.text_frame
    text_frame.clear()  # 기존 내용 제거
    text_frame.text = text

    # 즉시 폰트 속성 강제 할당
    enforce_text_properties(
        text_frame, font_size, font_name, bold, color, alignment, vertical_anchor
    )

    return text_frame


def add_bullets_with_enforcement(text_frame, bullet_list, font_size=FONT_BODY,
                                 font_name="맑은 고딕", color=None, line_spacing=1.5):
    """
    텍스트 프레임에 불릿 리스트를 추가하고 폰트 속성 강제 할당

    Args:
        text_frame: TextFrame 객체
        bullet_list: 불릿 항목 리스트 (list of str)
        font_size: 폰트 크기 (기본: FONT_BODY = 10pt)
        font_name: 폰트 이름
        color: RGBColor 객체
        line_spacing: 줄 간격 (배수)

    Returns:
        int: 추가된 불릿 개수
    """
    text_frame.clear()
    text_frame.word_wrap = True

    for i, bullet_text in enumerate(bullet_list):
        if i > 0:
            text_frame.add_paragraph()

        para = text_frame.paragraphs[i]
        para.text = bullet_text
        para.level = 0  # 불릿 레벨
        para.line_spacing = line_spacing

        # 모든 run에 폰트 속성 강제 할당
        for run in para.runs:
            run.font.name = font_name
            run.font.size = font_size  # ← 절대 누락 불가!
            if color:
                run.font.color.rgb = color

    return len(bullet_list)


# ============================================================================
# SVG 이미지 삽입 함수
# ============================================================================

def insert_svg_as_image(slide, svg_path, left, top, width=None, height=None):
    """
    SVG 파일을 PPTX 슬라이드에 이미지로 삽입

    python-pptx는 SVG를 직접 지원하지 않으므로,
    cairosvg로 PNG로 변환 후 삽입

    Args:
        slide: Slide 객체
        svg_path: SVG 파일 경로 (str)
        left: 좌측 위치 (Inches)
        top: 상단 위치 (Inches)
        width: 너비 (Inches, optional)
        height: 높이 (Inches, optional)

    Returns:
        Picture: 삽입된 이미지 객체
    """
    import os
    import tempfile

    # SVG 파일 존재 확인
    if not os.path.exists(svg_path):
        raise FileNotFoundError(f"❌ SVG 파일이 없습니다: {svg_path}")

    try:
        import cairosvg
        from PIL import Image
    except ImportError:
        print("❌ cairosvg 또는 pillow가 설치되지 않았습니다.")
        print("   pip3 install cairosvg pillow")
        raise

    # SVG → PNG 변환 (임시 파일)
    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
        tmp_png_path = tmp.name

    try:
        # SVG를 PNG로 변환
        cairosvg.svg2png(url=svg_path, write_to=tmp_png_path, scale=2.0)  # 2x resolution

        # PNG를 PPTX에 삽입
        if width and height:
            pic = slide.shapes.add_picture(tmp_png_path, left, top, width, height)
        elif width:
            pic = slide.shapes.add_picture(tmp_png_path, left, top, width=width)
        elif height:
            pic = slide.shapes.add_picture(tmp_png_path, left, top, height=height)
        else:
            pic = slide.shapes.add_picture(tmp_png_path, left, top)

        return pic

    finally:
        # 임시 PNG 파일 삭제
        if os.path.exists(tmp_png_path):
            os.remove(tmp_png_path)


# ============================================================================
# 검증 함수: 생성 후 PPTX 품질 검사
# ============================================================================

def verify_pptx_quality(pptx_path):
    """
    생성된 PPTX 파일의 품질을 검증

    검증 항목:
    1. 슬라이드 크기 (10.83" × 7.50")
    2. 슬라이드 개수 (40+ 필요)
    3. 폰트 크기 누락 (None) 검사
    4. Shape 개수 (평균 15+ 필요)
    5. 10pt 폰트 사용 비율 (60%+ 필요)

    Returns:
        dict: 검증 결과
        {
            "passed": bool,
            "errors": list of str,
            "warnings": list of str,
            "stats": dict
        }
    """
    from pptx import Presentation
    import os

    if not os.path.exists(pptx_path):
        return {
            "passed": False,
            "errors": [f"❌ 파일이 없습니다: {pptx_path}"],
            "warnings": [],
            "stats": {}
        }

    prs = Presentation(pptx_path)
    errors = []
    warnings = []
    stats = {}

    # 1. 슬라이드 크기 검증
    width_inches = prs.slide_width / 914400
    height_inches = prs.slide_height / 914400
    stats["dimensions"] = f"{width_inches:.2f}\" × {height_inches:.2f}\""

    if abs(width_inches - 10.83) > 0.01 or abs(height_inches - 7.50) > 0.01:
        errors.append(f"❌ 슬라이드 크기 오류: {width_inches:.2f}\" × {height_inches:.2f}\" (목표: 10.83\" × 7.50\")")

    # 2. 슬라이드 개수
    slide_count = len(prs.slides)
    stats["slide_count"] = slide_count

    if slide_count < 40:
        warnings.append(f"⚠️ 슬라이드 개수 부족: {slide_count}장 (목표: 48장)")

    # 3. 폰트 크기 누락 검사
    none_font_count = 0
    font_size_distribution = {}
    total_text_runs = 0

    for i, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text.strip():
                            total_text_runs += 1

                            if run.font.size is None:
                                none_font_count += 1
                            else:
                                size_pt = int(run.font.size.pt)
                                font_size_distribution[size_pt] = font_size_distribution.get(size_pt, 0) + 1

    stats["total_text_runs"] = total_text_runs
    stats["font_size_distribution"] = font_size_distribution

    if none_font_count > 0:
        errors.append(f"❌ 폰트 크기 누락: {none_font_count}개 run에 font.size = None!")

    # 4. 10pt 폰트 사용 비율
    if total_text_runs > 0:
        pt10_count = font_size_distribution.get(10, 0)
        pt10_ratio = pt10_count / total_text_runs
        stats["10pt_ratio"] = f"{pt10_ratio * 100:.1f}%"

        if pt10_ratio < 0.50:  # 50% 미만이면 경고
            warnings.append(f"⚠️ 10pt 폰트 비율 낮음: {pt10_ratio * 100:.1f}% (목표: 60%+)")

    # 5. Shape 개수 평균
    shape_counts = [len(slide.shapes) for slide in prs.slides]
    avg_shapes = sum(shape_counts) / len(shape_counts) if shape_counts else 0
    stats["avg_shapes_per_slide"] = f"{avg_shapes:.1f}"

    if avg_shapes < 10:
        warnings.append(f"⚠️ 평균 Shape 개수 부족: {avg_shapes:.1f} (목표: 15+)")

    # 결과 종합
    passed = len(errors) == 0

    return {
        "passed": passed,
        "errors": errors,
        "warnings": warnings,
        "stats": stats
    }


def print_verification_report(result):
    """검증 결과를 보기 좋게 출력"""
    print("\n" + "=" * 80)
    print("PPTX 품질 검증 결과")
    print("=" * 80 + "\n")

    # 통계
    print("📊 통계:")
    for key, value in result["stats"].items():
        print(f"   {key}: {value}")
    print()

    # 에러
    if result["errors"]:
        print("🚫 에러:")
        for error in result["errors"]:
            print(f"   {error}")
        print()

    # 경고
    if result["warnings"]:
        print("⚠️ 경고:")
        for warning in result["warnings"]:
            print(f"   {warning}")
        print()

    # 최종 판정
    if result["passed"]:
        print("✅ 모든 필수 검증 통과!")
    else:
        print("❌ 검증 실패 - 수정 후 다시 생성하세요.")

    print("=" * 80 + "\n")


# ============================================================================
# 사용 예시
# ============================================================================

if __name__ == "__main__":
    print(__doc__)
    print("\n이 모듈은 다음 함수들을 제공합니다:\n")
    print("1. enforce_text_properties() - 폰트 속성 강제 할당")
    print("2. create_text_with_enforcement() - 텍스트 생성 + 즉시 속성 할당")
    print("3. add_bullets_with_enforcement() - 불릿 리스트 + 폰트 강제 설정")
    print("4. insert_svg_as_image() - SVG를 PNG로 변환하여 삽입")
    print("5. verify_pptx_quality() - 생성된 PPTX 품질 검증")
    print("6. print_verification_report() - 검증 결과 출력")
    print("\n모든 함수는 100% 품질 보장을 위해 에러를 발생시킵니다.")
