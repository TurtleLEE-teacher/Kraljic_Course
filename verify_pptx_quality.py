#!/usr/bin/env python3
"""
PPTX Quality Verification Script
Enforces S4HANA design standards for Part 1-9 consistency

Usage:
    python3 verify_pptx_quality.py Part1_Session1_StrategicInventory.pptx
"""

import sys
from pptx import Presentation
from pptx.util import Inches


def verify_pptx(filepath):
    """Verify PPTX meets quality standards"""
    print(f"\n{'='*60}")
    print(f"PPTX Quality Verification: {filepath}")
    print(f"{'='*60}\n")

    try:
        prs = Presentation(filepath)
    except Exception as e:
        print(f"❌ FATAL: Cannot open file: {e}")
        return False

    failures = []
    warnings = []

    # Check 1: Slide dimensions (CRITICAL)
    expected_width = 914400 * 10.83  # EMUs
    expected_height = 914400 * 7.5

    if abs(prs.slide_width - expected_width) > 100:
        failures.append(
            f"❌ Width: {prs.slide_width/914400:.2f}\" (should be 10.83\")"
        )

    if abs(prs.slide_height - expected_height) > 100:
        failures.append(
            f"❌ Height: {prs.slide_height/914400:.2f}\" (should be 7.5\")"
        )

    # Check 2: Slide count
    slide_count = len(list(prs.slides))
    if slide_count < 20:
        failures.append(
            f"❌ Only {slide_count} slides (expected 20+ for comprehensive coverage)"
        )
    elif slide_count > 35:
        warnings.append(
            f"⚠️  {slide_count} slides (consider splitting if >35)"
        )

    # Check 3: Shapes per slide (content density)
    low_density_slides = []
    shape_counts = []

    all_slides = list(prs.slides)
    for i, slide in enumerate(all_slides[1:], 2):  # Skip cover slide
        shape_count = len(slide.shapes)
        shape_counts.append(shape_count)

        if shape_count < 10:
            low_density_slides.append(f"Slide {i}: {shape_count} shapes")

    avg_shapes = sum(shape_counts) / len(shape_counts) if shape_counts else 0

    if low_density_slides:
        failures.append(
            f"❌ Low shape count (target: 20+):\n   " +
            "\n   ".join(low_density_slides[:5])
        )
        if len(low_density_slides) > 5:
            failures.append(f"   ... and {len(low_density_slides) - 5} more")

    if avg_shapes < 15:
        failures.append(
            f"❌ Average shapes per slide: {avg_shapes:.1f} (should be 20+)"
        )

    # Check 4: Font sizes distribution
    font_sizes = {}
    total_runs = 0

    sample_slides = list(prs.slides)[:10]
    for slide in sample_slides:  # Sample first 10 slides
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            size = int(run.font.size.pt)
                            font_sizes[size] = font_sizes.get(size, 0) + 1
                            total_runs += 1

    if total_runs > 0:
        pt10_ratio = font_sizes.get(10, 0) / total_runs
        pt12_ratio = font_sizes.get(12, 0) / total_runs

        # 10pt should be PRIMARY (target: 65%, accept ≥40%)
        if pt10_ratio < 0.4:
            failures.append(
                f"❌ 10pt text ratio: {pt10_ratio*100:.1f}% (should be 60%+)"
            )
        elif pt10_ratio < 0.55:
            warnings.append(
                f"⚠️  10pt text ratio: {pt10_ratio*100:.1f}% (target: 65%+)"
            )

        # 12pt for bullets (target: 20-25%, accept 15-35%)
        if pt12_ratio > 0.4:
            warnings.append(
                f"⚠️  12pt text ratio: {pt12_ratio*100:.1f}% (should be 20-25%)"
            )

    # Check 5: Governing messages presence (sample check)
    slides_without_gov_msg = []

    sample_check_slides = list(prs.slides)[1:6]
    for i, slide in enumerate(sample_check_slides, 2):  # Check slides 2-6
        # Look for text boxes with 16pt Bold text (governing message signature)
        has_gov_msg = False

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if (run.font.size and
                            int(run.font.size.pt) == 16 and
                            run.font.bold):
                            has_gov_msg = True
                            break

        if not has_gov_msg:
            slides_without_gov_msg.append(i)

    if slides_without_gov_msg:
        warnings.append(
            f"⚠️  Possible missing governing messages (16pt Bold) on slides: {slides_without_gov_msg}"
        )

    # Print results
    print("=" * 60)
    print("VERIFICATION RESULTS")
    print("=" * 60)

    print(f"\n✓ Slide count: {slide_count}")
    print(f"✓ Dimensions: {prs.slide_width/914400:.2f}\" × {prs.slide_height/914400:.2f}\"")
    print(f"✓ Average shapes per slide: {avg_shapes:.1f}")

    if total_runs > 0:
        print(f"\n✓ Font size distribution (first 10 slides):")
        for size in sorted(font_sizes.keys()):
            ratio = font_sizes[size] / total_runs
            print(f"   {size}pt: {ratio*100:.1f}%")

    if warnings:
        print(f"\n⚠️  WARNINGS ({len(warnings)}):")
        for warning in warnings:
            print(f"   {warning}")

    if failures:
        print(f"\n❌ FAILURES ({len(failures)}):")
        for failure in failures:
            print(f"   {failure}")
        print(f"\n{'='*60}")
        print("🚫 QUALITY CHECK FAILED - DO NOT PROCEED")
        print("   Fix issues above and regenerate PPTX")
        print(f"{'='*60}\n")
        return False
    else:
        print(f"\n{'='*60}")
        print("✅ ALL QUALITY CHECKS PASSED")
        print(f"{'='*60}\n")
        return True


if __name__ == '__main__':
    if len(sys.argv) < 2:
        print("Usage: python3 verify_pptx_quality.py <pptx_file>")
        sys.exit(1)

    filepath = sys.argv[1]
    success = verify_pptx(filepath)
    sys.exit(0 if success else 1)
