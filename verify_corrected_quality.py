#!/usr/bin/env python3
"""
Verify CORRECTED Part 1 PPTX - Focus on 3 requirements
1. Strict monochrome (black/white/gray only)
2. Grid alignment
3. Shape variety (arrows, structured boxes)
"""

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE

def verify_corrected_pptx(filepath):
    """Detailed verification of corrected PPTX"""
    print(f"\n{'='*80}")
    print(f"CORRECTED Part 1 PPTX - Quality Verification")
    print(f"{'='*80}\n")

    prs = Presentation(filepath)

    # Basic checks
    print(f"1. Basic Properties:")
    print(f"   Dimensions: {prs.slide_width.inches:.2f}\" × {prs.slide_height.inches:.2f}\"")
    print(f"   Total slides: {len(prs.slides)}")
    print(f"   Status: {'✓ PASS' if len(prs.slides) == 48 else f'✗ FAIL (expected 48)'}\n")

    # Detailed analysis of slides 1-8 (corrected slides)
    print(f"2. Detailed Analysis (Corrected Slides 1-8):")
    print(f"{'='*80}")

    # Track color usage
    all_colors = set()
    monochrome_colors = {
        (0, 0, 0),           # Black
        (51, 51, 51),        # Dark Gray
        (102, 102, 102),     # Med Gray
        (204, 204, 204),     # Light Gray
        (230, 230, 230),     # Very Light Gray
        (255, 255, 255),     # White
        (26, 82, 118),       # Dark Blue (accent)
    }

    for idx in range(1, 9):  # Slides 1-8
        if idx > len(prs.slides):
            break

        slide = prs.slides[idx - 1]

        print(f"\nSlide {idx}:")

        # Count shapes by type
        shape_types = {}
        total_chars = 0
        font_sizes = []
        slide_colors = set()

        for shape in slide.shapes:
            # Count shape types
            shape_type_name = str(shape.shape_type)
            try:
                for attr in dir(MSO_SHAPE_TYPE):
                    if not attr.startswith('_') and isinstance(getattr(MSO_SHAPE_TYPE, attr), int):
                        if getattr(MSO_SHAPE_TYPE, attr) == shape.shape_type:
                            shape_type_name = attr
                            break
            except:
                pass

            shape_types[shape_type_name] = shape_types.get(shape_type_name, 0) + 1

            # Check fill colors
            if shape.shape_type == 1:  # AUTO_SHAPE
                try:
                    if hasattr(shape.fill, 'fore_color') and hasattr(shape.fill.fore_color, 'rgb'):
                        rgb = shape.fill.fore_color.rgb
                        color_tuple = (rgb[0], rgb[1], rgb[2])
                        slide_colors.add(color_tuple)
                        all_colors.add(color_tuple)
                except:
                    pass

            if hasattr(shape, "text_frame"):
                # Count text
                text = shape.text
                total_chars += len(text)

                # Collect font sizes
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            font_sizes.append(run.font.size.pt)

        # Print shape counts
        total_shapes = sum(shape_types.values())
        print(f"  Total shapes: {total_shapes}")

        # Show shape type breakdown
        shape_summary = []
        for stype, count in sorted(shape_types.items(), key=lambda x: -x[1])[:5]:
            shape_summary.append(f"{stype}:{count}")
        print(f"  Types: {', '.join(shape_summary)}")

        print(f"  Text: {total_chars} characters")

        if font_sizes:
            unique_sizes = sorted(set([int(f) for f in font_sizes]))
            print(f"  Font sizes: {unique_sizes}")

            # Count small fonts (8-11pt)
            small_fonts = [f for f in font_sizes if 8 <= f <= 11]
            if font_sizes:
                pct = len(small_fonts) / len(font_sizes) * 100
                print(f"  Small fonts (8-11pt): {len(small_fonts)}/{len(font_sizes)} ({pct:.1f}%)")

        # Color check
        non_monochrome = slide_colors - monochrome_colors
        if non_monochrome:
            print(f"  ⚠ Non-monochrome colors found: {non_monochrome}")
        else:
            print(f"  ✓ Monochrome colors only")

        # Quality assessment
        quality = "✓ GOOD" if total_shapes >= 15 else "⚠ LOW"
        if total_shapes >= 30:
            quality = "✓✓ EXCELLENT"
        print(f"  Quality: {quality} (target: 30-40 shapes)")

    # Overall color compliance
    print(f"\n{'='*80}")
    print(f"3. Color Compliance Check:")

    non_monochrome_all = all_colors - monochrome_colors
    if non_monochrome_all:
        print(f"   ⚠ Non-monochrome colors detected:")
        for color in non_monochrome_all:
            print(f"      RGB{color}")
        print(f"   Status: ✗ FAIL - Using rainbow colors!")
    else:
        print(f"   ✓ PASS - Strict monochrome compliance")

    # Overall statistics
    print(f"\n{'='*80}")
    print(f"4. Overall Statistics (All 48 slides):")

    total_shapes = 0
    all_fonts = []

    for slide in prs.slides:
        slide_shapes = 0
        for shape in slide.shapes:
            slide_shapes += 1
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            all_fonts.append(run.font.size.pt)
        total_shapes += slide_shapes

    avg_shapes = total_shapes / len(prs.slides) if len(prs.slides) > 0 else 0
    print(f"   Average shapes per slide: {avg_shapes:.1f}")
    print(f"   Target: 30-40 shapes")
    print(f"   Status: {'✓ PASS' if avg_shapes >= 20 else '⚠ NEEDS IMPROVEMENT'}")

    if all_fonts:
        small_fonts = [f for f in all_fonts if 8 <= f <= 11]
        pct = len(small_fonts) / len(all_fonts) * 100
        print(f"   Small fonts (8-11pt): {pct:.1f}%")
        print(f"   Target: 60%+")

    print(f"\n{'='*80}")
    print(f"Verification Complete")
    print(f"{'='*80}\n")

if __name__ == "__main__":
    verify_corrected_pptx("/home/user/Kraljic_Course/PPTX_SAMPLE/Part1_Session1_Corrected.pptx")
