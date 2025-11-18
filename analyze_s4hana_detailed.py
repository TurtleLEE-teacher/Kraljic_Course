#!/usr/bin/env python3
"""
Detailed analysis of S4HANA reference PPTX
- Font sizes (all unique sizes)
- Layout patterns (text box positions and sizes)
- Shape usage (rectangles, arrows, connectors)
- Content density (% of slide area filled)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import statistics

def analyze_s4hana(filepath):
    """Detailed analysis of S4HANA PPTX"""
    print(f"\n{'='*80}")
    print(f"S4HANA Reference PPTX - Detailed Analysis")
    print(f"{'='*80}\n")

    prs = Presentation(filepath)

    # Slide dimensions
    width_in = prs.slide_width.inches
    height_in = prs.slide_height.inches
    slide_area = width_in * height_in

    print(f"Slide Dimensions: {width_in:.2f}\" × {height_in:.2f}\"")
    print(f"Slide Area: {slide_area:.2f} sq in\n")

    # Collect statistics
    all_font_sizes = []
    all_fonts = set()
    shapes_by_type = {}
    content_densities = []

    print(f"{'='*80}")
    print(f"Analyzing {len(prs.slides)} slides...\n")

    for idx, slide in enumerate(prs.slides, 1):
        print(f"--- Slide {idx} ---")

        # Calculate content density
        filled_area = 0
        shape_count = {"text": 0, "shape": 0, "picture": 0, "table": 0, "group": 0}

        for shape in slide.shapes:
            # Shape area
            if hasattr(shape, "width") and hasattr(shape, "height"):
                shape_area = (shape.width.inches * shape.height.inches)
                filled_area += shape_area

            # Shape type
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                shape_count["text"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shape_count["picture"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shape_count["table"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shape_count["group"] += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shape_count["shape"] += 1
                shapes_by_type[shape.shape_type] = shapes_by_type.get(shape.shape_type, 0) + 1

            # Text analysis
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            all_font_sizes.append(run.font.size.pt)
                        if run.font.name:
                            all_fonts.add(run.font.name)

            # Print shape details for first 3 slides
            if idx <= 5:
                if hasattr(shape, "text") and shape.text.strip():
                    text_preview = shape.text.strip()[:40].replace("\n", " ")
                    print(f"  Shape: {shape.shape_type}")
                    if hasattr(shape, "left") and hasattr(shape, "top"):
                        print(f"    Position: ({shape.left.inches:.2f}\", {shape.top.inches:.2f}\")")
                    if hasattr(shape, "width") and hasattr(shape, "height"):
                        print(f"    Size: {shape.width.inches:.2f}\" × {shape.height.inches:.2f}\"")
                    print(f"    Text: \"{text_preview}...\"")

                    if hasattr(shape, "text_frame"):
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.size:
                                    print(f"    Font: {run.font.name}, {run.font.size.pt}pt, Bold: {run.font.bold}")
                                    break
                            break

        density = (filled_area / slide_area) * 100
        content_densities.append(density)

        print(f"  Shapes: Text={shape_count['text']}, Shape={shape_count['shape']}, "
              f"Picture={shape_count['picture']}, Table={shape_count['table']}, Group={shape_count['group']}")
        print(f"  Content Density: {density:.1f}%")
        print()

    # Summary statistics
    print(f"\n{'='*80}")
    print(f"Summary Statistics")
    print(f"{'='*80}\n")

    print(f"Font Sizes Used:")
    if all_font_sizes:
        unique_sizes = sorted(set(all_font_sizes))
        print(f"  Unique sizes: {', '.join([f'{s:.0f}pt' for s in unique_sizes])}")
        print(f"  Most common: {statistics.mode(all_font_sizes):.0f}pt")
        print(f"  Range: {min(all_font_sizes):.0f}pt - {max(all_font_sizes):.0f}pt")

    print(f"\nFonts Used:")
    for font in sorted(all_fonts):
        print(f"  - {font}")

    print(f"\nContent Density:")
    if content_densities:
        print(f"  Average: {statistics.mean(content_densities):.1f}%")
        print(f"  Median: {statistics.median(content_densities):.1f}%")
        print(f"  Min: {min(content_densities):.1f}%")
        print(f"  Max: {max(content_densities):.1f}%")

        # Slides with >85% density
        high_density = [i+1 for i, d in enumerate(content_densities) if d >= 85]
        print(f"  Slides with ≥85% density: {len(high_density)}/{len(content_densities)}")
        if high_density:
            print(f"    Slide numbers: {', '.join(map(str, high_density[:10]))}")

    print(f"\n{'='*80}\n")

if __name__ == "__main__":
    analyze_s4hana("/home/user/Kraljic_Course/PPTX_SAMPLE/S4HANA_PI단계_단계 종료보고_20230510_v.1.4.pptx")
